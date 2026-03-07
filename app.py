import streamlit as st
import fitz  # PyMuPDF
import google.generativeai as genai
import re
import json
import genanki
import tempfile
import os
import io
import uuid
import concurrent.futures
import random
import sys
import time
from datetime import datetime, timezone, timedelta
from pathlib import Path
from openai import OpenAI
from docx import Document
from docx.oxml import OxmlElement
from pptx import Presentation
from difflib import SequenceMatcher
import subprocess
import shutil
import base64
import zipfile
import xml.etree.ElementTree as ET
import importlib.util
import hashlib
import requests
from src.prompts import PROMPT_MCQ, PROMPT_CLOZE, PROMPT_SHORT, PROMPT_ESSAY
from src.repositories import load_json_file, save_json_file

# ============================================================================
# 감사 로그 (append-only JSONL)
# ============================================================================
def _hash_text(text: str) -> str:
    if not text:
        return ""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()[:16]

def append_audit_log(event: str, payload: dict, user_id=None):
    try:
        row = {
            "run_id": str(int(time.time() * 1000)),
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "event": event,
            **(payload or {}),
        }
        with open(get_audit_log_file(user_id=user_id), "a", encoding="utf-8") as f:
            f.write(json.dumps(row, ensure_ascii=False) + "\n")
    except Exception:
        # 감사 로그 실패는 앱 실행을 막지 않음
        pass

def _gemini_usage_tokens(response):
    usage = getattr(response, "usage_metadata", None)
    if not usage:
        return None
    total = getattr(usage, "total_token_count", None)
    if total is not None:
        return total
    prompt_t = getattr(usage, "prompt_token_count", None)
    cand_t = getattr(usage, "candidates_token_count", None)
    if prompt_t is None and cand_t is None:
        return None
    return (prompt_t or 0) + (cand_t or 0)

def _openai_usage_tokens(response):
    usage = getattr(response, "usage", None)
    if usage is None:
        return None
    if isinstance(usage, dict):
        return usage.get("total_tokens", None)
    return getattr(usage, "total_tokens", None)
# FSRS (optional)
try:
    from fsrs import Scheduler, Card, Rating, ReviewLog
    FSRS_AVAILABLE = True
except Exception:
    FSRS_AVAILABLE = False

FSRS_DEFAULT_PARAMETERS = (
    0.212, 1.2931, 2.3065, 8.2956, 6.4133, 0.8334, 3.0194, 0.001,
    1.8722, 0.1666, 0.796, 1.4835, 0.0614, 0.2629, 1.6483, 0.6014,
    1.8729, 0.5425, 0.0912, 0.0658, 0.1542,
)

MODE_MCQ = "📝 객관식 문제 (Case Study)"
MODE_CLOZE = "🧩 빈칸 뚫기 (Anki Cloze)"
MODE_SHORT = "🧠 단답형 문제"
MODE_ESSAY = "🧾 서술형 문제"

# ============================================================================
# 초기 설정
# ============================================================================
st.set_page_config(page_title="Axioma Qbank", page_icon="🧬", layout="wide")

def get_app_data_dir():
    env_dir = os.getenv("AXIOMA_QBANK_DATA_DIR", "").strip() or os.getenv("MEDTUTOR_DATA_DIR", "").strip()
    if env_dir:
        base = Path(env_dir).expanduser()
        try:
            base.mkdir(parents=True, exist_ok=True)
        except Exception:
            pass
        return base
    if getattr(sys, "frozen", False):
        new_base = Path.home() / "AxiomaQbank"
        old_base = Path.home() / "MedTutor"
        base = old_base if old_base.exists() and not new_base.exists() else new_base
        try:
            base.mkdir(parents=True, exist_ok=True)
        except Exception:
            pass
        return base
    base = Path.cwd()
    try:
        base.mkdir(parents=True, exist_ok=True)
    except Exception:
        pass
    return base

DATA_DIR = get_app_data_dir()
QUESTION_BANK_FILE = str(DATA_DIR / "questions.json")
EXAM_HISTORY_FILE = str(DATA_DIR / "exam_history.json")
USER_SETTINGS_FILE = str(DATA_DIR / "user_settings.json")
AUDIT_LOG_FILE = str(DATA_DIR / "audit_log.jsonl")
AUTH_USERS_FILE = str(DATA_DIR / "auth_users.json")
SUPABASE_URL = os.getenv("SUPABASE_URL", "").rstrip("/")
SUPABASE_ANON_KEY = os.getenv("SUPABASE_ANON_KEY", "")
SUPABASE_TABLE = "medtutor_user_data"

def sanitize_user_id(user_id):
    text = (user_id or "").strip()
    if not text:
        return "guest"
    return re.sub(r"[^a-zA-Z0-9._-]", "_", text)[:80] or "guest"

def is_valid_email(value):
    text = (value or "").strip()
    return bool(re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", text))

def get_current_user_id():
    return sanitize_user_id(st.session_state.get("auth_user_id", "guest"))

def get_user_data_dir(user_id=None):
    uid = sanitize_user_id(user_id) if user_id is not None else get_current_user_id()
    base = DATA_DIR / "users" / uid
    try:
        base.mkdir(parents=True, exist_ok=True)
    except Exception:
        pass
    return base

def get_question_bank_file(user_id=None):
    return str(get_user_data_dir(user_id) / "questions.json")

def get_exam_history_file(user_id=None):
    return str(get_user_data_dir(user_id) / "exam_history.json")

def get_user_settings_file(user_id=None):
    return str(get_user_data_dir(user_id) / "user_settings.json")

def get_audit_log_file(user_id=None):
    return str(get_user_data_dir(user_id) / "audit_log.jsonl")

MODEL_PRICING_USD_PER_1M = {
    "gpt-4o-mini": {"input": 0.15, "output": 0.60, "blended": 0.30},
    "gemini-2.5-flash-lite": {"input": 0.10, "output": 0.40, "blended": 0.20},
    "gemini-2.5-flash": {"input": 0.30, "output": 2.50, "blended": 0.90},
}

def get_configured_admin_users():
    raw = os.getenv("AXIOMA_ADMIN_USERS", "").strip()
    if not raw:
        raw = os.getenv("MEDTUTOR_ADMIN_USERS", "").strip()
    users = set()
    for token in raw.split(","):
        value = token.strip()
        if value:
            users.add(value.lower())
    return users

def is_admin_user():
    admins = get_configured_admin_users()
    if not admins:
        return False
    uid = str(st.session_state.get("auth_user_id", "")).strip().lower()
    email = str(st.session_state.get("auth_email", "")).strip().lower()
    return uid in admins or email in admins

def list_local_user_ids():
    users_dir = DATA_DIR / "users"
    if not users_dir.exists():
        return []
    result = []
    for item in users_dir.iterdir():
        if item.is_dir():
            result.append(item.name)
    return sorted(result)

def read_audit_rows_for_user(user_id):
    rows = []
    path = Path(get_audit_log_file(user_id))
    if not path.exists():
        return rows
    try:
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    row = json.loads(line)
                except Exception:
                    continue
                if isinstance(row, dict):
                    rows.append(row)
    except Exception:
        return []
    return rows

def summarize_usage_rows(rows):
    summary_by_model = {}
    for row in rows:
        model = str(row.get("model") or "unknown")
        event = str(row.get("event") or "")
        usage_tokens = row.get("usage_tokens")
        item = summary_by_model.setdefault(
            model,
            {
                "calls": 0,
                "tokens": 0,
                "gen_calls": 0,
                "grade_calls": 0,
            },
        )
        item["calls"] += 1
        if isinstance(usage_tokens, int):
            item["tokens"] += usage_tokens
        if event == "gen.question":
            item["gen_calls"] += 1
        if event.startswith("grade."):
            item["grade_calls"] += 1
    return summary_by_model

def estimate_cost_usd_from_summary(summary_by_model):
    total = 0.0
    breakdown = []
    for model, item in summary_by_model.items():
        tokens = int(item.get("tokens", 0))
        price = MODEL_PRICING_USD_PER_1M.get(model, {}).get("blended")
        if price is None:
            usd = None
        else:
            usd = (tokens / 1_000_000.0) * price
            total += usd
        breakdown.append(
            {
                "model": model,
                "calls": int(item.get("calls", 0)),
                "tokens": tokens,
                "est_cost_usd": usd,
                "gen_calls": int(item.get("gen_calls", 0)),
                "grade_calls": int(item.get("grade_calls", 0)),
            }
        )
    breakdown.sort(key=lambda x: (x["tokens"], x["calls"]), reverse=True)
    return total, breakdown

def is_supabase_enabled():
    return bool(SUPABASE_URL and SUPABASE_ANON_KEY)

def is_supabase_required():
    value = os.getenv("AXIOMA_REQUIRE_SUPABASE", "1").strip().lower()
    return value not in {"0", "false", "no", "off"}

if is_supabase_required() and not is_supabase_enabled():
    st.title("Axioma Qbank")
    st.error("이 배포는 Supabase 인증/저장을 필수로 사용합니다.")
    st.info("Secrets에 SUPABASE_URL, SUPABASE_ANON_KEY를 설정한 뒤 앱을 재시작하세요.")
    st.stop()

def _supabase_headers(access_token=None):
    token = access_token or SUPABASE_ANON_KEY
    return {
        "apikey": SUPABASE_ANON_KEY,
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
    }

def _supabase_error_message(response):
    try:
        data = response.json()
        return data.get("msg") or data.get("error_description") or data.get("message") or f"HTTP {response.status_code}"
    except Exception:
        return f"HTTP {response.status_code}"

def _default_remote_bundle():
    return {
        "questions": {"text": [], "cloze": []},
        "exam_history": [],
        "user_settings": {},
    }

def _normalize_remote_bundle(bundle):
    data = bundle if isinstance(bundle, dict) else {}
    questions = data.get("questions")
    exam_history = data.get("exam_history")
    user_settings = data.get("user_settings")
    return {
        "questions": questions if isinstance(questions, dict) else {"text": [], "cloze": []},
        "exam_history": exam_history if isinstance(exam_history, list) else [],
        "user_settings": user_settings if isinstance(user_settings, dict) else {},
    }

def supabase_sign_up(email, password):
    if not is_supabase_enabled():
        return False, "SUPABASE_URL / SUPABASE_ANON_KEY 설정이 필요합니다."
    payload = {"email": (email or "").strip(), "password": password or ""}
    if not payload["email"]:
        return False, "이메일을 입력해주세요."
    if len(payload["password"]) < 6:
        return False, "비밀번호는 6자 이상이어야 합니다."
    resp = requests.post(
        f"{SUPABASE_URL}/auth/v1/signup",
        headers=_supabase_headers(),
        json=payload,
        timeout=10,
    )
    if resp.status_code not in (200, 201):
        return False, _supabase_error_message(resp)
    return True, "회원가입이 완료되었습니다. 이메일 인증 설정이 켜져 있다면 인증 후 로그인하세요."

def supabase_sign_in(email, password):
    if not is_supabase_enabled():
        return False, "SUPABASE_URL / SUPABASE_ANON_KEY 설정이 필요합니다."
    payload = {"email": (email or "").strip(), "password": password or ""}
    if not payload["email"] or not payload["password"]:
        return False, "이메일과 비밀번호를 입력해주세요."
    resp = requests.post(
        f"{SUPABASE_URL}/auth/v1/token?grant_type=password",
        headers=_supabase_headers(),
        json=payload,
        timeout=10,
    )
    if resp.status_code != 200:
        return False, _supabase_error_message(resp)
    data = resp.json() or {}
    access_token = data.get("access_token")
    user = data.get("user") or {}
    user_id = user.get("id")
    email_value = user.get("email") or payload["email"]
    if not access_token or not user_id:
        return False, "로그인 토큰 정보를 읽지 못했습니다."
    return True, {"user_id": user_id, "email": email_value, "access_token": access_token}

def supabase_fetch_user_bundle(user_id, access_token):
    if not (is_supabase_enabled() and user_id and access_token):
        return None
    params = {"select": "questions,exam_history,user_settings", "user_id": f"eq.{user_id}"}
    resp = requests.get(
        f"{SUPABASE_URL}/rest/v1/{SUPABASE_TABLE}",
        headers=_supabase_headers(access_token),
        params=params,
        timeout=10,
    )
    if resp.status_code != 200:
        return None
    rows = resp.json() or []
    if not rows:
        return _default_remote_bundle()
    return _normalize_remote_bundle(rows[0])

def supabase_upsert_user_bundle(user_id, access_token, bundle):
    if not (is_supabase_enabled() and user_id and access_token):
        return False
    body = {
        "user_id": user_id,
        "questions": (bundle or {}).get("questions", {"text": [], "cloze": []}),
        "exam_history": (bundle or {}).get("exam_history", []),
        "user_settings": (bundle or {}).get("user_settings", {}),
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }
    headers = _supabase_headers(access_token)
    headers["Prefer"] = "resolution=merge-duplicates,return=representation"
    resp = requests.post(
        f"{SUPABASE_URL}/rest/v1/{SUPABASE_TABLE}",
        headers=headers,
        json=body,
        timeout=10,
    )
    return resp.status_code in (200, 201)

def use_remote_user_store(user_id=None):
    if not is_supabase_enabled():
        return False
    uid = sanitize_user_id(user_id) if user_id is not None else sanitize_user_id(st.session_state.get("auth_user_id", ""))
    token = st.session_state.get("auth_access_token")
    return bool(uid and uid != "guest" and token)

def load_remote_bundle(force=False):
    if not use_remote_user_store():
        return None
    uid = sanitize_user_id(st.session_state.get("auth_user_id", ""))
    token = st.session_state.get("auth_access_token")
    cache_key = f"remote_bundle:{uid}"
    cache = st.session_state.get("remote_bundle_cache", {})
    if not force and cache_key in cache:
        return _normalize_remote_bundle(cache[cache_key])
    bundle = supabase_fetch_user_bundle(uid, token)
    if bundle is None:
        return None
    if bundle == _default_remote_bundle():
        supabase_upsert_user_bundle(uid, token, bundle)
    cache[cache_key] = bundle
    st.session_state.remote_bundle_cache = cache
    return _normalize_remote_bundle(bundle)

def save_remote_bundle(bundle):
    if not use_remote_user_store():
        return False
    uid = sanitize_user_id(st.session_state.get("auth_user_id", ""))
    token = st.session_state.get("auth_access_token")
    norm = _normalize_remote_bundle(bundle)
    ok = supabase_upsert_user_bundle(uid, token, norm)
    if ok:
        cache = st.session_state.get("remote_bundle_cache", {})
        cache[f"remote_bundle:{uid}"] = norm
        st.session_state.remote_bundle_cache = cache
    return ok

def notify_remote_store_failure(message):
    try:
        st.session_state.last_action_notice = message
    except Exception:
        pass

PROMPT_VERSION = "v1"
GRADER_VERSION = "v1"
LLM_TEMPERATURE = 0.0
def _get_llm_seed():
    try:
        return int(os.getenv("GEN_SEED", "123"))
    except Exception:
        return None
LLM_SEED = _get_llm_seed()
def get_query_param(name, default=None):
    try:
        params = st.query_params
        if name in params:
            val = params[name]
            if isinstance(val, list):
                return val[0] if val else default
            return val
        return default
    except Exception:
        try:
            params = st.experimental_get_query_params()
            return params.get(name, [default])[0]
        except Exception:
            return default

def resolve_theme_mode_from_query(theme_param, default="Light"):
    value = str(theme_param or "").strip().lower()
    if value == "dark":
        return "Dark"
    if value == "light":
        return "Light"
    return default

def resolve_mobile_flag_from_query(mobile_param):
    value = str(mobile_param or "").strip().lower()
    return value in {"1", "true", "yes", "y", "mobile"}

safe_param = get_query_param("safe", None)
ping_param = get_query_param("ping", "0")
theme_param = get_query_param("theme", None)
mobile_param = get_query_param("mobile", "0")
resolved_theme_mode = resolve_theme_mode_from_query(theme_param, default="Light")
MOBILE_CLIENT = resolve_mobile_flag_from_query(mobile_param)

DEBUG_MODE = str(ping_param) == "1"
if DEBUG_MODE:
    st.write("✅ DEBUG: app.py loaded")
    st.write(f"Streamlit version: {st.__version__}")
    st.write(f"safe_param={safe_param}")
    st.stop()

LOCK_SAFE = str(safe_param) == "1"
LOCK_THEME = str(safe_param) == "0"

if "theme_enabled" not in st.session_state:
    st.session_state.theme_enabled = True if safe_param is None else LOCK_THEME
if "auth_user_id" not in st.session_state:
    st.session_state.auth_user_id = ""
if "auth_access_token" not in st.session_state:
    st.session_state.auth_access_token = ""
if "auth_email" not in st.session_state:
    st.session_state.auth_email = ""

# Session State 초기화
if "current_question_idx" not in st.session_state:
    st.session_state.current_question_idx = 0
if "exam_questions" not in st.session_state:
    st.session_state.exam_questions = []
if "user_answers" not in st.session_state:
    st.session_state.user_answers = {}
if "exam_started" not in st.session_state:
    st.session_state.exam_started = False
if "exam_finished" not in st.session_state:
    st.session_state.exam_finished = False
if "exam_mode" not in st.session_state:
    st.session_state.exam_mode = "시험모드"
if "exam_type" not in st.session_state:
    st.session_state.exam_type = "객관식"
if "auto_next" not in st.session_state:
    st.session_state.auto_next = False
if "auto_advance_guard" not in st.session_state:
    st.session_state.auto_advance_guard = None
if "revealed_answers" not in st.session_state:
    st.session_state.revealed_answers = set()
if "explanation_default" not in st.session_state:
    st.session_state.explanation_default = False
if "exam_stats_applied" not in st.session_state:
    st.session_state.exam_stats_applied = False
if "graded_questions" not in st.session_state:
    st.session_state.graded_questions = set()
# (trend_days retained for future use)
if "trend_days" not in st.session_state:
    st.session_state.trend_days = 14
if "wrong_priority" not in st.session_state:
    st.session_state.wrong_priority = "오답 횟수"
if "current_exam_meta" not in st.session_state:
    st.session_state.current_exam_meta = {}
if "exam_history_saved" not in st.session_state:
    st.session_state.exam_history_saved = False
if "gemini_model_id" not in st.session_state:
    st.session_state.gemini_model_id = "gemini-2.5-flash"
if "ai_model" not in st.session_state:
    st.session_state.ai_model = "🔵 Google Gemini"
if "api_key" not in st.session_state:
    st.session_state.api_key = None
if "openai_api_key" not in st.session_state:
    st.session_state.openai_api_key = None
if "dual_exam_text" not in st.session_state:
    st.session_state.dual_exam_text = ""
if "dual_exam_images" not in st.session_state:
    st.session_state.dual_exam_images = []
if "dual_exam_page_text" not in st.session_state:
    st.session_state.dual_exam_page_text = []
if "dual_match_scores" not in st.session_state:
    st.session_state.dual_match_scores = {}
if "wrong_weight_recent" not in st.session_state:
    st.session_state.wrong_weight_recent = 0.7
if "wrong_weight_count" not in st.session_state:
    st.session_state.wrong_weight_count = 0.3
if "theme_mode" not in st.session_state:
    st.session_state.theme_mode = resolved_theme_mode
if "theme_bg" not in st.session_state:
    st.session_state.theme_bg = "Gradient"
if "last_action_notice" not in st.session_state:
    st.session_state.last_action_notice = ""
if "generation_failure" not in st.session_state:
    st.session_state.generation_failure = ""
if "generation_preview_items" not in st.session_state:
    st.session_state.generation_preview_items = []
if "generation_preview_mode" not in st.session_state:
    st.session_state.generation_preview_mode = "📝 객관식 문제 (Case Study)"
if "generation_preview_subject" not in st.session_state:
    st.session_state.generation_preview_subject = "General"
if "generation_preview_unit" not in st.session_state:
    st.session_state.generation_preview_unit = "미분류"
if "generation_prewarm_cache" not in st.session_state:
    st.session_state["generation_prewarm_cache"] = {}
if "generation_prewarm_errors" not in st.session_state:
    st.session_state["generation_prewarm_errors"] = {}
if "generation_async_job" not in st.session_state:
    st.session_state["generation_async_job"] = None
if "export_docx_bytes" not in st.session_state:
    st.session_state.export_docx_bytes = b""
if "exam_mode_entry_anchor" not in st.session_state:
    st.session_state.exam_mode_entry_anchor = ""
if "heatmap_bins" not in st.session_state:
    st.session_state.heatmap_bins = [0, 1, 3, 6, 10]
if "heatmap_colors" not in st.session_state:
    st.session_state.heatmap_colors = ["#ffffff", "#d7f3f0", "#b2e9e3", "#7fd6cc", "#4fc1b6", "#1f8e86"]
if "profile_name" not in st.session_state:
    st.session_state.profile_name = "default"
if "select_placeholder_exam" not in st.session_state:
    st.session_state.select_placeholder_exam = "선택하세요"
if "select_placeholder_study" not in st.session_state:
    st.session_state.select_placeholder_study = "선택하세요"
if "past_exam_text" not in st.session_state:
    st.session_state.past_exam_text = ""
if "past_exam_items" not in st.session_state:
    st.session_state.past_exam_items = []
if "past_exam_file" not in st.session_state:
    st.session_state.past_exam_file = ""
if "past_exam_images" not in st.session_state:
    st.session_state.past_exam_images = []
if "image_display_width" not in st.session_state:
    st.session_state.image_display_width = 520
if "past_exam_anchors" not in st.session_state:
    st.session_state.past_exam_anchors = {}
if "user_data_cache" not in st.session_state:
    st.session_state["user_data_cache"] = {}
if "home_visual_loaded" not in st.session_state:
    st.session_state.home_visual_loaded = False

def reset_runtime_state_for_auth_change():
    volatile_keys = [
        "exam_questions",
        "user_answers",
        "revealed_answers",
        "graded_questions",
        "current_exam_meta",
        "exam_started",
        "exam_finished",
        "exam_history_saved",
        "generation_preview_items",
        "generation_failure",
        "generation_prewarm_cache",
        "generation_prewarm_errors",
        "generation_async_job",
        "last_action_notice",
        "past_exam_items",
        "past_exam_images",
        "past_exam_text",
        "fsrs_settings_initialized",
        "remote_bundle_cache",
        "user_data_cache",
    ]
    for key in volatile_keys:
        if key in st.session_state:
            del st.session_state[key]

def _user_data_cache_key(kind, user_id=None):
    if user_id is not None:
        uid = sanitize_user_id(user_id)
        scope = "local"
    elif use_remote_user_store():
        uid = sanitize_user_id(st.session_state.get("auth_user_id", ""))
        scope = "remote"
    else:
        uid = get_current_user_id()
        scope = "local"
    return f"{kind}:{scope}:{uid}"

def _get_user_data_cache(kind, user_id=None):
    cache = st.session_state.get("user_data_cache", {})
    return cache.get(_user_data_cache_key(kind, user_id=user_id))

def _set_user_data_cache(kind, value, user_id=None):
    cache = st.session_state.get("user_data_cache", {})
    cache[_user_data_cache_key(kind, user_id=user_id)] = value
    st.session_state["user_data_cache"] = cache
    return value

def _get_or_load_user_data(kind, loader, user_id=None, force=False):
    if not force:
        cached = _get_user_data_cache(kind, user_id=user_id)
        if cached is not None:
            return cached
    data = loader()
    return _set_user_data_cache(kind, data, user_id=user_id)

def build_upload_signature(file_name, file_bytes):
    data = file_bytes if isinstance(file_bytes, (bytes, bytearray)) else b""
    ext = Path(file_name or "").suffix.lower()
    digest = hashlib.sha256(data).hexdigest()[:16]
    return f"{ext}:{len(data)}:{digest}"

def make_uploaded_file_from_bytes(file_name, file_bytes):
    proxy = io.BytesIO(file_bytes if isinstance(file_bytes, (bytes, bytearray)) else b"")
    proxy.name = file_name or "uploaded.bin"
    return proxy

def _prewarm_cache_key(kind, signature):
    return f"{kind}:{signature}"

def get_generation_prewarm_text(kind, signature):
    key = _prewarm_cache_key(kind, signature)
    return st.session_state.get("generation_prewarm_cache", {}).get(key)

def set_generation_prewarm_text(kind, signature, text):
    key = _prewarm_cache_key(kind, signature)
    cache = st.session_state.get("generation_prewarm_cache", {})
    cache[key] = text
    st.session_state["generation_prewarm_cache"] = cache
    errors = st.session_state.get("generation_prewarm_errors", {})
    if key in errors:
        del errors[key]
    st.session_state["generation_prewarm_errors"] = errors
    return text

def get_generation_prewarm_error(kind, signature):
    key = _prewarm_cache_key(kind, signature)
    return st.session_state.get("generation_prewarm_errors", {}).get(key)

def set_generation_prewarm_error(kind, signature, error_text):
    key = _prewarm_cache_key(kind, signature)
    errors = st.session_state.get("generation_prewarm_errors", {})
    errors[key] = error_text
    st.session_state["generation_prewarm_errors"] = errors

def clear_generation_prewarm_error(kind, signature):
    key = _prewarm_cache_key(kind, signature)
    errors = st.session_state.get("generation_prewarm_errors", {})
    if key in errors:
        del errors[key]
    st.session_state["generation_prewarm_errors"] = errors

@st.cache_resource(show_spinner=False)
def get_generation_executor():
    return concurrent.futures.ThreadPoolExecutor(max_workers=2)

def get_generation_runtime_context():
    return {
        "gemini_model_id": get_gemini_model_id(),
        "audit_user_id": get_current_user_id(),
    }

def estimate_generation_runtime_minutes(total_bytes, num_files, num_items, has_style_file=False):
    mb = float(total_bytes or 0) / (1024.0 * 1024.0)
    files = max(1, int(num_files or 1))
    requested = max(1, int(num_items or 1))
    base = 0.8 + (0.28 * files) + (0.055 * requested) + (0.07 * mb)
    if requested >= 20:
        base += 0.8
    if requested >= 30:
        base += 0.8
    if has_style_file:
        base += 0.4
    return max(1.0, round(base, 1))

def start_generation_async_job(
    raw_text,
    mode,
    ai_model,
    num_items,
    chunk_size,
    overlap,
    api_key,
    openai_api_key,
    style_text,
    subject,
    unit,
    resolved_flavor=None,
    mix_basic_ratio=70,
    runtime_context=None,
):
    context = runtime_context if isinstance(runtime_context, dict) else {}
    executor = get_generation_executor()
    future = executor.submit(
        generate_content_in_chunks,
        raw_text,
        mode,
        ai_model,
        num_items,
        chunk_size,
        overlap,
        api_key,
        openai_api_key,
        style_text,
        False,
        context.get("gemini_model_id"),
        context.get("audit_user_id"),
        resolved_flavor,
        mix_basic_ratio,
    )
    return {
        "id": str(uuid.uuid4()),
        "status": "running",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "future": future,
        "mode": mode,
        "subject": subject,
        "unit": unit,
        "num_items": int(num_items),
        "resolved_flavor": resolved_flavor or "",
        "mix_basic_ratio": int(mix_basic_ratio or 70),
    }

def update_generation_async_job_state(job):
    if not isinstance(job, dict):
        return None
    status = job.get("status")
    if status in ("done", "error", "cancelled"):
        return job
    future = job.get("future")
    if future is None:
        job["status"] = "error"
        job["error"] = "백그라운드 작업 객체를 찾을 수 없습니다."
        job["completed_at"] = datetime.now(timezone.utc).isoformat()
        return job
    if not future.done():
        job["status"] = "running"
        return job
    try:
        result = future.result()
        job["result"] = result if isinstance(result, list) else []
        job["status"] = "done"
    except Exception as e:
        job["status"] = "error"
        job["error"] = str(e)
    job["completed_at"] = datetime.now(timezone.utc).isoformat()
    return job

def load_generation_queue_items():
    settings = load_user_settings()
    items = settings.get("generation_queue_v1")
    if not isinstance(items, list):
        return []
    out = []
    for item in items:
        if isinstance(item, dict) and item.get("id"):
            out.append(item)
    return out

def save_generation_queue_items(items):
    settings = load_user_settings()
    settings["generation_queue_v1"] = items if isinstance(items, list) else []
    return save_user_settings(settings)

def build_generation_queue_item(
    source_name,
    source_signature,
    raw_text,
    style_text,
    flavor_choice,
    resolved_flavor,
    mix_basic_ratio,
    mode,
    num_items,
    subject,
    unit,
    ai_model,
    chunk_size,
    overlap,
    quality_filter,
    min_length,
):
    return {
        "id": str(uuid.uuid4()),
        "source_name": source_name,
        "source_signature": source_signature or "",
        "status": "queued",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "raw_text": raw_text,
        "style_text": style_text or "",
        "flavor_choice": str(flavor_choice or ""),
        "resolved_flavor": str(resolved_flavor or ""),
        "mix_basic_ratio": int(mix_basic_ratio or 70),
        "mode": mode,
        "num_items": int(num_items),
        "subject": subject,
        "unit": unit,
        "ai_model": ai_model,
        "chunk_size": int(chunk_size),
        "overlap": int(overlap),
        "quality_filter": bool(quality_filter),
        "min_length": int(min_length),
    }

def is_duplicate_generation_queue_item(
    queue_items,
    source_signature,
    flavor_choice,
    mode,
    num_items,
    subject,
    unit,
):
    items = queue_items if isinstance(queue_items, list) else []
    sig = str(source_signature or "")
    for item in items:
        if item.get("status") not in {"queued", "running"}:
            continue
        if str(item.get("source_signature") or "") != sig:
            continue
        if str(item.get("mode") or "") != str(mode or ""):
            continue
        if str(item.get("flavor_choice") or "") != str(flavor_choice or ""):
            continue
        if int(item.get("num_items", 0)) != int(num_items or 0):
            continue
        if str(item.get("subject") or "") != str(subject or ""):
            continue
        if str(item.get("unit") or "") != str(unit or ""):
            continue
        return True
    return False

def remove_generation_queue_job(queue_items, job_id):
    items = queue_items if isinstance(queue_items, list) else []
    out = [x for x in items if str(x.get("id")) != str(job_id)]
    return (len(out) < len(items)), out

def _drop_generation_job_payload(item):
    if not isinstance(item, dict):
        return item
    slim = dict(item)
    slim.pop("raw_text", None)
    slim.pop("style_text", None)
    slim.pop("result", None)
    return slim

def start_next_generation_queue_job_if_idle(queue_items, api_key=None, openai_api_key=None, runtime_context=None):
    items = queue_items if isinstance(queue_items, list) else []
    active = st.session_state.get("generation_async_job")
    if isinstance(active, dict) and active.get("status") == "running":
        return items, False
    for idx, item in enumerate(items):
        if item.get("status") != "queued":
            continue
        job = start_generation_async_job(
            raw_text=item.get("raw_text", ""),
            mode=item.get("mode", MODE_MCQ),
            ai_model=item.get("ai_model", st.session_state.get("ai_model", "🔵 Google Gemini")),
            num_items=int(item.get("num_items", 10)),
            chunk_size=int(item.get("chunk_size", 8000)),
            overlap=int(item.get("overlap", 500)),
            api_key=api_key,
            openai_api_key=openai_api_key,
            style_text=item.get("style_text", ""),
            subject=item.get("subject", "General"),
            unit=item.get("unit", "미분류"),
            resolved_flavor=item.get("resolved_flavor", ""),
            mix_basic_ratio=int(item.get("mix_basic_ratio", 70)),
            runtime_context=runtime_context,
        )
        job["queue_id"] = item.get("id")
        st.session_state["generation_async_job"] = job
        item["status"] = "running"
        item["started_at"] = datetime.now(timezone.utc).isoformat()
        items[idx] = item
        return items, True
    return items, False

def reconcile_generation_queue_with_async(queue_items, default_quality_filter=True, default_min_length=30):
    items = queue_items if isinstance(queue_items, list) else []
    notices = []
    async_job = st.session_state.get("generation_async_job")
    if not isinstance(async_job, dict):
        return items, notices
    async_job = update_generation_async_job_state(async_job)
    st.session_state["generation_async_job"] = async_job
    status = async_job.get("status")
    if status == "running":
        return items, notices

    queue_id = str(async_job.get("queue_id") or "")
    target_idx = -1
    for idx, item in enumerate(items):
        if str(item.get("id")) == queue_id:
            target_idx = idx
            break
    if target_idx < 0:
        st.session_state["generation_async_job"] = None
        return items, notices

    target = dict(items[target_idx])
    target["completed_at"] = datetime.now(timezone.utc).isoformat()

    if status == "done":
        result = async_job.get("result") or []
        if result and isinstance(result, list):
            saved_count = add_questions_to_bank(
                result,
                target.get("mode", MODE_MCQ),
                target.get("subject", "General"),
                target.get("unit", "미분류"),
                quality_filter=bool(target.get("quality_filter", default_quality_filter)),
                min_length=int(target.get("min_length", default_min_length)),
            )
            target["status"] = "done"
            target["result_count"] = len(result)
            target["saved_count"] = int(saved_count)
            dropped = max(0, int(target["result_count"]) - int(saved_count))
            notices.append(
                f"생성 완료: {target.get('source_name', '')} "
                f"(요청 {target['result_count']}개 / 저장 {saved_count}개 / 중복·필터 제외 {dropped}개)"
            )
        else:
            target["status"] = "failed"
            target["error"] = "생성 결과가 비어 있습니다."
            notices.append(f"생성 실패: {target.get('source_name', '')} (결과 없음)")
    elif status == "cancelled":
        target["status"] = "cancelled"
        target["error"] = async_job.get("error", "사용자 취소")
        notices.append(f"작업 취소: {target.get('source_name', '')}")
    else:
        target["status"] = "failed"
        target["error"] = async_job.get("error", "알 수 없는 오류")
        notices.append(f"생성 실패: {target.get('source_name', '')}")

    items[target_idx] = _drop_generation_job_payload(target)
    st.session_state["generation_async_job"] = None
    return items, notices

def revive_stale_running_queue_items(queue_items):
    items = queue_items if isinstance(queue_items, list) else []
    async_job = st.session_state.get("generation_async_job")
    if isinstance(async_job, dict) and async_job.get("status") == "running":
        return items, False
    changed = False
    for idx, item in enumerate(items):
        if item.get("status") == "running":
            restored = dict(item)
            restored["status"] = "queued"
            restored.pop("started_at", None)
            items[idx] = restored
            changed = True
    return items, changed

# ============================================================================
# JSON 데이터 관리 함수
# ============================================================================
def load_auth_users():
    if os.path.exists(AUTH_USERS_FILE):
        try:
            with open(AUTH_USERS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                return data if isinstance(data, dict) else {}
        except Exception:
            return {}
    return {}

def save_auth_users(data):
    with open(AUTH_USERS_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def _hash_password(password, salt_hex):
    raw = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), bytes.fromhex(salt_hex), 120000)
    return raw.hex()

def register_user_account(user_id, password):
    if is_supabase_required() and not is_supabase_enabled():
        return False, "Supabase 설정이 필요합니다. 운영자에게 문의하세요."
    if is_supabase_enabled():
        return supabase_sign_up(user_id, password)
    uid = sanitize_user_id(user_id)
    if uid == "guest":
        return False, "아이디를 입력해주세요."
    if len(password or "") < 6:
        return False, "비밀번호는 6자 이상이어야 합니다."
    users = load_auth_users()
    if uid in users:
        return False, "이미 존재하는 아이디입니다."
    salt_hex = os.urandom(16).hex()
    users[uid] = {
        "salt": salt_hex,
        "password_hash": _hash_password(password, salt_hex),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    save_auth_users(users)
    get_user_data_dir(uid)
    return True, "회원가입이 완료되었습니다."

def authenticate_user_account(user_id, password):
    if is_supabase_required() and not is_supabase_enabled():
        return False, "Supabase 설정이 필요합니다. 운영자에게 문의하세요."
    if is_supabase_enabled():
        ok, payload = supabase_sign_in(user_id, password)
        if not ok:
            return False, payload
        st.session_state.auth_access_token = payload["access_token"]
        st.session_state.auth_email = payload["email"]
        return True, payload["user_id"]
    uid = sanitize_user_id(user_id)
    users = load_auth_users()
    row = users.get(uid)
    if not isinstance(row, dict):
        return False, "아이디 또는 비밀번호가 올바르지 않습니다."
    salt_hex = row.get("salt", "")
    expected = row.get("password_hash", "")
    if not salt_hex or not expected:
        return False, "계정 정보가 손상되었습니다."
    current = _hash_password(password or "", salt_hex)
    if current != expected:
        return False, "아이디 또는 비밀번호가 올바르지 않습니다."
    return True, uid

def load_questions(user_id=None) -> dict:
    """questions.json 파일 로드"""
    cached = _get_user_data_cache("questions", user_id=user_id)
    if cached is not None:
        return ensure_question_ids(cached)
    if user_id is None and is_supabase_required():
        if use_remote_user_store():
            bundle = load_remote_bundle()
            if bundle is not None:
                data = ensure_question_ids(bundle.get("questions", {"text": [], "cloze": []}))
                return _set_user_data_cache("questions", data, user_id=user_id)
            notify_remote_store_failure("⚠️ Supabase에서 문항 데이터를 불러오지 못했습니다.")
        return _set_user_data_cache("questions", {"text": [], "cloze": []}, user_id=user_id)
    if user_id is None and use_remote_user_store():
        bundle = load_remote_bundle()
        if bundle is not None:
            data = ensure_question_ids(bundle.get("questions", {"text": [], "cloze": []}))
            return _set_user_data_cache("questions", data, user_id=user_id)
    question_bank_file = get_question_bank_file(user_id)
    data = load_json_file(question_bank_file, {"text": [], "cloze": []})
    # 마이그레이션: 기존 형식 확인 및 필요시 변환
    if data and isinstance(data.get("text"), list) and len(data.get("text", [])) > 0:
        first = data["text"][0]
        if isinstance(first, dict) and "content" in first and "type" not in first:
            # 기존 형식 (content 필드) -> 새 형식으로 마이그레이션
            migrate_old_format(data, user_id=user_id)
            return load_questions(user_id=user_id)  # 다시 로드
    data = ensure_question_ids(data)
    return _set_user_data_cache("questions", data, user_id=user_id)

def migrate_old_format(data: dict, user_id=None):
    """기존 형식의 questions.json을 새 형식으로 마이그레이션"""
    try:
        migrated_text = []
        migrated_cloze = []
        
        for item in data.get("text", []):
            if isinstance(item, dict) and "content" in item:
                # 기존 형식에서 파싱
                parsed = extract_mcq_components(item["content"])
                if parsed:
                    parsed["subject"] = item.get("subject", "General")
                    parsed["date_added"] = item.get("date_added", datetime.now().isoformat())
                    migrated_text.append(parsed)
        
        for item in data.get("cloze", []):
            if isinstance(item, dict) and "content" in item:
                # Cloze 기존 형식 파싱
                content = item["content"]
                if '{{c1::' in content:
                    m = re.search(r'\{\{c1::(.+?)\}\}', content)
                    if m:
                        answer = m.group(1).strip()
                        front = re.sub(r'\{\{c1::.+?\}\}', '____', content)
                        migrated_cloze.append({
                            "type": "cloze",
                            "front": front,
                            "answer": answer,
                            "explanation": "",
                            "subject": item.get("subject", "General"),
                            "date_added": item.get("date_added", datetime.now().isoformat())
                        })
        
        # 새 형식으로 저장
        data["text"] = migrated_text
        data["cloze"] = migrated_cloze
        save_questions(data, user_id=user_id)
        
        import sys
        print(f"[MIGRATION] {len(migrated_text)}개 MCQ, {len(migrated_cloze)}개 Cloze 마이그레이션 완료", file=sys.stderr)
    except Exception as e:
        import sys
        print(f"[MIGRATION ERROR] {str(e)}", file=sys.stderr)

def save_questions(data: dict, user_id=None):
    """questions.json 파일 저장"""
    if user_id is None and is_supabase_required():
        if not use_remote_user_store():
            notify_remote_store_failure("⚠️ Supabase 로그인 상태가 아니어서 저장할 수 없습니다.")
            return False
        bundle = load_remote_bundle() or _default_remote_bundle()
        bundle["questions"] = data
        if save_remote_bundle(bundle):
            _set_user_data_cache("questions", data, user_id=user_id)
            return True
        notify_remote_store_failure("⚠️ Supabase 저장 실패로 문항 저장이 취소되었습니다.")
        return False
    if user_id is None and use_remote_user_store():
        bundle = load_remote_bundle() or _default_remote_bundle()
        bundle["questions"] = data
        if save_remote_bundle(bundle):
            _set_user_data_cache("questions", data, user_id=user_id)
            return True
    question_bank_file = get_question_bank_file(user_id)
    if not save_json_file(question_bank_file, data):
        return False
    _set_user_data_cache("questions", data, user_id=user_id)
    return True

def load_exam_history(user_id=None):
    cached = _get_user_data_cache("exam_history", user_id=user_id)
    if cached is not None:
        return cached
    if user_id is None and is_supabase_required():
        if use_remote_user_store():
            bundle = load_remote_bundle()
            if bundle is not None:
                data = bundle.get("exam_history", [])
                if not isinstance(data, list):
                    data = []
                return _set_user_data_cache("exam_history", data, user_id=user_id)
            notify_remote_store_failure("⚠️ Supabase에서 시험 기록을 불러오지 못했습니다.")
        return _set_user_data_cache("exam_history", [], user_id=user_id)
    if user_id is None and use_remote_user_store():
        bundle = load_remote_bundle()
        if bundle is not None:
            data = bundle.get("exam_history", [])
            if not isinstance(data, list):
                data = []
            return _set_user_data_cache("exam_history", data, user_id=user_id)
    exam_history_file = get_exam_history_file(user_id)
    data = load_json_file(exam_history_file, [])
    if not isinstance(data, list):
        data = []
    return _set_user_data_cache("exam_history", data, user_id=user_id)

def save_exam_history(items, user_id=None):
    if user_id is None and is_supabase_required():
        if not use_remote_user_store():
            notify_remote_store_failure("⚠️ Supabase 로그인 상태가 아니어서 시험 기록을 저장할 수 없습니다.")
            return False
        bundle = load_remote_bundle() or _default_remote_bundle()
        bundle["exam_history"] = items if isinstance(items, list) else []
        if save_remote_bundle(bundle):
            _set_user_data_cache("exam_history", bundle["exam_history"], user_id=user_id)
            return True
        notify_remote_store_failure("⚠️ Supabase 저장 실패로 시험 기록 저장이 취소되었습니다.")
        return False
    if user_id is None and use_remote_user_store():
        bundle = load_remote_bundle() or _default_remote_bundle()
        bundle["exam_history"] = items if isinstance(items, list) else []
        if save_remote_bundle(bundle):
            _set_user_data_cache("exam_history", bundle["exam_history"], user_id=user_id)
            return True
    exam_history_file = get_exam_history_file(user_id)
    normalized = items if isinstance(items, list) else []
    if not save_json_file(exam_history_file, normalized):
        return False
    _set_user_data_cache("exam_history", normalized, user_id=user_id)
    return True

def add_exam_history(session, user_id=None):
    history = load_exam_history(user_id=user_id)
    history.insert(0, session)
    save_exam_history(history[:200], user_id=user_id)
    return history

def clear_question_bank(mode="all", user_id=None):
    data = load_questions(user_id=user_id)
    if mode == "mcq":
        data["text"] = []
    elif mode == "cloze":
        data["cloze"] = []
    else:
        data = {"text": [], "cloze": []}
    save_questions(data, user_id=user_id)
    return data

def clear_exam_history(user_id=None):
    save_exam_history([], user_id=user_id)

def load_user_settings(user_id=None):
    cached = _get_user_data_cache("user_settings", user_id=user_id)
    if cached is not None:
        return cached
    if user_id is None and is_supabase_required():
        if use_remote_user_store():
            bundle = load_remote_bundle()
            if bundle is not None:
                data = bundle.get("user_settings", {})
                if not isinstance(data, dict):
                    data = {}
                return _set_user_data_cache("user_settings", data, user_id=user_id)
            notify_remote_store_failure("⚠️ Supabase에서 사용자 설정을 불러오지 못했습니다.")
        return _set_user_data_cache("user_settings", {}, user_id=user_id)
    if user_id is None and use_remote_user_store():
        bundle = load_remote_bundle()
        if bundle is not None:
            data = bundle.get("user_settings", {})
            if not isinstance(data, dict):
                data = {}
            return _set_user_data_cache("user_settings", data, user_id=user_id)
    user_settings_file = get_user_settings_file(user_id)
    data = load_json_file(user_settings_file, {})
    if not isinstance(data, dict):
        data = {}
    return _set_user_data_cache("user_settings", data, user_id=user_id)

def save_user_settings(data, user_id=None):
    if user_id is None and is_supabase_required():
        if not use_remote_user_store():
            notify_remote_store_failure("⚠️ Supabase 로그인 상태가 아니어서 설정을 저장할 수 없습니다.")
            return False
        bundle = load_remote_bundle() or _default_remote_bundle()
        bundle["user_settings"] = data if isinstance(data, dict) else {}
        if save_remote_bundle(bundle):
            _set_user_data_cache("user_settings", bundle["user_settings"], user_id=user_id)
            return True
        notify_remote_store_failure("⚠️ Supabase 저장 실패로 설정 저장이 취소되었습니다.")
        return False
    if user_id is None and use_remote_user_store():
        bundle = load_remote_bundle() or _default_remote_bundle()
        bundle["user_settings"] = data if isinstance(data, dict) else {}
        if save_remote_bundle(bundle):
            _set_user_data_cache("user_settings", bundle["user_settings"], user_id=user_id)
            return True
    user_settings_file = get_user_settings_file(user_id)
    normalized = data if isinstance(data, dict) else {}
    if not save_json_file(user_settings_file, normalized):
        return False
    _set_user_data_cache("user_settings", normalized, user_id=user_id)
    return True

def load_fsrs_settings():
    data = load_user_settings()
    default = {
        "desired_retention": 0.9,
        "learning_steps": [1, 10],
        "relearning_steps": [10],
        "maximum_interval": 36500,
        "enable_fuzzing": True,
        "parameters": list(FSRS_DEFAULT_PARAMETERS),
    }
    fsrs = data.get("fsrs_settings")
    if not isinstance(fsrs, dict):
        return default
    merged = {**default, **fsrs}
    return merged

def save_fsrs_settings(settings):
    data = load_user_settings()
    data["fsrs_settings"] = settings
    save_user_settings(data)

def get_gemini_model_id():
    return st.session_state.get("gemini_model_id") or "gemini-2.5-flash"

def _steps_to_timedelta(steps):
    out = []
    for s in steps:
        try:
            val = int(s)
            if val > 0:
                out.append(timedelta(minutes=val))
        except Exception:
            continue
    return tuple(out)

def get_fsrs_scheduler():
    if not FSRS_AVAILABLE:
        return None
    settings = load_fsrs_settings()
    params = settings.get("parameters") or list(FSRS_DEFAULT_PARAMETERS)
    try:
        return Scheduler(
            parameters=params,
            desired_retention=float(settings.get("desired_retention", 0.9)),
            learning_steps=_steps_to_timedelta(settings.get("learning_steps", [1, 10])),
            relearning_steps=_steps_to_timedelta(settings.get("relearning_steps", [10])),
            maximum_interval=int(settings.get("maximum_interval", 36500)),
            enable_fuzzing=bool(settings.get("enable_fuzzing", True)),
        )
    except Exception:
        return Scheduler()

# FSRS settings -> session state (initialize once)
if "fsrs_settings_initialized" not in st.session_state:
    _fsrs_settings = load_fsrs_settings()
    st.session_state.fsrs_settings_initialized = True
    st.session_state.fsrs_desired_retention = _fsrs_settings.get("desired_retention", 0.9)
    st.session_state.fsrs_learning_steps_text = ",".join(map(str, _fsrs_settings.get("learning_steps", [1, 10])))
    st.session_state.fsrs_relearning_steps_text = ",".join(map(str, _fsrs_settings.get("relearning_steps", [10])))
    st.session_state.fsrs_max_interval = _fsrs_settings.get("maximum_interval", 36500)
    st.session_state.fsrs_enable_fuzzing = _fsrs_settings.get("enable_fuzzing", True)
    st.session_state.fsrs_params_text = json.dumps(_fsrs_settings.get("parameters", list(FSRS_DEFAULT_PARAMETERS)))

def apply_profile_settings(profile_name):
    data = load_user_settings()
    prof = data.get(profile_name)
    if not isinstance(prof, dict):
        return False

    default_bins = [0, 1, 3, 6, 10]
    default_colors = ["#ffffff", "#d7f3f0", "#b2e9e3", "#7fd6cc", "#4fc1b6", "#1f8e86"]

    bins = prof.get("heatmap_bins")
    if isinstance(bins, list) and len(bins) >= 5:
        try:
            parsed_bins = [int(bins[i]) for i in range(5)]
            if parsed_bins[0] != 0:
                parsed_bins[0] = 0
            valid = all(parsed_bins[i] < parsed_bins[i + 1] for i in range(4))
            st.session_state.heatmap_bins = parsed_bins if valid else default_bins
        except Exception:
            st.session_state.heatmap_bins = default_bins
    else:
        st.session_state.heatmap_bins = default_bins

    colors = prof.get("heatmap_colors")
    if isinstance(colors, list) and len(colors) >= 6:
        normalized_colors = []
        for i in range(6):
            color = colors[i]
            if isinstance(color, str) and re.match(r"^#[0-9a-fA-F]{6}$", color):
                normalized_colors.append(color)
            else:
                normalized_colors.append(default_colors[i])
        st.session_state.heatmap_colors = normalized_colors
    else:
        st.session_state.heatmap_colors = default_colors

    st.session_state.select_placeholder_exam = str(
        prof.get("select_placeholder_exam", st.session_state.select_placeholder_exam)
    )
    st.session_state.select_placeholder_study = str(
        prof.get("select_placeholder_study", st.session_state.select_placeholder_study)
    )
    return True

def persist_profile_settings(profile_name):
    data = load_user_settings()
    safe_bins = st.session_state.heatmap_bins
    safe_colors = st.session_state.heatmap_colors
    if not (isinstance(safe_bins, list) and len(safe_bins) >= 5):
        safe_bins = [0, 1, 3, 6, 10]
    if not (isinstance(safe_colors, list) and len(safe_colors) >= 6):
        safe_colors = ["#ffffff", "#d7f3f0", "#b2e9e3", "#7fd6cc", "#4fc1b6", "#1f8e86"]
    data[profile_name] = {
        "heatmap_bins": safe_bins[:5],
        "heatmap_colors": safe_colors[:6],
        "select_placeholder_exam": st.session_state.select_placeholder_exam,
        "select_placeholder_study": st.session_state.select_placeholder_study,
    }
    save_user_settings(data)

def ensure_question_ids(data: dict) -> dict:
    """모든 문항에 고유 ID 부여"""
    updated = False
    for item in data.get("text", []) + data.get("cloze", []):
        if isinstance(item, dict) and "id" not in item:
            item["id"] = str(uuid.uuid4())
            updated = True
    if updated:
        save_questions(data)
    return data

def _normalize_text_for_dedupe(value):
    return " ".join(str(value or "").strip().lower().split())

def build_question_dedupe_key(item, mode=MODE_MCQ):
    if not isinstance(item, dict):
        return ""
    if mode == MODE_MCQ:
        stem = _normalize_text_for_dedupe(item.get("problem") or item.get("front"))
        options = item.get("options") or []
        normalized_options = "|".join(_normalize_text_for_dedupe(opt) for opt in options)
        answer = str(item.get("answer") or "").strip()
        return f"mcq::{stem}::{normalized_options}::{answer}"
    front = _normalize_text_for_dedupe(item.get("front") or item.get("problem"))
    answer = _normalize_text_for_dedupe(item.get("answer"))
    response_type = str(item.get("response_type") or "").strip().lower()
    return f"cloze::{response_type}::{front}::{answer}"

def add_questions_to_bank(questions_data, mode, subject="General", unit="미분류", quality_filter=True, min_length=20, batch_id=None):
    """생성된 문제를 question bank에 추가 (구조화된 JSON 형식)
    
    Args:
        questions_data: 다음 중 하나
            - 구조화된 dict의 리스트: [{"problem": ..., "options": [...], "answer": 1, "explanation": ...}]
            - 문자열: 기존 호환성을 위함
        mode: 모드 (객관식/빈칸/단답형/서술형)
        subject: 과목명
        quality_filter: 품질 필터링 여부
        min_length: 최소 길이
    
    Returns:
        추가된 문제 개수
    """
    bank = load_questions()
    
    # 문자열이면 파싱 (기존 호환성)
    if isinstance(questions_data, str):
        parsed_questions = parse_generated_text_to_structured(questions_data, mode)
    else:
        parsed_questions = questions_data if isinstance(questions_data, list) else [questions_data]
    
    added_count = 0
    target_key = "text" if mode == MODE_MCQ else "cloze"
    existing_keys = set()
    for existing_item in bank.get(target_key, []):
        key = build_question_dedupe_key(existing_item, mode)
        if key:
            existing_keys.add(key)
    if not batch_id:
        batch_id = datetime.now().strftime("%Y%m%d-%H%M%S") + "-" + uuid.uuid4().hex[:6]

    for q_data in parsed_questions:
        if not q_data:
            continue
        
        # 품질 필터링
        if quality_filter:
            if mode == MODE_MCQ:
                problem_text = q_data.get("problem", "")
                if len(problem_text) < min_length:
                    continue
            else:
                front_text = q_data.get("front", "")
                if len(front_text) < min_length:
                    continue
        
        # 메타데이터 추가
        q_data["subject"] = q_data.get("subject") or subject
        q_data["unit"] = q_data.get("unit") or unit
        q_data["date_added"] = datetime.now().isoformat()
        if "id" not in q_data:
            q_data["id"] = str(uuid.uuid4())
        q_data["batch_id"] = q_data.get("batch_id") or batch_id
        
        key = build_question_dedupe_key(q_data, mode)
        if key and key in existing_keys:
            continue
        if key:
            existing_keys.add(key)

        bank[target_key].append(q_data)
        
        added_count += 1
    
    save_questions(bank)
    return added_count

def add_questions_to_bank_auto(items, subject="General", unit="미분류", quality_filter=True, min_length=20, batch_id=None):
    """MCQ/Cloze 혼합 입력 자동 분류 후 저장"""
    if not batch_id:
        batch_id = datetime.now().strftime("%Y%m%d-%H%M%S") + "-" + uuid.uuid4().hex[:6]
    mcq_items = []
    cloze_items = []
    for item in items:
        if not isinstance(item, dict):
            continue
        item["subject"] = item.get("subject") or subject
        item["unit"] = item.get("unit") or unit
        item["batch_id"] = item.get("batch_id") or batch_id
        if item.get("type") == "cloze":
            cloze_items.append(item)
        else:
            mcq_items.append(item)
    added = 0
    if mcq_items:
        added += add_questions_to_bank(mcq_items, MODE_MCQ, subject, unit, quality_filter, min_length, batch_id=batch_id)
    if cloze_items:
        added += add_questions_to_bank(cloze_items, MODE_CLOZE, subject, unit, quality_filter, min_length, batch_id=batch_id)
    return added


def parse_free_response_items(text, response_type="short"):
    items = []

    def append_item(obj):
        if not isinstance(obj, dict):
            return
        front = (obj.get("front") or obj.get("question") or obj.get("problem") or "").strip()
        answer = (obj.get("answer") or obj.get("reference_answer") or obj.get("model_answer") or "").strip()
        explanation = (obj.get("explanation") or obj.get("rationale") or "").strip()
        if front and answer:
            items.append({
                "type": "cloze",
                "response_type": response_type,
                "front": front,
                "answer": answer,
                "explanation": explanation,
            })

    parsed_json = _parse_json_from_text(text)
    if isinstance(parsed_json, dict):
        parsed_json = [parsed_json]
    if isinstance(parsed_json, list):
        for obj in parsed_json:
            append_item(obj)
        if items:
            return items

    blocks = re.split(r"\n-{3,}\n", text)
    for block in blocks:
        line = block.strip()
        if not line:
            continue
        if "\t" in line:
            cols = [c.strip() for c in line.split("\t")]
            if len(cols) >= 2 and cols[0] and cols[1]:
                items.append({
                    "type": "cloze",
                    "response_type": response_type,
                    "front": cols[0],
                    "answer": cols[1],
                    "explanation": cols[2] if len(cols) > 2 else "",
                })
            continue

        lines = [x.strip() for x in line.splitlines() if x.strip()]
        if len(lines) < 2:
            continue
        front = re.sub(r"^(문항|문제|Q)\s*[:：]\s*", "", lines[0], flags=re.IGNORECASE).strip()
        answer = re.sub(r"^(정답|답|A)\s*[:：]\s*", "", lines[1], flags=re.IGNORECASE).strip()
        explanation = ""
        if len(lines) > 2:
            explanation = re.sub(r"^(해설|설명)\s*[:：]\s*", "", "\n".join(lines[2:]), flags=re.IGNORECASE).strip()
        if front and answer:
            items.append({
                "type": "cloze",
                "response_type": response_type,
                "front": front,
                "answer": answer,
                "explanation": explanation,
            })
    return items


def parse_generated_text_to_structured(text, mode):
    """생성된 텍스트를 구조화된 형식으로 파싱
    
    Returns:
        구조화된 dict의 리스트
    """
    results = []
    mode_mcq = globals().get("MODE_MCQ", "📝 객관식 문제 (Case Study)")
    mode_cloze = globals().get("MODE_CLOZE", "🧩 빈칸 뚫기 (Anki Cloze)")
    mode_short = globals().get("MODE_SHORT", "🧠 단답형 문제")
    mode_essay = globals().get("MODE_ESSAY", "🧾 서술형 문제")
    
    if mode == mode_mcq:
        # 1) JSON 형식 우선 파싱 (Gemini/OpenAI JSON 대응)
        # 전체 텍스트가 JSON 배열/객체인 경우
        try:
            stripped = text.strip()
            if stripped.startswith("{") or stripped.startswith("["):
                parsed = json.loads(stripped)
                if isinstance(parsed, dict):
                    parsed = [parsed]
                if isinstance(parsed, list):
                    for item in parsed:
                        norm = normalize_mcq_item(item)
                        if norm:
                            results.append(norm)
                    if results:
                        return results
        except Exception:
            pass

        # 복수 JSON 블록이 섞여 있는 경우를 탐지
        try:
            decoder = json.JSONDecoder()
            idx = 0
            stripped = text.strip()
            while idx < len(stripped):
                if stripped[idx] not in "{[":
                    idx += 1
                    continue
                try:
                    obj, next_idx = decoder.raw_decode(stripped[idx:])
                    idx += next_idx
                    if isinstance(obj, dict):
                        obj = [obj]
                    if isinstance(obj, list):
                        for item in obj:
                            norm = normalize_mcq_item(item)
                            if norm:
                                results.append(norm)
                except Exception:
                    idx += 1
            if results:
                return results
        except Exception:
            pass

        # TSV 또는 '---' 구분자로 된 MCQ 파싱
        items = text.split("\n---\n")
        
        for item in items:
            item = item.strip()
            if not item or len(item) < 50:
                continue
            
            # TSV 형식: problem_text\texplanation
            parts = item.split('\t')
            problem_part = parts[0].strip() if parts else ""
            explanation_part = parts[1].strip() if len(parts) > 1 else ""
            
            if not problem_part:
                continue
            
            # 정답과 선지 추출
            parsed = extract_mcq_components(problem_part)
            if parsed:
                parsed["explanation"] = explanation_part
                results.append(parsed)
    elif mode == mode_cloze:
        # Cloze 형식: 한 줄에 하나씩
        lines = text.split('\n')
        for line in lines:
            line = line.strip()
            if not line or '{{c1::' not in line:
                continue
            
            # 해설 분리
            explanation = ""
            if '\t' in line:
                line, explanation = line.split('\t', 1)
            
            # 정답 추출
            m = re.search(r'\{\{c1::(.+?)\}\}', line)
            if not m:
                continue
            
            answer = m.group(1).strip()
            front = re.sub(r'\{\{c1::.+?\}\}', '____', line)
            
            results.append({
                "type": "cloze",
                "response_type": "cloze",
                "front": front,
                "answer": answer,
                "explanation": explanation
            })
    elif mode == mode_short:
        results = parse_free_response_items(text, response_type="short")
    elif mode == mode_essay:
        results = parse_free_response_items(text, response_type="essay")
    
    return results


def extract_mcq_components(problem_text):
    """MCQ 텍스트에서 문제, 선지, 정답을 추출
    
    Returns:
        {"type": "mcq", "problem": ..., "options": [...], "answer": ..., "explanation": ""}
        또는 None (파싱 실패 시)
    """
    try:
        # 정답 추출
        ans_match = re.search(r"정답:\s*\{\{c1::([1-5①②③④⑤]+)\}\}", problem_text)
        if not ans_match:
            return None
        
        ans_str = ans_match.group(1).strip()
        circ_to_num = {'①': '1', '②': '2', '③': '3', '④': '4', '⑤': '5'}
        answer_num = int(circ_to_num.get(ans_str, ans_str))
        
        # 선지 추출: ① ... ② ... 형식
        options = []
        opt_pattern = r'(?:①|②|③|④|⑤)\s*([^①②③④⑤\n]+?)(?=(?:①|②|③|④|⑤|$))'
        matches = re.findall(opt_pattern, problem_text)
        options = [opt.strip() for opt in matches if opt.strip()]
        
        if len(options) < 3:  # 최소 3개 이상 필요
            return None
        
        # 선지를 5개로 정규화 (부족하면 채우기)
        while len(options) < 5:
            options.append(f"보기 {len(options) + 1}")
        options = options[:5]  # 5개 초과면 자르기
        
        # 문제 텍스트 정제: 정답/선지 제거 후 스템만 남기기
        problem_clean = re.sub(r'정답:\s*\{\{c1::.+?\}\}', '', problem_text).strip()
        # 선지 시작 위치 이전만 스템으로 사용
        first_opt = re.search(r'(①|②|③|④|⑤)', problem_clean)
        if first_opt:
            stem = problem_clean[:first_opt.start()].strip()
        else:
            stem = problem_clean
        stem = re.sub(r'\s+', ' ', stem)
        if not stem:
            stem = problem_clean
        
        return {
            "type": "mcq",
            "problem": stem,
            "options": options,
            "answer": answer_num,
            "explanation": ""
        }
    except Exception as e:
        import sys
        print(f"[EXTRACT ERROR] {str(e)}", file=sys.stderr)
        return None

def parse_mcq_content(q_data: dict) -> dict:
    """저장된 MCQ 데이터를 시험 표시용으로 변환
    
    Args:
        q_data: {"type": "mcq", "problem": ..., "options": [...], "answer": ..., "explanation": ...}
    
    Returns:
        {"type": "mcq", "front": ..., "problem": ..., "options": [...], "correct": ..., "explanation": ...}
    """
    stem = sanitize_mcq_problem_text(q_data.get("problem", ""))
    return {
        "type": "mcq",
        "raw": stem,
        "front": stem,
        "problem": stem,
        "options": q_data.get("options", []),
        "correct": q_data.get("answer"),  # 숫자 형식: 1-5
        "explanation": q_data.get("explanation", ""),
        "subject": q_data.get("subject"),
        "unit": q_data.get("unit"),
        "difficulty": q_data.get("difficulty"),
        "id": q_data.get("id"),
        "fsrs": q_data.get("fsrs"),
        "note": q_data.get("note", ""),
        "images": q_data.get("images", []),
        "stats": q_data.get("stats", {}),
        "bookmarked": bool(q_data.get("bookmarked", False)),
    }

def sanitize_mcq_problem_text(problem_text):
    text = re.sub(r"\s+", " ", str(problem_text or "")).strip()
    if not text:
        return ""

    # 중복 [문제] 마커가 붙는 경우 첫 문항만 유지
    second_marker = text.find("[문제]", len("[문제]"))
    if second_marker != -1:
        text = text[:second_marker].strip()

    # 물음표 뒤에 공백 없이 다른 문항이 붙은 경우(예: "...것은?TPN)...") 첫 문항으로 절단
    hard_concat = re.search(r"\?[^\s\"'”’)\]}]", text)
    if hard_concat:
        text = text[: hard_concat.start() + 1].strip()

    return text

def parse_cloze_content(q_data: dict) -> dict:
    """저장된 Cloze 데이터를 시험 표시용으로 변환
    
    Args:
        q_data: {"type": "cloze", "front": ..., "answer": ..., "explanation": ...}
    
    Returns:
        {"type": "cloze", "front": ..., "raw": ..., "answer": ..., "explanation": ...}
    """
    return {
        "type": "cloze",
        "raw": q_data.get("front", ""),
        "front": q_data.get("front", ""),
        "answer": q_data.get("answer", ""),
        "response_type": q_data.get("response_type", "cloze"),
        "explanation": q_data.get("explanation", ""),
        "subject": q_data.get("subject"),
        "unit": q_data.get("unit"),
        "difficulty": q_data.get("difficulty"),
        "id": q_data.get("id"),
        "fsrs": q_data.get("fsrs"),
        "note": q_data.get("note", ""),
        "images": q_data.get("images", []),
        "stats": q_data.get("stats", {}),
        "bookmarked": bool(q_data.get("bookmarked", False)),
    }

def get_question_stats():
    """저장된 문제 통계"""
    bank = load_questions()
    return {
        "total_text": len(bank.get("text", [])),
        "total_cloze": len(bank.get("cloze", []))
    }

def fuzzy_match(user_answer, correct_answer, threshold=0.8):
    """Cloze 답변 유사도 비교"""
    user_clean = re.sub(r'[^\w가-힣]', '', str(user_answer).lower())
    correct_clean = re.sub(r'[^\w가-힣]', '', correct_answer.lower())
    
    if user_clean == correct_clean:
        return True
    ratio = SequenceMatcher(None, user_clean, correct_clean).ratio()
    return ratio >= threshold

def calculate_quality_score(item_text, mode):
    """항목의 품질 점수 계산 (0~1.0)"""
    score = 0.4
    text = item_text.strip()
    text_len = len(text)
    
    # 길이 점수
    if 80 < text_len < 500:
        score += 0.25
    elif 50 < text_len < 700:
        score += 0.15
    
    # 형식 점수
    if mode == "📝 객관식 문제 (Case Study)":
        if "정답:" in text:
            score += 0.15
        options = len(re.findall(r"①|②|③|④|⑤", text))
        if options >= 3:
            score += 0.15
    else:  # Cloze
        if "{{c1::" in text:
            score += 0.3
    
    # 의학 용어 점수
    medical_keywords = ["증상", "진단", "치료", "질병", "검사", "수치", "질환", "증후군"]
    kw_count = sum(1 for kw in medical_keywords if kw in text)
    if kw_count >= 2:
        score += 0.15
    elif kw_count >= 1:
        score += 0.08
    
    if text.endswith((".", "。")):
        score += 0.08
    
    complex_chars = text.count(",") + text.count(";") + text.count("(")
    if 2 <= complex_chars <= 8:
        score += 0.05
    
    return min(max(score, 0.0), 1.0)

def auto_tag(item_text):
    """휴리스틱 기반 난이도/카테고리 태깅"""
    txt = item_text.lower()
    
    # 카테고리
    categories = []
    if any(k in txt for k in ["심장", "심근", "부정맥", "협심증"]):
        categories.append("cardio")
    if any(k in txt for k in ["폐", "호흡", "기관지", "천식"]):
        categories.append("pulmonary")
    if any(k in txt for k in ["신경", "뇌", "척추", "신경계"]):
        categories.append("neuro")
    if any(k in txt for k in ["암", "종양", "신생물"]):
        categories.append("oncology")
    if any(k in txt for k in ["신장", "신부전", "사구체"]):
        categories.append("nephro")
    if not categories:
        categories.append("general")
    
    # 난이도
    length = len(item_text)
    complexity = item_text.count(";") + item_text.count(",")
    if length < 150 and complexity < 3:
        difficulty = "⭐ 쉬움"
    elif length < 350 and complexity < 6:
        difficulty = "⭐⭐ 중간"
    else:
        difficulty = "⭐⭐⭐ 어려움"
    
    return difficulty, categories

def is_answer_correct(q, user_ans):
    if q.get("type") == "mcq":
        correct_choice = q.get("correct")
        return bool(correct_choice and user_ans == correct_choice)
    response_type = q.get("response_type", "cloze")
    if response_type == "essay":
        ai_grade = q.get("_ai_grade")
        return bool(isinstance(ai_grade, dict) and ai_grade.get("is_correct") is True)
    correct_text = q.get("answer")
    return bool(correct_text and isinstance(user_ans, str) and fuzzy_match(user_ans, correct_text))

def parse_iso_datetime(value):
    if not value:
        return None
    try:
        if isinstance(value, (int, float)):
            return datetime.fromtimestamp(value, tz=timezone.utc)
        if isinstance(value, str):
            v = value.replace("Z", "+00:00")
            return datetime.fromisoformat(v)
    except Exception:
        return None
    return None

def get_fsrs_report(questions, now=None):
    if not FSRS_AVAILABLE:
        return None
    check_time = now or datetime.now(timezone.utc)
    total = len(questions)
    stats = get_fsrs_stats(questions, now=check_time)
    review_count_7d = 0
    rating_counts = {"Again": 0, "Hard": 0, "Good": 0, "Easy": 0}
    intervals = []
    last_review = None
    for q in questions:
        fsrs = q.get("fsrs") or {}
        card_data = fsrs.get("card")
        if card_data:
            try:
                card = Card.from_json(card_data)
                if hasattr(card, "interval"):
                    intervals.append(float(card.interval))
            except Exception:
                pass
        # last_rating
        last_rating = fsrs.get("last_rating")
        if last_rating in rating_counts:
            rating_counts[last_rating] += 1

        # logs
        for log in fsrs.get("logs", []) or []:
            if isinstance(log, dict):
                for key in ("review_datetime", "reviewed_at", "time", "date", "review"):
                    dt = parse_iso_datetime(log.get(key))
                    if dt:
                        if dt.tzinfo is None:
                            dt = dt.replace(tzinfo=timezone.utc)
                        if (check_time - dt).days <= 7:
                            review_count_7d += 1
                        if last_review is None or dt > last_review:
                            last_review = dt
                        break
                rating = log.get("rating")
                if isinstance(rating, str) and rating in rating_counts:
                    rating_counts[rating] += 1
    avg_interval = sum(intervals) / len(intervals) if intervals else 0
    return {
        "total": total,
        "stats": stats,
        "review_count_7d": review_count_7d,
        "avg_interval": avg_interval,
        "last_review": last_review.isoformat() if last_review else None,
        "rating_counts": rating_counts,
    }

def update_question_stats(q_id, is_correct):
    bank = load_questions()
    now = datetime.now(timezone.utc).isoformat()
    for key in ("text", "cloze"):
        for item in bank.get(key, []):
            if item.get("id") == q_id:
                stats = item.get("stats") or {}
                stats["right"] = int(stats.get("right", 0))
                stats["wrong"] = int(stats.get("wrong", 0))
                if is_correct:
                    stats["right"] += 1
                else:
                    stats["wrong"] += 1
                stats["last_attempt"] = now
                history = stats.get("history") or []
                history.append({"time": now, "correct": bool(is_correct)})
                stats["history"] = history[-200:]
                item["stats"] = stats
                save_questions(bank)
                append_audit_log("grade.answer", {
                    "question_id": q_id,
                    "correct": bool(is_correct),
                    "score": 1 if is_correct else 0,
                    "grader_version": GRADER_VERSION,
                })
                return stats
    return None

def update_question_note(q_id, note_text):
    bank = load_questions()
    for key in ("text", "cloze"):
        for item in bank.get(key, []):
            if item.get("id") == q_id:
                item["note"] = note_text
                save_questions(bank)
                return True
    return False

def update_question_bookmark(q_id, bookmarked):
    bank = load_questions()
    for key in ("text", "cloze"):
        for item in bank.get(key, []):
            if item.get("id") == q_id:
                item["bookmarked"] = bool(bookmarked)
                save_questions(bank)
                return True
    return False

def update_question_by_id(q_id, patch):
    if not q_id or not isinstance(patch, dict):
        return False
    bank = load_questions()
    for key in ("text", "cloze"):
        for item in bank.get(key, []):
            if item.get("id") == q_id:
                allowed = {
                    "subject", "unit", "problem", "options", "answer", "front",
                    "explanation", "difficulty", "note", "image"
                }
                item.update({k: v for k, v in patch.items() if k in allowed})
                save_questions(bank)
                return True
    return False

def delete_mcq_by_ids(ids):
    if not ids:
        return 0
    data = load_questions()
    before = len(data.get("text", []))
    data["text"] = [q for q in data.get("text", []) if q.get("id") not in ids]
    save_questions(data)
    return before - len(data.get("text", []))

def delete_mcq_by_batch(batch_id):
    if not batch_id:
        return 0
    data = load_questions()
    before = len(data.get("text", []))
    data["text"] = [q for q in data.get("text", []) if (q.get("batch_id") or "legacy") != batch_id]
    save_questions(data)
    return before - len(data.get("text", []))

def delete_questions_by_subject_units(subject_unit_map, mode="all"):
    if not isinstance(subject_unit_map, dict) or not subject_unit_map:
        return 0

    normalized = {}
    for subj, units in subject_unit_map.items():
        subj_name = str(subj or "General")
        if not isinstance(units, (list, tuple, set)):
            continue
        cleaned = {str(u or "미분류") for u in units if str(u or "").strip()}
        if cleaned:
            normalized[subj_name] = cleaned
    if not normalized:
        return 0

    keys = ["text", "cloze"] if mode == "all" else (["text"] if mode == "mcq" else ["cloze"])
    data = load_questions()
    before = sum(len(data.get(k, [])) for k in keys)

    def _keep(item):
        subj = str(item.get("subject") or "General")
        unit = str(item.get("unit") or "미분류")
        units = normalized.get(subj)
        if not units:
            return True
        return unit not in units

    for key in keys:
        data[key] = [q for q in data.get(key, []) if _keep(q)]

    after = sum(len(data.get(k, [])) for k in keys)
    deleted = before - after
    if deleted > 0:
        save_questions(data)
    return deleted

def get_mcq_batches(questions):
    batches = {}
    for q in questions:
        b = q.get("batch_id") or "legacy"
        batches[b] = batches.get(b, 0) + 1
    return batches

def get_wrong_note_stats(questions):
    wrong_items = []
    total_wrong = 0
    for q in questions:
        stats = q.get("stats") or {}
        wrong = int(stats.get("wrong", 0))
        if wrong > 0:
            wrong_items.append(q)
            total_wrong += wrong
    return wrong_items, total_wrong

def sort_wrong_first(questions, mode="오답 횟수", weight_recent=0.7, weight_count=0.3):
    def last_wrong_time(q):
        stats = q.get("stats") or {}
        hist = stats.get("history") or []
        latest = None
        for entry in hist:
            if not isinstance(entry, dict):
                continue
            if entry.get("correct") is True:
                continue
            dt = parse_iso_datetime(entry.get("time"))
            if dt:
                if latest is None or dt > latest:
                    latest = dt
        return latest or datetime.min.replace(tzinfo=timezone.utc)

    def score(q):
        stats = q.get("stats") or {}
        wrong = int(stats.get("wrong", 0))
        right = int(stats.get("right", 0))
        total = wrong + right
        rate = wrong / total if total > 0 else 0
        if mode == "오답률":
            return (rate, wrong)
        if mode == "최근 오답":
            # 최근 오답일수록 높은 점수
            last_dt = last_wrong_time(q)
            days_since = (datetime.now(timezone.utc) - last_dt).days if last_dt else 9999
            recency_score = 1 / (1 + max(days_since, 0))
            combined = weight_recent * recency_score + weight_count * wrong
            return (combined, recency_score, wrong)
        return (wrong, rate)

    return sorted(questions, key=score, reverse=True)

def compute_recent_accuracy(questions, days=7, now=None):
    check_time = now or datetime.now(timezone.utc)
    cutoff = check_time - timedelta(days=days)
    correct = 0
    total = 0
    for q in questions:
        stats = q.get("stats") or {}
        hist = stats.get("history") or []
        for entry in hist:
            if not isinstance(entry, dict):
                continue
            dt = parse_iso_datetime(entry.get("time"))
            if not dt:
                continue
            if dt.tzinfo is None:
                dt = dt.replace(tzinfo=timezone.utc)
            if dt >= cutoff:
                total += 1
                if entry.get("correct") is True:
                    correct += 1
    accuracy = (correct / total * 100) if total > 0 else None
    return {"correct": correct, "total": total, "accuracy": accuracy}

def compute_accuracy_trend(questions, days=14, now=None):
    check_time = now or datetime.now(timezone.utc)
    start = (check_time - timedelta(days=days - 1)).date()
    buckets = {}
    for i in range(days):
        d = start + timedelta(days=i)
        buckets[d.isoformat()] = {"correct": 0, "total": 0}
    for q in questions:
        stats = q.get("stats") or {}
        hist = stats.get("history") or []
        for entry in hist:
            if not isinstance(entry, dict):
                continue
            dt = parse_iso_datetime(entry.get("time"))
            if not dt:
                continue
            if dt.tzinfo is None:
                dt = dt.replace(tzinfo=timezone.utc)
            dkey = dt.date().isoformat()
            if dkey in buckets:
                buckets[dkey]["total"] += 1
                if entry.get("correct") is True:
                    buckets[dkey]["correct"] += 1
    series = []
    for dkey in sorted(buckets.keys()):
        total = buckets[dkey]["total"]
        acc = (buckets[dkey]["correct"] / total * 100) if total > 0 else 0
        series.append({"date": dkey, "accuracy": acc})
    return series

def compute_overall_accuracy(questions):
    right = 0
    wrong = 0
    for q in questions:
        stats = q.get("stats") or {}
        right += int(stats.get("right", 0))
        wrong += int(stats.get("wrong", 0))
    total = right + wrong
    if total == 0:
        return None
    accuracy = right / total * 100
    return {"correct": right, "wrong": wrong, "total": total, "accuracy": accuracy}

def fsrs_group_report(questions, group_key, now=None):
    if not FSRS_AVAILABLE:
        return []
    check_time = now or datetime.now(timezone.utc)
    groups = {}
    for q in questions:
        key = (q.get(group_key) or "General")
        g = groups.setdefault(key, {"due": 0, "overdue": 0, "future": 0, "new": 0, "total": 0})
        g["total"] += 1
        fsrs = q.get("fsrs") or {}
        card_data = fsrs.get("card")
        if not card_data:
            g["new"] += 1
            g["due"] += 1
            continue
        try:
            card = Card.from_json(card_data)
            if card.due <= check_time:
                g["due"] += 1
                if card.due < check_time:
                    g["overdue"] += 1
            else:
                g["future"] += 1
        except Exception:
            g["due"] += 1
    rows = []
    for k, v in sorted(groups.items(), key=lambda x: x[0]):
        rows.append({"그룹": k, **v})
    return rows

def apply_mcq_shortcut(idx):
    val = (st.session_state.get(f"shortcut_{idx}") or "").strip().upper()
    if not val:
        return
    letters = ["A", "B", "C", "D", "E"]
    sel = None
    if val in letters:
        sel = letters.index(val)
    elif val.isdigit():
        n = int(val)
        if 1 <= n <= 5:
            sel = n - 1
    labels = st.session_state.get(f"labels_real_{idx}") or []
    if sel is not None and 0 <= sel < len(labels):
        st.session_state[f"q_{idx}"] = labels[sel]

def goto_prev_question():
    st.session_state.current_question_idx = max(0, st.session_state.current_question_idx - 1)

def goto_next_question():
    total = len(st.session_state.get("exam_questions", []))
    if total:
        st.session_state.current_question_idx = min(total - 1, st.session_state.current_question_idx + 1)

def finish_exam_session():
    st.session_state.exam_finished = True

def get_unique_subjects(questions):
    subjects = sorted({(q.get("subject") or "General") for q in questions})
    return subjects

def get_unit_name(q):
    return q.get("unit") or q.get("chapter") or q.get("topic") or "미분류"

def get_units_by_subject(questions):
    mapping = {}
    for q in questions:
        subj = (q.get("subject") or "General")
        unit = get_unit_name(q)
        mapping.setdefault(subj, set()).add(unit)
    return {k: sorted(v) for k, v in mapping.items()}

def filter_questions_by_subject(questions, selected_subjects):
    if not selected_subjects:
        return questions
    return [q for q in questions if (q.get("subject") or "General") in selected_subjects]

def filter_questions_by_subject_unit(questions, selected_subjects, selected_units):
    if not selected_subjects and not selected_units:
        return questions
    filtered = []
    for q in questions:
        subj = q.get("subject") or "General"
        unit = get_unit_name(q)
        if selected_subjects and subj not in selected_subjects:
            continue
        if selected_units and unit not in selected_units:
            continue
        filtered.append(q)
    return filtered

def filter_questions_by_subject_unit_hierarchy(questions, selected_subjects, unit_filter_by_subject):
    """분과 단위 계층 필터(과목별 단원 선택)"""
    if not questions:
        return []
    if not selected_subjects:
        return []
    out = []
    for q in questions:
        subj = q.get("subject") or "General"
        if subj not in selected_subjects:
            continue
        unit = get_unit_name(q)
        allowed_units = unit_filter_by_subject.get(subj)
        if not allowed_units:
            continue
        if unit in allowed_units:
            out.append(q)
    return out

def collect_subject_unit_map(questions):
    """분과별 단원 목록 생성"""
    mapping = {}
    for q in questions:
        subj = q.get("subject") or "General"
        unit = get_unit_name(q)
        mapping.setdefault(subj, set()).add(unit)
    return {k: sorted(v) for k, v in mapping.items()}

def collect_export_questions(questions, selected_subjects, unit_filter_by_subject, include_all_units=True, randomize=False, random_seed=None):
    if include_all_units:
        if selected_subjects:
            items = filter_questions_by_subject(questions, selected_subjects)
        else:
            items = list(questions)
    else:
        items = filter_questions_by_subject_unit_hierarchy(questions, selected_subjects, unit_filter_by_subject)
    out = list(items)
    deduped = []
    seen = set()
    for q in out:
        qid = str(q.get("id") or "").strip()
        if qid:
            key = f"id::{qid}"
        elif q.get("type") == "cloze":
            front = " ".join(str(q.get("front") or "").split())
            answer = " ".join(str(q.get("answer") or "").split())
            key = f"cloze::{front}::{answer}"
        else:
            stem = " ".join(str(q.get("problem") or "").split())
            options = "|".join(" ".join(str(opt).split()) for opt in (q.get("options") or []))
            answer = str(q.get("answer") or "")
            key = f"mcq::{stem}::{options}::{answer}"
        if key in seen:
            continue
        seen.add(key)
        deduped.append(q)
    out = deduped
    if randomize and len(out) > 1:
        rng = random.Random(random_seed)
        rng.shuffle(out)
    return out

def get_question_attempt_summary(question):
    stats = question.get("stats") or {}
    right = int(stats.get("right", 0))
    wrong = int(stats.get("wrong", 0))
    attempts = right + wrong
    history = stats.get("history") or []
    last = history[-1] if history else {}
    last_correct = bool(last.get("correct")) if isinstance(last, dict) and "correct" in last else None
    last_time = last.get("time") if isinstance(last, dict) else None
    return {
        "right": right,
        "wrong": wrong,
        "attempts": attempts,
        "last_correct": last_correct,
        "last_time": last_time,
    }

def select_learning_session_questions(questions, learning_mode="탐색형", num_questions=10, random_seed=None):
    items = list(questions or [])
    if not items:
        return []
    if str(learning_mode).startswith("탐색형"):
        return sorted(
            items,
            key=lambda q: (
                q.get("subject") or "General",
                get_unit_name(q),
                str(q.get("id") or ""),
                str(q.get("problem") or q.get("front") or ""),
            ),
        )
    n = min(max(1, int(num_questions or 1)), len(items))
    if n >= len(items):
        return items
    rng = random.Random(random_seed) if random_seed is not None else random
    return rng.sample(items, n)

def _exam_group_key(question, group_mode="분과+단원"):
    subject = question.get("subject") or "General"
    if group_mode == "분과":
        return subject
    unit = get_unit_name(question)
    return f"{subject}::{unit}"

def select_exam_questions_balanced(
    questions,
    num_questions,
    distribution_mode="비례(보유 문항 기준)",
    group_mode="분과+단원",
    random_seed=None,
):
    items = list(questions or [])
    if not items:
        return []
    target_n = min(max(1, int(num_questions or 1)), len(items))
    rng = random.Random(random_seed) if random_seed is not None else random

    groups = {}
    for q in items:
        key = _exam_group_key(q, group_mode=group_mode)
        groups.setdefault(key, []).append(q)
    for key in groups:
        rng.shuffle(groups[key])

    keys = list(groups.keys())
    if not keys:
        return []
    available = {k: len(groups[k]) for k in keys}

    if distribution_mode == "균등(선택 그룹 기준)":
        weights = {k: 1.0 for k in keys}
    else:
        weights = {k: float(available[k]) for k in keys}

    total_weight = sum(weights.values()) or 1.0
    raw_alloc = {k: (target_n * weights[k] / total_weight) for k in keys}
    alloc = {k: min(int(raw_alloc[k]), available[k]) for k in keys}
    remaining = target_n - sum(alloc.values())

    order = sorted(keys, key=lambda k: (raw_alloc[k] - int(raw_alloc[k]), available[k]), reverse=True)
    while remaining > 0:
        progressed = False
        for k in order:
            if alloc[k] < available[k]:
                alloc[k] += 1
                remaining -= 1
                progressed = True
                if remaining <= 0:
                    break
        if not progressed:
            break

    selected = []
    for k in keys:
        selected.extend(groups[k][: alloc[k]])

    if len(selected) < target_n:
        chosen = {id(x) for x in selected}
        leftovers = [q for q in items if id(q) not in chosen]
        need = target_n - len(selected)
        if leftovers:
            rng.shuffle(leftovers)
            selected.extend(leftovers[:need])
    return selected[:target_n]


def summarize_subject_review_status(questions):
    """분과별 복습 상태(복습대상/연체/단원 수) 요약"""
    if not questions:
        return []
    now = datetime.now(timezone.utc)

    # 오답문항(통계 기반)
    wrong_by_subject = {}
    for q in questions:
        subj = q.get("subject") or "General"
        if int((q.get("stats") or {}).get("wrong", 0)) > 0:
            wrong_by_subject[subj] = wrong_by_subject.get(subj, 0) + 1

    if FSRS_AVAILABLE:
        rows = fsrs_group_report(questions, "subject", now=now)
        out = []
        for row in rows:
            subject_name = row.get("그룹") or "General"
            out.append({
                "분과": subject_name,
                "총문항": row.get("total", 0),
                "복습대상": row.get("due", 0),
                "연체": row.get("overdue", 0),
                "미래": row.get("future", 0),
                "신규": row.get("new", 0),
                "오답문항": wrong_by_subject.get(subject_name, 0),
            })
        return sorted(out, key=lambda x: (x["복습대상"], x["총문항"]), reverse=True)

    summary = {}
    for q in questions:
        subj = q.get("subject") or "General"
        row = summary.setdefault(subj, {"분과": subj, "총문항": 0, "복습대상": 0, "연체": 0, "미래": 0, "신규": 0})
        row["총문항"] += 1

        due_at = (q.get("srs") or {}).get("due")
        try:
            if due_at:
                due_dt = datetime.fromisoformat(str(due_at).replace("Z", "+00:00"))
                if due_dt.tzinfo is None:
                    due_dt = due_dt.replace(tzinfo=timezone.utc)
                if due_dt <= now:
                    row["복습대상"] += 1
            else:
                row["복습대상"] += 1
        except Exception:
            # 파싱 실패 시 기본적으로 복습 대상 처리(사용자에게 노출용으로는 안전한 기본값)
            row["복습대상"] += 1

        row["오답문항"] = wrong_by_subject.get(subj, row.get("오답문항", 0))
    # 기본 SRS는 연체/미래/신규를 따로 추적하지 않음
    return sorted(summary.values(), key=lambda x: (x["복습대상"], x["총문항"]), reverse=True)

def build_exam_payload(raw_items, exam_type):
    """문항 목록을 시험 진행용 payload로 변환"""
    parsed = []
    for raw in raw_items:
        if exam_type == "객관식":
            parsed_item = parse_mcq_content(raw)
        else:
            parsed_item = parse_cloze_content(raw)
        if parsed_item:
            parsed.append(parsed_item)
    return parsed

def start_exam_session_from_items(raw_items, exam_type, mode):
    """문항 리스트로 시험/학습 세션을 즉시 시작"""
    parsed = build_exam_payload(raw_items, exam_type)
    if not parsed:
        return 0
    st.session_state.exam_questions = parsed
    st.session_state.current_question_idx = 0
    st.session_state.user_answers = {}
    st.session_state.exam_started = True
    st.session_state.exam_finished = False
    st.session_state.exam_mode = mode
    st.session_state.exam_type = exam_type
    st.session_state.auto_next = False
    st.session_state.auto_advance_guard = None
    st.session_state.revealed_answers = set()
    st.session_state.exam_stats_applied = False
    st.session_state.graded_questions = set()
    st.session_state.exam_history_saved = False
    st.session_state.current_exam_meta = {
        "mode": mode,
        "type": exam_type,
        "subjects": sorted({(q.get("subject") or "General") for q in raw_items}),
        "units": sorted({get_unit_name(q) for q in raw_items}),
        "num_questions": len(parsed),
        "started_at": datetime.now(timezone.utc).isoformat()
    }
    return len(parsed)

def normalize_mcq_item(item):
    if not isinstance(item, dict):
        return None
    if "content" in item and "problem" not in item:
        parsed = extract_mcq_components(item.get("content", ""))
        if parsed:
            parsed["explanation"] = item.get("explanation", "")
            parsed["subject"] = item.get("subject")
            parsed["unit"] = item.get("unit")
            parsed["difficulty"] = item.get("difficulty")
            parsed["id"] = item.get("id")
            parsed["fsrs"] = item.get("fsrs")
            return parsed
    problem = (item.get("problem") or "").strip()
    options = item.get("options") or []
    answer = item.get("answer", 1)
    explanation = item.get("explanation", "")
    if not problem or not isinstance(options, list):
        return None
    # 옵션 길이 5로 정규화
    options = [str(opt).strip() for opt in options if str(opt).strip()]
    while len(options) < 5:
        options.append(f"보기 {len(options) + 1}")
    options = options[:5]
    try:
        answer_num = int(answer)
    except Exception:
        answer_num = 1
    if answer_num < 1 or answer_num > 5:
        answer_num = 1
    return {
        "type": "mcq",
        "problem": problem,
        "options": options,
        "answer": answer_num,
        "explanation": explanation,
        "subject": item.get("subject"),
        "unit": item.get("unit"),
        "difficulty": item.get("difficulty"),
        "id": item.get("id"),
        "fsrs": item.get("fsrs"),
    }

def normalize_cloze_item(item):
    if not isinstance(item, dict):
        return None
    if "content" in item and "front" not in item:
        # 구버전 content 필드
        content = item.get("content", "")
        if "{{c1::" in content:
            m = re.search(r'\{\{c1::(.+?)\}\}', content)
            if m:
                answer = m.group(1).strip()
                front = re.sub(r'\{\{c1::.+?\}\}', '____', content)
                return {
                    "type": "cloze",
                    "response_type": item.get("response_type", "cloze"),
                    "front": front,
                    "answer": answer,
                    "explanation": item.get("explanation", ""),
                    "subject": item.get("subject"),
                    "unit": item.get("unit"),
                    "difficulty": item.get("difficulty"),
                    "id": item.get("id"),
                    "fsrs": item.get("fsrs"),
                }
        return None
    front = (item.get("front") or "").strip()
    answer = (item.get("answer") or "").strip()
    explanation = item.get("explanation", "")
    response_type = item.get("response_type", "cloze")
    if response_type not in {"cloze", "short", "essay"}:
        response_type = "cloze"
    if not front or not answer:
        return None
    return {
        "type": "cloze",
        "response_type": response_type,
        "front": front,
        "answer": answer,
        "explanation": explanation,
        "subject": item.get("subject"),
        "unit": item.get("unit"),
        "difficulty": item.get("difficulty"),
        "id": item.get("id"),
        "fsrs": item.get("fsrs"),
    }

def format_explanation_text(text):
    if not text:
        return ""
    if "|" in text:
        parts = [p.strip() for p in re.split(r"\s*\|\s*", text) if p.strip()]
        if len(parts) > 1:
            return "\n".join([f"- {p}" for p in parts])
    return text

def _set_row_cant_split(row):
    tr_pr = row._tr.get_or_add_trPr()
    if not any(child.tag.endswith("cantSplit") for child in tr_pr):
        tr_pr.append(OxmlElement("w:cantSplit"))

def build_docx_question_sheet(items, title="Axioma Qbank 문제집"):
    doc = Document()
    doc.add_heading(title, level=1)
    doc.add_paragraph("좌측: 문항 | 우측: 정답 및 해설")
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.autofit = True
    table.rows[0].cells[0].text = "문항"
    table.rows[0].cells[1].text = "정답 & 해설"

    letters = ["A", "B", "C", "D", "E"]
    for i, item in enumerate(items, 1):
        row = table.add_row()
        _set_row_cant_split(row)
        left = row.cells[0]
        right = row.cells[1]

        stem = (item.get("problem") or item.get("front") or item.get("raw") or "").strip()
        left.text = f"문항 {i}\n{stem}"

        if item.get("type") == "mcq":
            opts = item.get("options") or []
            if opts:
                left.add_paragraph("")
            for j, opt in enumerate(opts[:5]):
                left.add_paragraph(f"{letters[j]}. {opt}")
            correct = item.get("answer") or item.get("correct")
            right.text = "정답"
            if isinstance(correct, int) and 1 <= correct <= 5:
                right.add_paragraph(letters[correct - 1])
            else:
                right.add_paragraph(str(correct))
        else:
            right.text = "정답"
            right.add_paragraph(str(item.get("answer", "")))

        explanation = (item.get("explanation") or "").strip()
        if explanation:
            right.add_paragraph("")
            right.add_paragraph("해설")
            for line in format_explanation_text(explanation).splitlines():
                right.add_paragraph(line)

    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    return out.getvalue()


def _to_markdown_table(rows):
    """pyarrow 의존성 없이 간단한 표 렌더링."""
    if not rows:
        st.caption("표시할 데이터가 없습니다.")
        return
    if hasattr(rows, "to_dict"):
        try:
            rows = rows.to_dict(orient="records")
        except Exception:
            rows = list(rows)
    if not isinstance(rows, list):
        rows = list(rows)
    if not rows:
        st.caption("표시할 데이터가 없습니다.")
        return

    first = rows[0]
    if not isinstance(first, dict):
        try:
            rows = [dict(item) for item in rows]
        except Exception:
            st.caption("표 형식 변환에 실패했습니다.")
            return

    if not rows:
        st.caption("표시할 데이터가 없습니다.")
        return

    headers = list(rows[0].keys())

    def _fmt(v):
        if v is None:
            return ""
        if isinstance(v, (list, tuple)):
            return ", ".join([str(x) for x in v])
        return str(v).replace("|", "\\|").replace("\n", "<br>")

    header_line = "| " + " | ".join(headers) + " |\n"
    sep_line = "|" + "|".join([" --- " for _ in headers]) + "| \n"
    body_lines = []
    for row in rows:
        vals = [_fmt(row.get(h, "")) for h in headers]
        body_lines.append("| " + " | ".join(vals) + " |")
    table = header_line + sep_line + "\n".join(body_lines)
    st.markdown(table)


def safe_dataframe(data, fallback_to_markdown=True, *args, **kwargs):
    """Use st.dataframe when available, fallback to markdown table if it fails."""
    try:
        return st.dataframe(data, *args, **kwargs)
    except Exception:
        if not fallback_to_markdown:
            raise
        _to_markdown_table(data)

def _is_option_line(line):
    if re.match(r"^\s*[①②③④⑤]", line):
        return True
    if re.match(r"^\s*[1-5][).]", line):
        return True
    return False

def _answer_token_to_num(token):
    token = str(token).strip()
    circled = {"①": 1, "②": 2, "③": 3, "④": 4, "⑤": 5}
    if token in circled:
        return circled[token]
    if token.isdigit():
        n = int(token)
        if 1 <= n <= 5:
            return n
    token = token.upper()
    if token in ["A", "B", "C", "D", "E"]:
        return ord(token) - ord("A") + 1
    return None

def preclean_exam_text(text):
    if not text:
        return ""
    lines = [l.rstrip() for l in text.splitlines()]

    # Find first probable question line
    q_re = re.compile(r"^\s*(?:문항|문제|Question|Q)?\s*\d{1,3}\s*[).]")
    q_alt = re.compile(r"[①②③④⑤]")
    first_idx = None
    for i, line in enumerate(lines):
        if q_re.match(line.strip()) or q_alt.search(line):
            first_idx = i
            break
    if first_idx is not None:
        lines = lines[first_idx:]

    # Remove page-only lines like "- 3 -" or empty separators
    cleaned = []
    for line in lines:
        s = line.strip()
        if not s:
            cleaned.append("")
            continue
        if re.match(r"^[-–—]{2,}$", s):
            cleaned.append("")
            continue
        if re.match(r"^[-–—]?\s*\d+\s*[-–—]?$", s):
            # page number lines
            cleaned.append("")
            continue
        cleaned.append(line)

    # Merge standalone number lines with the following text line
    merged = []
    i = 0
    num_re = re.compile(r"^\s*\d{1,3}\s*[).]?\s*$")
    while i < len(cleaned):
        line = cleaned[i]
        if num_re.match(line.strip()):
            j = i + 1
            while j < len(cleaned) and not cleaned[j].strip():
                j += 1
            if j < len(cleaned):
                merged.append(f"{line.strip()} {cleaned[j].strip()}".strip())
                i = j + 1
                continue
        merged.append(line)
        i += 1

    # Normalize excessive spaces
    merged = [re.sub(r"[ \t]+", " ", l).strip() for l in merged]
    return "\n".join([l for l in merged if l is not None]).strip()

def parse_exam_text_fuzzy(text, preclean=True):
    """기출문제 원문을 최대한 파싱해 MCQ/Cloze로 변환 (베타)"""
    if not text:
        return []
    if preclean:
        text = preclean_exam_text(text) or text

    def insert_breaks(raw):
        # Insert line breaks before common question markers to improve splitting
        raw = re.sub(r"(?<!\n)(Question\s*\d+\s*[).])", r"\n\1", raw, flags=re.IGNORECASE)
        raw = re.sub(r"(?<!\n)(문항\s*\d+\s*[).])", r"\n\1", raw)
        raw = re.sub(r"(?<!\n)(문제\s*\d+\s*[).])", r"\n\1", raw)
        raw = re.sub(r"(?<!\n)(Q\s*\d+\s*[).])", r"\n\1", raw, flags=re.IGNORECASE)
        return raw

    def split_exam_blocks_simple(raw):
        raw = insert_breaks(raw)
        pattern = re.compile(r"(?m)^\s*(?:문항|문제|Question|Q)?\s*(\d{1,3})\s*[).]\s*", re.IGNORECASE)
        matches = list(pattern.finditer(raw))
        if matches:
            blocks = []
            for i, m in enumerate(matches):
                start = m.start()
                end = matches[i + 1].start() if i + 1 < len(matches) else len(raw)
                blocks.append(raw[start:end].strip())
            return blocks
        blocks = [b.strip() for b in re.split(r"\n-{3,}\n", raw) if b.strip()]
        return blocks if blocks else [raw.strip()]

    def split_blocks(raw):
        raw = insert_breaks(raw)
        pattern = re.compile(r"(?m)^\s*(?:문항|문제|Question|Q)?\s*(\d{1,3})\s*[).]\s*", re.IGNORECASE)
        matches = list(pattern.finditer(raw))
        if matches:
            blocks = []
            for i, m in enumerate(matches):
                start = m.start()
                end = matches[i + 1].start() if i + 1 < len(matches) else len(raw)
                blocks.append((raw[start:end].strip(), int(m.group(1))))
            return blocks
        # fallback: split by long dashes or blank lines
        blocks = [b.strip() for b in re.split(r"\n-{3,}\n", raw) if b.strip()]
        return [(b, None) for b in blocks] if blocks else [(raw.strip(), None)]

    def extract_answer_and_explanation(block):
        ans = None
        exp_lines = []
        capturing = False
        for line in block.splitlines():
            line = line.strip()
            if not line:
                continue
            if re.match(r"^\s*(?:문항|문제|Question|Q)?\s*\d{1,3}\s*[).]\s*", line, re.IGNORECASE):
                if capturing:
                    break
            m = re.match(r"^(정답|답)\s*[:：]?\s*(.+)$", line)
            if m:
                ans = m.group(2).strip()
                capturing = True
                continue
            m2 = re.match(r"^(해설|설명)\s*[:：]?\s*(.+)$", line)
            if m2:
                capturing = True
                exp_lines.append(m2.group(2).strip())
                continue
            if capturing:
                if _is_option_line(line):
                    continue
                exp_lines.append(line)
        exp = "\n".join([l for l in exp_lines if l]).strip()
        return ans, exp

    items = []
    for block, qnum in split_blocks(text):
        if not block:
            continue
        source_page = None
        for line in block.splitlines():
            m_page = re.match(r"^===\s*페이지\s*(\d+)\s*===", line.strip())
            if m_page:
                source_page = int(m_page.group(1))
        ans_token, explanation = extract_answer_and_explanation(block)
        # remove answer/explanation lines for stem/options parsing
        cleaned = "\n".join(
            [ln for ln in block.splitlines() if not re.match(r"^\s*(정답|답|해설|설명)\s*[:：]", ln.strip())]
        ).strip()

        # try circled options
        if "①" in cleaned:
            parts = re.split(r"[①②③④⑤]", cleaned)
            stem = parts[0].strip()
            stem = re.sub(r"^\s*(?:문항\s*)?\d+\s*[).]\s*", "", stem).strip()
            options = [p.strip() for p in parts[1:] if p.strip()]
            if len(options) >= 3:
                answer_num = _answer_token_to_num(ans_token) or 1
                items.append({
                    "type": "mcq",
                    "problem": stem,
                    "options": options[:5],
                    "answer": answer_num,
                    "explanation": explanation,
                    "page": source_page,
                    "qnum": qnum,
                })
                continue

        # try numbered options (1) 2) ...
        opt_lines = re.findall(r"(?m)^\s*[1-5][).]\s*(.+)$", cleaned)
        if len(opt_lines) >= 3:
            stem = re.split(r"(?m)^\s*[1-5][).]\s*", cleaned)[0].strip()
            stem = re.sub(r"^\s*(?:문항\s*)?\d+\s*[).]\s*", "", stem).strip()
            answer_num = _answer_token_to_num(ans_token) or 1
            items.append({
                "type": "mcq",
                "problem": stem,
                "options": [o.strip() for o in opt_lines][:5],
                "answer": answer_num,
                "explanation": explanation,
                "page": source_page,
                "qnum": qnum,
            })
            continue

        # fallback to cloze if answer exists
        if ans_token:
            answer_text = str(ans_token).strip()
            stem = re.sub(r"^\s*(?:문항\s*)?\d+\s*[).]\s*", "", cleaned).strip()
            if stem and answer_text:
                items.append({
                    "type": "cloze",
                    "front": stem,
                    "answer": answer_text,
                    "explanation": explanation,
                    "page": source_page,
                    "qnum": qnum,
                })
                continue
    return clean_parsed_items(items)

def split_exam_blocks(raw):
    if not raw:
        return []
    raw = re.sub(r"(?<!\n)(Question\s*\d+\s*[).])", r"\n\1", raw, flags=re.IGNORECASE)
    raw = re.sub(r"(?<!\n)(문항\s*\d+\s*[).])", r"\n\1", raw)
    raw = re.sub(r"(?<!\n)(문제\s*\d+\s*[).])", r"\n\1", raw)
    raw = re.sub(r"(?<!\n)(Q\s*\d+\s*[).])", r"\n\1", raw, flags=re.IGNORECASE)
    pattern = re.compile(r"(?m)^\s*(?:문항|문제|Question|Q)?\s*(\d{1,3})\s*[).]\s*", re.IGNORECASE)
    matches = list(pattern.finditer(raw))
    if matches:
        blocks = []
        for i, m in enumerate(matches):
            start = m.start()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(raw)
            blocks.append(raw[start:end].strip())
        return blocks
    blocks = [b.strip() for b in re.split(r"\n-{3,}\n", raw) if b.strip()]
    return blocks if blocks else [raw.strip()]

def parse_answer_map_from_text(text):
    answer_map = {}
    for block in split_exam_blocks(text):
        if not block:
            continue
        m = re.match(r"^\s*(?:문항|문제|Question|Q)?\s*(\d{1,3})\s*[).]", block.strip(), re.IGNORECASE)
        qnum = int(m.group(1)) if m else None
        ans = None
        exp_lines = []
        for line in block.splitlines():
            l = line.strip()
            if not l:
                continue
            m_ans = re.search(r"(정답|답)\s*[:：]?\s*([①②③④⑤1-5])", l)
            if m_ans:
                ans = m_ans.group(2)
                rest = l[m_ans.end():].strip()
                if rest:
                    exp_lines.append(rest)
                continue
            m_ans2 = re.search(r"▶\s*([①②③④⑤1-5])", l)
            if m_ans2 and ans is None:
                ans = m_ans2.group(1)
                rest = l[m_ans2.end():].strip()
                if rest:
                    exp_lines.append(rest)
                continue
            m_qans = re.match(r"^\s*\d{1,3}\s*[).]?\s*([①②③④⑤1-5])\b\s*(.*)$", l)
            if m_qans and ans is None:
                ans = m_qans.group(1)
                if m_qans.group(2).strip():
                    exp_lines.append(m_qans.group(2).strip())
                continue
            if ans is None and re.match(r"^[①②③④⑤1-5]$", l):
                ans = l
                continue
            if ans is not None:
                if re.match(r"^\s*(?:문항|문제|Question|Q)?\s*\d{1,3}\s*[).]", l, re.IGNORECASE):
                    break
                exp_lines.append(l)
        if qnum and ans:
            answer_map[qnum] = {"answer": ans, "explanation": "\n".join(exp_lines).strip()}
    return answer_map

def parse_pdf_layout(pdf_bytes):
    items_all = []
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        for page_idx in range(doc.page_count):
            page = doc.load_page(page_idx)
            width = page.rect.width
            data = page.get_text("dict")
            lines = []
            for block in data.get("blocks", []):
                for line in block.get("lines", []):
                    text = "".join(span.get("text", "") for span in line.get("spans", []))
                    text = text.strip()
                    if not text:
                        continue
                    x0, y0, x1, y1 = line.get("bbox", [0, 0, 0, 0])
                    lines.append({"text": text, "x0": x0, "x1": x1, "y0": y0})

            if not lines:
                continue

            centers = [((l["x0"] + l["x1"]) / 2) for l in lines]
            left_lines = [l for l, c in zip(lines, centers) if c < width * 0.45]
            right_lines = [l for l, c in zip(lines, centers) if c > width * 0.55]
            middle_lines = [l for l, c in zip(lines, centers) if width * 0.45 <= c <= width * 0.55]
            marker_lines = [l for l in middle_lines if re.match(r"^\s*\d{1,3}\s*[).]?\s*$", l["text"])]
            two_col = len(left_lines) >= 5 and len(right_lines) >= 5

            def merge_number_lines(ls, tol=4.0):
                num_re = re.compile(r"^\s*\d{1,3}\s*[).]?\s*$")
                merged = set()
                for i, num_line in enumerate(ls):
                    if not num_re.match(num_line["text"]):
                        continue
                    # find closest non-number line within tolerance
                    candidates = []
                    for j, other in enumerate(ls):
                        if i == j or num_re.match(other["text"]):
                            continue
                        dy = abs(other["y0"] - num_line["y0"])
                        if dy <= tol:
                            candidates.append((dy, j, other))
                    if candidates:
                        _, j, target = min(candidates, key=lambda x: x[0])
                        prefix = num_line["text"].strip()
                        if not target["text"].strip().startswith(prefix):
                            target["text"] = f"{prefix} {target['text']}".strip()
                        merged.add(i)
                return [l for idx, l in enumerate(ls) if idx not in merged]

            def build_text(ls):
                ls_sorted = sorted(ls, key=lambda x: (x["y0"], x["x0"]))
                text = "\n".join([l["text"] for l in ls_sorted])
                return f"=== 페이지 {page_idx + 1} ===\n" + text

            if two_col:
                left_text = build_text(merge_number_lines(left_lines + marker_lines))
                right_text = build_text(merge_number_lines(right_lines + marker_lines))
                items = parse_exam_text_fuzzy(left_text)
                ans_map = parse_answer_map_from_text(right_text)
                for idx, it in enumerate(items):
                    if not it.get("page"):
                        it["page"] = page_idx + 1
                    qnum = it.get("qnum")
                    if qnum in ans_map:
                        ans_token = ans_map[qnum].get("answer")
                        exp = ans_map[qnum].get("explanation") or ""
                    else:
                        # fallback: 순서 기반 매칭
                        keys = sorted(ans_map.keys())
                        ans_token = ans_map.get(keys[idx], {}).get("answer") if idx < len(keys) else None
                        exp = ans_map.get(keys[idx], {}).get("explanation") if idx < len(keys) else ""
                    if it.get("type") == "mcq" and ans_token:
                        it["answer"] = _answer_token_to_num(ans_token) or it.get("answer")
                    elif it.get("type") == "cloze" and ans_token:
                        it["answer"] = it.get("answer") or ans_token
                    if exp and not it.get("explanation"):
                        it["explanation"] = exp
                items_all.extend(items)
            else:
                full_text = build_text(lines)
                items = parse_exam_text_fuzzy(full_text)
                for it in items:
                    if not it.get("page"):
                        it["page"] = page_idx + 1
                items_all.extend(items)
        doc.close()
    except Exception:
        return []
    return clean_parsed_items(items_all)

def parse_pdf_layout_ai(pdf_bytes, ai_model, api_key=None, openai_api_key=None, hint_text=""):
    items_all = []
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        for page_idx in range(doc.page_count):
            page = doc.load_page(page_idx)
            width = page.rect.width
            data = page.get_text("dict")
            lines = []
            for block in data.get("blocks", []):
                for line in block.get("lines", []):
                    text = "".join(span.get("text", "") for span in line.get("spans", []))
                    text = text.strip()
                    if not text:
                        continue
                    x0, y0, x1, y1 = line.get("bbox", [0, 0, 0, 0])
                    lines.append({"text": text, "x0": x0, "x1": x1, "y0": y0})
            if not lines:
                continue
            centers = [((l["x0"] + l["x1"]) / 2) for l in lines]
            left_lines = [l for l, c in zip(lines, centers) if c < width * 0.45]
            right_lines = [l for l, c in zip(lines, centers) if c > width * 0.55]
            middle_lines = [l for l, c in zip(lines, centers) if width * 0.45 <= c <= width * 0.55]
            marker_lines = [l for l in middle_lines if re.match(r"^\s*\d{1,3}\s*[).]?\s*$", l["text"])]

            def build_text(ls):
                ls_sorted = sorted(ls, key=lambda x: (x["y0"], x["x0"]))
                text = "\n".join([l["text"] for l in ls_sorted])
                return f"=== 페이지 {page_idx + 1} ===\n" + text

            left_text = build_text(left_lines + marker_lines)
            right_text = build_text(right_lines + marker_lines)
            ai_items = ai_parse_exam_layout(
                left_text,
                right_text,
                ai_model=ai_model,
                api_key=api_key,
                openai_api_key=openai_api_key,
                hint_text=hint_text
            )
            for it in ai_items:
                it["page"] = page_idx + 1
            items_all.extend(ai_items)
        doc.close()
    except Exception:
        return []
    return clean_parsed_items(items_all)

def extract_pdf_page_texts(pdf_bytes):
    texts = []
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        for i in range(doc.page_count):
            page = doc.load_page(i)
            page_text = page.get_text().strip()
            texts.append(page_text)
        doc.close()
    except Exception:
        return []
    return texts

def match_questions_to_pages(items, page_texts):
    scores = {}
    if not items or not page_texts:
        return scores
    page_tokens = [_tokenize_for_match(t) for t in page_texts]
    for idx, item in enumerate(items):
        stem = (item.get("problem") or item.get("front") or "")
        tokens = _tokenize_for_match(stem)
        if not tokens:
            continue
        best_page = None
        best_score = 0.0
        for p_idx, pt in enumerate(page_tokens):
            inter = tokens & pt
            score = len(inter) / max(1, len(tokens))
            if score > best_score:
                best_score = score
                best_page = p_idx + 1
        if best_page:
            scores[idx] = {"page": best_page, "score": best_score}
            item["page"] = best_page
    return scores

def parse_qa_to_cloze(text):
    """정답: 패턴을 이용해 Q/A를 Cloze 형태로 변환"""
    results = []
    lines = [l.strip() for l in text.splitlines()]
    buffer_lines = []
    last_item = None
    for line in lines:
        if not line:
            continue
        if re.match(r"^(해설|설명)\s*[:：]", line):
            explanation = re.split(r"[:：]", line, 1)[1].strip()
            if last_item:
                last_item["explanation"] = explanation
            continue
        m = re.match(r"^(정답|답)\s*[:：]\s*(.+)$", line)
        if m:
            answer = m.group(2).strip()
            question = " ".join(buffer_lines).strip()
            if question and answer:
                last_item = {
                    "type": "cloze",
                    "front": question,
                    "answer": answer,
                    "explanation": ""
                }
                results.append(last_item)
            buffer_lines = []
        else:
            buffer_lines.append(line)
    return results

def apply_theme(theme_mode, bg_mode):
    color_scheme = "dark" if theme_mode == "Dark" else "light"
    if theme_mode == "Dark":
        base_bg = "#0f172a"
        surface = "#111827"
        surface_2 = "#0b1220"
        text = "#e5e7eb"
        subtext = "#94a3b8"
        accent = "#38bdf8"
        accent2 = "#22d3ee"
        border = "#1f2937"
    else:
        base_bg = "#f8fafc"
        surface = "#ffffff"
        surface_2 = "#f1f5f9"
        text = "#111827"
        subtext = "#4b5563"
        accent = "#0ea5a4"
        accent2 = "#14b8a6"
        border = "#e5e7eb"

    if bg_mode == "None":
        bg = "none"
    elif bg_mode == "Grid":
        bg = "radial-gradient(circle, rgba(148,163,184,0.2) 1px, transparent 1px)"
    else:
        if theme_mode == "Dark":
            bg = "radial-gradient(1200px 600px at 10% 0%, rgba(59,130,246,0.15), transparent 60%), linear-gradient(180deg, #0b1220 0%, #0f172a 100%)"
        else:
            bg = "radial-gradient(1200px 600px at 10% 0%, rgba(20,184,166,0.12), transparent 60%), linear-gradient(180deg, #f8fafc 0%, #eef2ff 100%)"

    st.markdown(
        f"""
        <style>
        :root {{
            --bg: {base_bg};
            --surface: {surface};
            --surface-2: {surface_2};
            --text: {text};
            --muted: {subtext};
            --accent: {accent};
            --accent-2: {accent2};
            --border: {border};
        }}
        html, body {{
            color: var(--text);
            color-scheme: {color_scheme};
        }}
        .stApp {{
            background-color: var(--bg);
            background-image: {bg};
            color: var(--text);
            color-scheme: {color_scheme};
        }}
        [data-testid="stAppViewContainer"] {{
            color: var(--text) !important;
        }}
        [data-testid="stAppViewContainer"] p,
        [data-testid="stAppViewContainer"] span,
        [data-testid="stAppViewContainer"] label,
        [data-testid="stAppViewContainer"] li,
        [data-testid="stAppViewContainer"] strong,
        [data-testid="stAppViewContainer"] h1,
        [data-testid="stAppViewContainer"] h2,
        [data-testid="stAppViewContainer"] h3,
        [data-testid="stAppViewContainer"] h4,
        [data-testid="stAppViewContainer"] h5,
        [data-testid="stAppViewContainer"] h6 {{
            color: var(--text) !important;
            opacity: 1 !important;
        }}
        [data-testid="stSidebar"] {{
            background: var(--surface);
            border-right: 1px solid var(--border);
        }}
        [data-testid="stSidebar"] * {{
            color: var(--text) !important;
            opacity: 1 !important;
        }}
        [data-testid="stHeader"] {{
            background: transparent;
        }}
        .stButton>button {{
            background: var(--accent);
            color: #ffffff !important;
            border: 1px solid var(--accent);
            border-radius: 10px;
        }}
        .stButton>button * {{
            color: #ffffff !important;
        }}
        .stButton>button:hover {{
            background: var(--accent-2);
            color: #ffffff !important;
        }}
        div[data-baseweb="input"] > div,
        div[data-baseweb="select"] > div,
        .stTextArea textarea {{
            background: var(--surface-2);
            border: 1px solid var(--border);
            color: var(--text) !important;
            -webkit-text-fill-color: var(--text) !important;
        }}
        input::placeholder,
        textarea::placeholder {{
            color: var(--muted) !important;
            opacity: 1 !important;
        }}
        .stTabs [data-baseweb="tab-list"] {{
            background: var(--surface);
            border: 1px solid var(--border);
            border-radius: 10px;
        }}
        .stTabs [data-baseweb="tab"] {{
            background: transparent !important;
            border-radius: 8px;
            transition: color 0.22s ease, background-color 0.22s ease;
        }}
        .stTabs [data-baseweb="tab"]:hover {{
            background: var(--surface-2) !important;
        }}
        .stTabs [aria-selected="true"] {{
            background: transparent !important;
            color: var(--accent) !important;
            font-weight: 700;
        }}
        .stAlert, .stExpander, .stMetric {{
            background: var(--surface);
            border: 1px solid var(--border);
            color: var(--text);
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )

def should_apply_custom_theme(theme_enabled, theme_mode):
    # Dark mode should be visible even when the custom-theme switch is off.
    return bool(theme_enabled) or str(theme_mode) == "Dark"

def apply_mobile_exam_styles():
    st.markdown(
        """
        <style>
        [data-testid="stRadio"] [role="radiogroup"] > label {
            padding: 0.55rem 0.6rem;
            border: 1px solid var(--border);
            border-radius: 10px;
            margin-bottom: 0.35rem;
            background: var(--surface);
        }
        [data-testid="stSelectbox"] label p,
        [data-testid="stRadio"] label p,
        [data-testid="stTextInput"] label p {
            font-size: 1rem !important;
        }
        .stButton > button {
            min-height: 44px;
            font-size: 1rem;
        }
        input[type="text"] {
            font-size: 16px !important;
        }
        .mobile-exam-caption {
            font-size: 0.95rem;
            color: var(--muted);
            margin-bottom: 0.3rem;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

def render_auth_landing_page():
    st.markdown(
        """
        <style>
        .auth-shell {
            max-width: 560px;
            margin: 1.5rem auto 0.5rem auto;
        }
        .auth-brand {
            text-align: center;
            font-size: 2rem;
            font-weight: 700;
            color: var(--text);
            letter-spacing: -0.02em;
            margin-bottom: 0.2rem;
        }
        .auth-subtitle {
            text-align: center;
            color: var(--muted);
            margin-bottom: 1rem;
        }
        .auth-card {
            background: var(--surface);
            border: 1px solid var(--border);
            border-radius: 16px;
            padding: 1rem 1rem 0.6rem 1rem;
            box-shadow: 0 10px 30px rgba(15, 23, 42, 0.08);
        }
        .auth-help {
            color: var(--muted);
            font-size: 0.9rem;
            text-align: center;
            margin-top: 0.5rem;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )
    st.markdown("<div class='auth-shell'>", unsafe_allow_html=True)
    st.markdown("<div class='auth-brand'>Axioma Qbank</div>", unsafe_allow_html=True)
    st.markdown("<div class='auth-subtitle'>이메일 계정으로 로그인하거나 새 계정을 만드세요.</div>", unsafe_allow_html=True)
    st.markdown("<div class='auth-card'>", unsafe_allow_html=True)

    if is_supabase_required() or is_supabase_enabled():
        tab_login, tab_signup = st.tabs(["Log in", "Create account"])
        with tab_login:
            with st.form("auth_login_form_main", clear_on_submit=False):
                login_email = st.text_input("EMAIL ADDRESS", key="auth_login_email_main")
                login_password = st.text_input("PASSWORD", type="password", key="auth_login_password_main")
                login_submit = st.form_submit_button("Log in", use_container_width=True)
            if login_submit:
                if not is_valid_email(login_email):
                    st.error("이메일 형식을 확인해주세요.")
                else:
                    ok, result = authenticate_user_account(login_email, login_password)
                    if ok:
                        reset_runtime_state_for_auth_change()
                        st.session_state.auth_user_id = result
                        st.rerun()
                    else:
                        st.error(result)
        with tab_signup:
            with st.form("auth_signup_form_main", clear_on_submit=True):
                signup_email = st.text_input("EMAIL ADDRESS", key="auth_signup_email_main")
                signup_password = st.text_input("PASSWORD (6자 이상)", type="password", key="auth_signup_password_main")
                signup_password_confirm = st.text_input("CONFIRM PASSWORD", type="password", key="auth_signup_password_confirm_main")
                signup_submit = st.form_submit_button("Create account", use_container_width=True)
            if signup_submit:
                if not is_valid_email(signup_email):
                    st.error("이메일 형식을 확인해주세요.")
                elif signup_password != signup_password_confirm:
                    st.error("비밀번호 확인이 일치하지 않습니다.")
                else:
                    ok, message = register_user_account(signup_email, signup_password)
                    if ok:
                        st.success(message)
                    else:
                        st.error(message)
        st.markdown("<div class='auth-help'>Supabase Auth 모드: 이메일/비밀번호 로그인</div>", unsafe_allow_html=True)
    else:
        tab_login, tab_signup = st.tabs(["로그인", "회원가입"])
        with tab_login:
            with st.form("auth_login_form_main_local", clear_on_submit=False):
                login_user_id = st.text_input("아이디", key="auth_login_user_id_main")
                login_password = st.text_input("비밀번호", type="password", key="auth_login_password_main_local")
                login_submit = st.form_submit_button("로그인", use_container_width=True)
            if login_submit:
                ok, result = authenticate_user_account(login_user_id, login_password)
                if ok:
                    reset_runtime_state_for_auth_change()
                    st.session_state.auth_user_id = result
                    st.rerun()
                else:
                    st.error(result)
        with tab_signup:
            with st.form("auth_signup_form_main_local", clear_on_submit=True):
                signup_user_id = st.text_input("새 아이디", key="auth_signup_user_id_main")
                signup_password = st.text_input("새 비밀번호 (6자 이상)", type="password", key="auth_signup_password_main_local")
                signup_submit = st.form_submit_button("회원가입", use_container_width=True)
            if signup_submit:
                ok, message = register_user_account(signup_user_id, signup_password)
                if ok:
                    st.success(message)
                else:
                    st.error(message)
        st.markdown("<div class='auth-help'>로컬 계정 모드</div>", unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

def show_action_notice():
    msg = st.session_state.get("last_action_notice", "")
    if msg:
        st.success(msg)
        st.session_state.last_action_notice = ""

def render_copyright_ack(scope_key: str):
    st.info("업로드 자료는 권리를 보유하거나 사용 허락을 받은 자료만 사용하세요. 원문 파일은 영구 저장하지 않고 세션 처리 후 폐기됩니다.")
    ack_rights = st.checkbox(
        "업로드 자료에 대한 이용 권리를 보유/허락받았음을 확인합니다.",
        key=f"copyright_ack_rights_{scope_key}",
    )
    ack_no_redistribute = st.checkbox(
        "타인의 저작물을 무단 재배포하지 않겠습니다.",
        key=f"copyright_ack_no_redistribute_{scope_key}",
    )
    return bool(ack_rights and ack_no_redistribute)

def render_generation_recovery_panel():
    if not st.session_state.get("generation_failure"):
        return
    with st.container():
        st.markdown("### ⚠️ 문제 생성 실패")
        st.error(st.session_state.generation_failure)
        st.caption("아래 버튼으로 바로 복구/초기화를 할 수 있습니다.")
        colr1, colr2 = st.columns(2)
        with colr1:
            if st.button("🔁 동일 조건 재실행", use_container_width=True, key="failure_retry_btn"):
                st.session_state.generation_failure = ""
                st.rerun()
        with colr2:
            if st.button("🧹 알림 지우기", use_container_width=True, key="failure_clear_btn"):
                st.session_state.generation_failure = ""

def compute_activity_heatmap(questions, days=365, now=None):
    check_time = now or datetime.now(timezone.utc)
    start = (check_time - timedelta(days=days - 1)).date()
    buckets = {}
    for i in range(days):
        d = start + timedelta(days=i)
        buckets[d.isoformat()] = {"count": 0, "correct": 0, "total": 0}
    for q in questions:
        stats = q.get("stats") or {}
        hist = stats.get("history") or []
        for entry in hist:
            if not isinstance(entry, dict):
                continue
            dt = parse_iso_datetime(entry.get("time"))
            if not dt:
                continue
            if dt.tzinfo is None:
                dt = dt.replace(tzinfo=timezone.utc)
            dkey = dt.date().isoformat()
            if dkey in buckets:
                buckets[dkey]["count"] += 1
                buckets[dkey]["total"] += 1
                if entry.get("correct") is True:
                    buckets[dkey]["correct"] += 1
    rows = []
    for dkey, val in buckets.items():
        d = datetime.fromisoformat(dkey).date()
        week_index = (d - start).days // 7
        rows.append({
            "date": d,
            "dow": d.weekday(),
            "week_index": week_index,
            "count": val["count"],
            "accuracy": (val["correct"] / val["total"] * 100) if val["total"] > 0 else 0
        })
    return rows

def fsrs_due(item, now=None):
    if not FSRS_AVAILABLE:
        return True
    try:
        fsrs = item.get("fsrs") or {}
        card_data = fsrs.get("card")
        if not card_data:
            return True
        card = Card.from_json(card_data)
        check_time = now or datetime.now(timezone.utc)
        return card.due <= check_time
    except Exception:
        return True

def simple_srs_due(item, now=None):
    check_time = now or datetime.now(timezone.utc)
    srs = item.get("srs") or {}
    due = parse_iso_datetime(srs.get("due"))
    return due is None or due <= check_time

def srs_due(item, now=None):
    if FSRS_AVAILABLE:
        return fsrs_due(item, now=now)
    return simple_srs_due(item, now=now)

def apply_simple_srs_rating(q_id, rating_label):
    bank = load_questions()
    now = datetime.now(timezone.utc)
    # base intervals in days
    base = {"Again": 1, "Hard": 2, "Good": 4, "Easy": 7}
    for key in ("text", "cloze"):
        for item in bank.get(key, []):
            if item.get("id") == q_id:
                srs = item.get("srs") or {}
                interval = int(srs.get("interval", 1))
                factor = {"Again": 0.5, "Hard": 1.2, "Good": 2.0, "Easy": 3.0}.get(rating_label, 2.0)
                new_interval = max(1, int(interval * factor))
                # if first time, use base
                if not srs:
                    new_interval = base.get(rating_label, 4)
                due = now + timedelta(days=new_interval)
                srs.update({
                    "interval": new_interval,
                    "due": due.isoformat(),
                    "last_rating": rating_label,
                    "last_review": now.isoformat(),
                })
                item["srs"] = srs
                save_questions(bank)
                return srs
    return None

def apply_srs_rating(q_id, rating):
    if FSRS_AVAILABLE:
        return apply_fsrs_rating(q_id, rating)
    # rating can be string label
    label = rating if isinstance(rating, str) else str(rating)
    return apply_simple_srs_rating(q_id, label)

def get_fsrs_queue(questions, now=None, limit=50):
    if not FSRS_AVAILABLE:
        return []
    check_time = now or datetime.now(timezone.utc)
    due_items = []
    for q in questions:
        fsrs = q.get("fsrs") or {}
        card_data = fsrs.get("card")
        if not card_data:
            due_items.append((q, check_time))
            continue
        try:
            card = Card.from_json(card_data)
            due_time = card.due
        except Exception:
            due_time = check_time
        if due_time <= check_time:
            due_items.append((q, due_time))
    due_items.sort(key=lambda x: x[1])
    return due_items[:limit]

def get_fsrs_stats(questions, now=None):
    if not FSRS_AVAILABLE:
        return None
    check_time = now or datetime.now(timezone.utc)
    due = 0
    overdue = 0
    future = 0
    new = 0
    for q in questions:
        fsrs = q.get("fsrs") or {}
        card_data = fsrs.get("card")
        if not card_data:
            new += 1
            due += 1
            continue
        try:
            card = Card.from_json(card_data)
            if card.due <= check_time:
                due += 1
                if card.due < check_time:
                    overdue += 1
            else:
                future += 1
        except Exception:
            due += 1
    return {
        "due": due,
        "overdue": overdue,
        "future": future,
        "new": new,
    }

def apply_fsrs_rating(q_id, rating):
    if not FSRS_AVAILABLE:
        return None
    bank = load_questions()
    now = datetime.now(timezone.utc)
    for key in ("text", "cloze"):
        for item in bank.get(key, []):
            if item.get("id") == q_id:
                card_data = (item.get("fsrs") or {}).get("card")
                if card_data:
                    try:
                        card = Card.from_json(card_data)
                    except Exception:
                        card = Card()
                else:
                    card = Card()
                scheduler = get_fsrs_scheduler() or Scheduler()
                card, log = scheduler.review_card(card, rating, now)
                fsrs = item.get("fsrs") or {}
                fsrs["card"] = card.to_json()
                fsrs["last_review"] = now.isoformat()
                fsrs["last_rating"] = rating.name if hasattr(rating, "name") else str(rating)
                fsrs["due"] = card.due.isoformat()
                logs = fsrs.get("logs", [])
                try:
                    logs.append(log.to_json())
                except Exception:
                    pass
                fsrs["logs"] = logs[-50:]
                item["fsrs"] = fsrs
                save_questions(bank)
                return fsrs
    return None

# ============================================================================
# 텍스트 추출 함수
# ============================================================================
@st.cache_resource(show_spinner=False)
def get_easyocr_reader(langs):
    try:
        import easyocr
    except Exception:
        return None
    return easyocr.Reader(list(langs), gpu=False)

def available_ocr_engines():
    engines = []
    if importlib.util.find_spec("easyocr") is not None:
        engines.append("easyocr")
    return engines

def ocr_page_image_bytes(image_bytes, engine="easyocr", langs=("ko", "en")):
    if engine != "easyocr":
        raise ValueError(f"지원하지 않는 OCR 엔진: {engine}")
    reader = get_easyocr_reader(tuple(langs))
    if reader is None:
        raise ValueError("easyocr 미설치")
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
        tmp.write(image_bytes)
        tmp_path = tmp.name
    try:
        results = reader.readtext(tmp_path, detail=1, paragraph=False)
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
    if not results:
        return ""
    def bbox_key(item):
        bbox = item[0] if isinstance(item, (list, tuple)) and item else None
        if not bbox:
            return (0, 0)
        ys = [p[1] for p in bbox]
        xs = [p[0] for p in bbox]
        return (min(ys), min(xs))
    results = sorted(results, key=bbox_key)
    lines = [r[1].strip() for r in results if len(r) > 1 and str(r[1]).strip()]
    return "\n".join(lines)

def ocr_pdf_bytes(pdf_bytes, engine="easyocr", langs=("ko", "en"), max_pages=0, zoom=2.0):
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    texts = []
    total_pages = doc.page_count
    limit = total_pages if max_pages in (0, None) else min(total_pages, max_pages)
    for i in range(limit):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        image_bytes = pix.tobytes("png")
        page_text = ocr_page_image_bytes(image_bytes, engine=engine, langs=langs)
        if page_text.strip():
            texts.append(f"=== 페이지 {i + 1} ===")
            texts.append(page_text)
            texts.append("")
    doc.close()
    return "\n".join(texts).strip()

def _ocr_pdf_page_with_ai(page_image_bytes, ai_model, api_key=None, openai_api_key=None):
    prompt = (
        "이 페이지에 보이는 텍스트를 가능한 한 정확하게 추출해서 JSON이나 해설 없이 \n"
        "순수 텍스트만 줄바꿈 유지 형식으로 반환하세요.\n"
        "표/문항 번호/선지 등은 모두 읽을 수 있는 그대로 보존하세요."
    )
    try:
        if ai_model == "🔵 Google Gemini":
            if not api_key:
                return ""
            genai.configure(api_key=api_key)
            requested_model = get_gemini_model_id()
            model_candidates = [requested_model]
            if requested_model != "gemini-2.0-flash":
                model_candidates.append("gemini-2.0-flash")

            # google-generativeai 버전별 허용 입력 포맷이 달라서 순차 시도
            img_pil = None
            try:
                img_pil = Image.open(io.BytesIO(page_image_bytes))
            except Exception:
                img_pil = None

            payloads = []
            if img_pil is not None:
                payloads.append([prompt, img_pil])
            payloads.append([prompt, {"mime_type": "image/png", "data": page_image_bytes}])
            payloads.append([prompt, page_image_bytes])

            for model_name in model_candidates:
                try:
                    model = genai.GenerativeModel(model_name)
                    for payload in payloads:
                        try:
                            response = model.generate_content(payload)
                            text = (getattr(response, "text", "") or "").strip()
                            if text:
                                return text
                        except Exception:
                            continue
                except Exception:
                    continue
            return ""
        if not openai_api_key:
            return ""
        client = OpenAI(api_key=openai_api_key)
        image_url = data_uri_from_bytes(page_image_bytes, ext="png")
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": image_url}},
                    ],
                }
            ],
            temperature=0,
        )
        return (response.choices[0].message.content or "").strip()
    except Exception:
        return ""


def ocr_pdf_bytes_with_ai(pdf_bytes, ai_model, api_key=None, openai_api_key=None, max_pages=0, zoom=2.0):
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    texts = []
    total_pages = doc.page_count
    limit = total_pages if max_pages in (0, None) else min(total_pages, max_pages)
    for i in range(limit):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        image_bytes = pix.tobytes("png")
        page_text = _ocr_pdf_page_with_ai(image_bytes, ai_model, api_key=api_key, openai_api_key=openai_api_key)
        if page_text.strip():
            texts.append(f"=== 페이지 {i + 1} ===")
            texts.append(page_text)
            texts.append("")
    doc.close()
    return "\n".join(texts).strip()

def data_uri_from_bytes(data, ext):
    ext = ext.lower().replace(".", "")
    if ext in ("jpg", "jpeg"):
        mime = "image/jpeg"
    elif ext == "png":
        mime = "image/png"
    elif ext == "bmp":
        mime = "image/bmp"
    elif ext == "gif":
        mime = "image/gif"
    else:
        mime = "application/octet-stream"
    b64 = base64.b64encode(data).decode("utf-8")
    return f"data:{mime};base64,{b64}"

def data_uri_to_bytes(uri):
    if not uri:
        return b""
    m = re.match(r"^data:.*?;base64,(.*)$", uri)
    if not m:
        return b""
    try:
        return base64.b64decode(m.group(1))
    except Exception:
        return b""

def extract_images_from_pdf_bytes(pdf_bytes, max_images=80, min_kb=20):
    images = []
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        seen = set()
        for page_idx in range(doc.page_count):
            page = doc.load_page(page_idx)
            for img in page.get_images(full=True):
                xref = img[0]
                base = doc.extract_image(xref)
                if not base or "image" not in base:
                    continue
                data = base["image"]
                if len(data) < min_kb * 1024:
                    continue
                rect = None
                try:
                    rect = page.get_image_bbox(xref)
                except Exception:
                    rect = None
                h = hashlib.sha1(data).hexdigest()
                if h in seen:
                    continue
                seen.add(h)
                ext = base.get("ext", "png")
                images.append({
                    "data_uri": data_uri_from_bytes(data, ext),
                    "ext": ext,
                    "page": page_idx + 1,
                    "y": rect.y0 if rect else None,
                    "y1": rect.y1 if rect else None,
                })
                if len(images) >= max_images:
                    break
            if len(images) >= max_images:
                break
        doc.close()
    except Exception:
        return []
    return images

def extract_pdf_question_anchors(pdf_bytes):
    anchors = {}
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        q_pattern = re.compile(r"^\s*(?:문항\s*)?(\d{1,3})\s*[).]")
        for page_idx in range(doc.page_count):
            page = doc.load_page(page_idx)
            page_anchors = []
            data = page.get_text("dict")
            for block in data.get("blocks", []):
                for line in block.get("lines", []):
                    line_text = "".join(span.get("text", "") for span in line.get("spans", []))
                    if not line_text:
                        continue
                    m = q_pattern.match(line_text.strip())
                    if m:
                        qnum = int(m.group(1))
                        y = line.get("bbox", [0, 0, 0, 0])[1]
                        page_anchors.append({"qnum": qnum, "y": y})
            if page_anchors:
                # de-duplicate by qnum, keep first occurrence
                seen = set()
                uniq = []
                for a in sorted(page_anchors, key=lambda x: x["y"]):
                    if a["qnum"] in seen:
                        continue
                    seen.add(a["qnum"])
                    uniq.append(a)
                anchors[page_idx + 1] = uniq
        doc.close()
    except Exception:
        return {}
    return anchors

def extract_images_from_hwp_bytes(hwp_bytes, max_images=80, min_kb=10):
    tmp_path = None
    odt_path = None
    images = []
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".hwp") as tmp:
            tmp.write(hwp_bytes)
            tmp_path = tmp.name
        with tempfile.NamedTemporaryFile(delete=False, suffix=".odt") as tmp_odt:
            odt_path = tmp_odt.name

        if shutil.which("hwp5odt"):
            result = subprocess.run(["hwp5odt", "--output", odt_path, tmp_path], capture_output=True, text=True)
            if result.returncode != 0:
                return []
        else:
            return []

        with zipfile.ZipFile(odt_path) as zf:
            for name in zf.namelist():
                if not name.startswith("bindata/"):
                    continue
                data = zf.read(name)
                if len(data) < min_kb * 1024:
                    continue
                ext = os.path.splitext(name)[1].lstrip(".") or "png"
                images.append({
                    "data_uri": data_uri_from_bytes(data, ext),
                    "ext": ext,
                    "page": None,
                })
                if len(images) >= max_images:
                    break
    except Exception:
        return []
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except Exception:
                pass
        if odt_path and os.path.exists(odt_path):
            try:
                os.unlink(odt_path)
            except Exception:
                pass
    return images

def _tokenize_for_match(text):
    if not text:
        return set()
    tokens = re.findall(r"[A-Za-z가-힣0-9]{2,}", text.lower())
    return set(tokens)

def clean_parsed_items(items, min_stem_len=15):
    cleaned = []
    for item in items or []:
        if not isinstance(item, dict):
            continue
        typ = item.get("type")
        if typ not in ("mcq", "cloze"):
            continue
        stem = (item.get("problem") if typ == "mcq" else item.get("front")) or ""
        stem = stem.strip()
        if not stem:
            continue
        if re.match(r"^(정답|답|해설|설명)\b", stem):
            continue
        if len(stem) < min_stem_len:
            if typ == "mcq" and len(item.get("options") or []) >= 3:
                pass
            else:
                continue
        if typ == "mcq":
            if len(item.get("options") or []) < 3:
                continue
        if typ == "cloze" and not str(item.get("answer", "")).strip():
            continue
        cleaned.append(item)
    return cleaned

def ocr_images_for_matching(images, engine="easyocr", langs=("ko", "en"), max_images=30, min_len=3):
    if not images:
        return images
    count = 0
    for img in images:
        if count >= max_images:
            break
        if img.get("ocr_text"):
            continue
        data = data_uri_to_bytes(img.get("data_uri", ""))
        if not data:
            continue
        try:
            text = ocr_page_image_bytes(data, engine=engine, langs=langs)
        except Exception:
            text = ""
        if text and len(text.strip()) >= min_len:
            img["ocr_text"] = text
        else:
            img["ocr_text"] = ""
        count += 1
    return images

def ai_match_images_to_items(items, images, ai_model, api_key=None, openai_api_key=None, max_images=10):
    if not items or not images or max_images <= 0:
        return items
    # group items by page
    page_map = {}
    for idx, item in enumerate(items):
        page = item.get("page")
        page_map.setdefault(page, []).append((idx, item))

    processed = 0
    for img in images:
        if processed >= max_images:
            break
        if img.get("matched"):
            continue
        page = img.get("page")
        candidates = page_map.get(page) or []
        if not candidates:
            continue
        # build candidate list
        lines = []
        for idx, item in candidates:
            stem = item.get("problem") or item.get("front") or ""
            stem = stem.replace("\n", " ").strip()
            if len(stem) > 160:
                stem = stem[:160] + "..."
            lines.append(f"{idx}: {stem}")
        prompt = (
            "You are matching a medical exam image to the most relevant question stem. "
            "Choose the single best question index from the list below. "
            "If none match, return -1. Return ONLY the index number.\n\n"
            "Questions:\n" + "\n".join(lines)
        )
        matched_idx = None
        try:
            if ai_model == "🔵 Google Gemini":
                if not api_key:
                    continue
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel(get_gemini_model_id())
                img_bytes = data_uri_to_bytes(img.get("data_uri", ""))
                response = model.generate_content([prompt, img_bytes])
                text = (response.text or "").strip()
            else:
                if not openai_api_key:
                    continue
                client = OpenAI(api_key=openai_api_key)
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "user", "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": img.get("data_uri", "")}},
                        ]}
                    ],
                    temperature=0
                )
                text = (response.choices[0].message.content or "").strip()
            m = re.search(r"-?\\d+", text)
            if m:
                matched_idx = int(m.group(0))
        except Exception:
            matched_idx = None

        if matched_idx is None or matched_idx < 0 or matched_idx >= len(items):
            processed += 1
            continue
        if items[matched_idx].get("images"):
            # avoid overwriting existing images
            processed += 1
            continue
        items[matched_idx].setdefault("images", [])
        items[matched_idx]["images"].append(img.get("data_uri"))
        img["matched"] = True
        processed += 1

    return items

def generate_explanations_ai(items, ai_model, api_key=None, openai_api_key=None, max_items=20):
    if not items or max_items <= 0:
        return items
    count = 0
    for item in items:
        if item.get("explanation"):
            continue
        if count >= max_items:
            break
        stem = item.get("problem") or item.get("front") or ""
        opts = item.get("options") or []
        answer = item.get("answer")
        if item.get("type") == "mcq":
            answer_text = None
            if isinstance(answer, int) and 1 <= answer <= len(opts):
                answer_text = opts[answer - 1]
            prompt = (
                "다음 객관식 문제의 해설을 2~4문장으로 작성하세요. "
                "정답 근거와 핵심 포인트만 간단히 설명하세요.\n\n"
                f"문항: {stem}\n"
                f"선지: {opts}\n"
                f"정답: {answer}"
            )
        else:
            prompt = (
                "다음 주관식/빈칸 문제의 해설을 2~4문장으로 작성하세요. "
                "정답 근거와 핵심 포인트만 간단히 설명하세요.\n\n"
                f"문항: {stem}\n"
                f"정답: {answer}"
            )
        try:
            if ai_model == "🔵 Google Gemini":
                if not api_key:
                    continue
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel(get_gemini_model_id())
                generation_config = {
                    "temperature": LLM_TEMPERATURE,
                    "top_p": 1.0,
                }
                response = model.generate_content(prompt, generation_config=generation_config)
                text = (response.text or "").strip()
                append_audit_log("gen.explanation.batch", {
                    "model": get_gemini_model_id(),
                    "temperature": LLM_TEMPERATURE,
                    "seed": None,
                    "prompt_hash": _hash_text(prompt),
                    "prompt_text": prompt,
                    "output_text": text,
                    "usage_tokens": _gemini_usage_tokens(response),
                    "prompt_version": PROMPT_VERSION,
                })
            else:
                if not openai_api_key:
                    continue
                client = OpenAI(api_key=openai_api_key)
                openai_params = {
                    "model": "gpt-4o-mini",
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": LLM_TEMPERATURE,
                    "max_tokens": 300,
                }
                if LLM_SEED is not None:
                    openai_params["seed"] = LLM_SEED
                response = client.chat.completions.create(**openai_params)
                text = (response.choices[0].message.content or "").strip()
                append_audit_log("gen.explanation.batch", {
                    "model": "gpt-4o-mini",
                    "temperature": LLM_TEMPERATURE,
                    "seed": LLM_SEED,
                    "prompt_hash": _hash_text(prompt),
                    "prompt_text": prompt,
                    "output_text": text,
                    "usage_tokens": _openai_usage_tokens(response),
                    "prompt_version": PROMPT_VERSION,
                })
            if text:
                item["explanation"] = text
                count += 1
        except Exception:
            continue
    return items

def generate_single_explanation_ai(item, ai_model, api_key=None, openai_api_key=None, return_error=False):
    if not item:
        return ("", "빈 문항") if return_error else ""
    stem = item.get("problem") or item.get("front") or item.get("raw") or ""
    opts = item.get("options") or item.get("choices") or []
    answer = item.get("answer")
    if answer is None:
        answer = item.get("correct")
    if item.get("type") == "mcq":
        prompt = (
            "다음 객관식 문제의 해설을 2~4문장으로 작성하세요. "
            "정답 근거와 핵심 포인트만 간단히 설명하세요.\n\n"
            f"문항: {stem}\n"
            f"선지: {opts}\n"
            f"정답: {answer}"
        )
    else:
        prompt = (
            "다음 주관식/빈칸 문제의 해설을 2~4문장으로 작성하세요. "
            "정답 근거와 핵심 포인트만 간단히 설명하세요.\n\n"
            f"문항: {stem}\n"
            f"정답: {answer}"
        )
    try:
        if ai_model == "🔵 Google Gemini":
            if not api_key:
                return ("", "Gemini API 키 없음") if return_error else ""
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(get_gemini_model_id())
            generation_config = {
                "temperature": LLM_TEMPERATURE,
                "top_p": 1.0,
            }
            response = model.generate_content(prompt, generation_config=generation_config)
            text = (response.text or "").strip()
            append_audit_log("gen.explanation.single", {
                "model": get_gemini_model_id(),
                "temperature": LLM_TEMPERATURE,
                "seed": None,
                "prompt_hash": _hash_text(prompt),
                "prompt_text": prompt,
                "output_text": text,
                "usage_tokens": _gemini_usage_tokens(response),
                "prompt_version": PROMPT_VERSION,
            })
            return (text, "") if return_error else text
        else:
            if not openai_api_key:
                return ("", "OpenAI API 키 없음") if return_error else ""
            client = OpenAI(api_key=openai_api_key)
            openai_params = {
                "model": "gpt-4o-mini",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": LLM_TEMPERATURE,
                "max_tokens": 300,
            }
            if LLM_SEED is not None:
                openai_params["seed"] = LLM_SEED
            response = client.chat.completions.create(**openai_params)
            text = (response.choices[0].message.content or "").strip()
            append_audit_log("gen.explanation.single", {
                "model": "gpt-4o-mini",
                "temperature": LLM_TEMPERATURE,
                "seed": LLM_SEED,
                "prompt_hash": _hash_text(prompt),
                "prompt_text": prompt,
                "output_text": text,
                "usage_tokens": _openai_usage_tokens(response),
                "prompt_version": PROMPT_VERSION,
            })
            return (text, "") if return_error else text
    except Exception as e:
        return ("", str(e)) if return_error else ""

def grade_essay_answer_ai(item, user_answer, ai_model, api_key=None, openai_api_key=None):
    question_text = (item.get("front") or item.get("problem") or "").strip()
    reference_answer = (item.get("answer") or "").strip()
    explanation = (item.get("explanation") or "").strip()
    if not question_text or not user_answer:
        return None, "질문 또는 응답이 비어 있습니다."
    prompt = (
        "다음 서술형 답안을 채점하세요. 반드시 JSON으로만 답하세요.\n"
        "JSON 형식: {\"score\": 0-100, \"is_correct\": true/false, \"feedback\": \"...\", \"key_points\": [\"...\"]}\n"
        f"[문항]\n{question_text}\n\n"
        f"[모범답안]\n{reference_answer}\n\n"
        f"[해설]\n{explanation}\n\n"
        f"[학생답안]\n{user_answer}"
    )
    try:
        if ai_model == "🔵 Google Gemini":
            if not api_key:
                return None, "Gemini API 키가 필요합니다."
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(get_gemini_model_id())
            response = model.generate_content(
                prompt,
                generation_config={"temperature": LLM_TEMPERATURE, "top_p": 1.0}
            )
            raw = response.text or ""
            usage_tokens = _gemini_usage_tokens(response)
            model_name = get_gemini_model_id()
        else:
            if not openai_api_key:
                return None, "OpenAI API 키가 필요합니다."
            client = OpenAI(api_key=openai_api_key)
            params = {
                "model": "gpt-4o-mini",
                "messages": [
                    {"role": "system", "content": "의학교육 채점자 역할로 JSON만 출력하세요."},
                    {"role": "user", "content": prompt},
                ],
                "temperature": LLM_TEMPERATURE,
                "max_tokens": 700,
            }
            if LLM_SEED is not None:
                params["seed"] = LLM_SEED
            response = client.chat.completions.create(**params)
            raw = (response.choices[0].message.content or "").strip()
            usage_tokens = _openai_usage_tokens(response)
            model_name = "gpt-4o-mini"
        parsed = _parse_json_from_text(raw)
        if not isinstance(parsed, dict):
            return None, "채점 응답 파싱 실패"
        score = parsed.get("score", 0)
        try:
            score = int(float(score))
        except Exception:
            score = 0
        score = max(0, min(100, score))
        result = {
            "score": score,
            "is_correct": bool(parsed.get("is_correct", False)),
            "feedback": str(parsed.get("feedback", "")).strip(),
            "key_points": parsed.get("key_points", []) if isinstance(parsed.get("key_points", []), list) else [],
        }
        append_audit_log("grade.essay", {
            "model": model_name,
            "temperature": LLM_TEMPERATURE,
            "seed": LLM_SEED if ai_model != "🔵 Google Gemini" else None,
            "prompt_hash": _hash_text(prompt),
            "prompt_text": prompt,
            "output_text": raw,
            "usage_tokens": usage_tokens,
            "grader_version": GRADER_VERSION,
        })
        return result, ""
    except Exception as e:
        return None, str(e)

def update_question_explanation(q_id, explanation_text):
    if not q_id:
        return False
    bank = load_questions()
    for key in ("text", "cloze"):
        for item in bank.get(key, []):
            if item.get("id") == q_id:
                item["explanation"] = explanation_text
                save_questions(bank)
                return True
    return False

def _extract_json_candidates(raw):
    if not raw:
        return []
    raw = raw.strip()
    candidates = []
    fence = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", raw)
    if fence:
        candidates.append(fence.group(1).strip())
    candidates.append(raw)
    arr = re.search(r"\[\s*\{[\s\S]+?\}\s*\]", raw)
    if arr:
        candidates.append(arr.group(0))
    obj = re.search(r"\{[\s\S]+\}", raw)
    if obj:
        candidates.append(obj.group(0))
    return candidates

def _parse_json_from_text(raw):
    for cand in _extract_json_candidates(raw):
        try:
            data = json.loads(cand)
            return data
        except Exception:
            continue
    return None

def ai_parse_exam_layout(left_text, right_text, ai_model, api_key=None, openai_api_key=None, hint_text=""):
    if not left_text or len(left_text.strip()) < 20:
        return []
    prompt = (
        "아래 LEFT/RIGHT 텍스트에서 시험 문항을 JSON 배열로 추출하세요. 오직 JSON만 출력하세요.\n"
        "LEFT에는 문항/선지가 있고, RIGHT에는 정답/해설(또는 요약)이 있습니다.\n"
        "RIGHT는 '▶ ⑤' 또는 '정답: ⑤' 같은 형식일 수 있으니 이를 정답으로 사용하세요.\n"
        "문항 번호가 보이면 qnum에 넣고, 없으면 순서대로 매칭하세요.\n"
        "형식:\n"
        "{\n"
        "  \"type\": \"mcq\" 또는 \"cloze\",\n"
        "  \"problem\": (mcq용 질문 본문),\n"
        "  \"front\": (cloze용 질문 본문),\n"
        "  \"options\": [\"선지1\", \"선지2\", ...] (mcq일 때만),\n"
        "  \"answer\": 정답 (mcq는 1-5 정수, cloze는 문자열),\n"
        "  \"explanation\": 해설(없으면 \"\"),\n"
        "  \"qnum\": 문항 번호(있으면 숫자)\n"
        "}\n"
        "[LEFT]\n"
    )
    if hint_text:
        prompt = f"[문서 구조 힌트]\n{hint_text}\n\n" + prompt
    prompt += left_text[:20000] + "\n\n[RIGHT]\n" + (right_text[:20000] if right_text else "")
    try:
        if ai_model == "🔵 Google Gemini":
            if not api_key:
                return []
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(get_gemini_model_id())
            response = model.generate_content(prompt)
            raw = response.text or ""
        else:
            if not openai_api_key:
                return []
            client = OpenAI(api_key=openai_api_key)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.2,
                max_tokens=4000
            )
            raw = response.choices[0].message.content or ""
        data = _parse_json_from_text(raw)
        if isinstance(data, dict):
            data = data.get("items") or data.get("questions") or data.get("data") or []
        if not isinstance(data, list):
            return []
        return clean_parsed_items(data)
    except Exception:
        return []

def ai_parse_exam_text(text, ai_model, api_key=None, openai_api_key=None, max_items=60, hint_text="", return_raw=False):
    if not text or len(text.strip()) < 20:
        return ([], "") if return_raw else []
    prompt = (
        "아래 텍스트에서 시험 문항을 JSON 배열로 추출하세요. 오직 JSON만 출력하세요.\n"
        "각 항목 형식:\n"
        "{\n"
        "  \"type\": \"mcq\" 또는 \"cloze\",\n"
        "  \"problem\": (mcq용 질문 본문),\n"
        "  \"front\": (cloze용 질문 본문),\n"
        "  \"options\": [\"선지1\", \"선지2\", ...] (mcq일 때만),\n"
        "  \"answer\": 정답 (mcq는 1-5 정수, cloze는 문자열),\n"
        "  \"explanation\": 해설(없으면 \"\"),\n"
        "  \"page\": 페이지 번호(텍스트에 '=== 페이지 N ===' 표기가 있으면 활용),\n"
        "  \"qnum\": 문항 번호(있으면 숫자)\n"
        "}\n"
        f"최대 {max_items}개까지만 출력하세요.\n"
        "문항이 겹치지 않도록 정확히 분리하세요.\n\n"
        "[원문]\n"
    )
    if hint_text:
        prompt = f"[문서 구조 힌트]\n{hint_text}\n\n" + prompt
    try:
        if ai_model == "🔵 Google Gemini":
            if not api_key:
                return ([], "") if return_raw else []
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(get_gemini_model_id())
            response = model.generate_content(prompt + text[:30000])
            raw = response.text or ""
        else:
            if not openai_api_key:
                return ([], "") if return_raw else []
            client = OpenAI(api_key=openai_api_key)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt + text[:30000]}],
                temperature=0.2,
                max_tokens=4000
            )
            raw = response.choices[0].message.content or ""

        data = _parse_json_from_text(raw)
        if data is None:
            return ([], raw) if return_raw else []
        if isinstance(data, dict):
            data = data.get("items") or data.get("questions") or data.get("data") or []
        if not isinstance(data, list):
            return ([], raw) if return_raw else []
        items = clean_parsed_items(data)
        return (items, raw) if return_raw else items
    except Exception:
        return ([], "") if return_raw else []

def ai_parse_exam_block(block_text, ai_model, api_key=None, openai_api_key=None, hint_text="", return_raw=False):
    if not block_text or len(block_text.strip()) < 10:
        return (None, "") if return_raw else None
    prompt = (
        "아래 텍스트에서 문항 1개를 JSON 객체로 추출하세요. 오직 JSON만 출력하세요.\n"
        "형식:\n"
        "{\n"
        "  \"type\": \"mcq\" 또는 \"cloze\",\n"
        "  \"problem\": (mcq용 질문 본문),\n"
        "  \"front\": (cloze용 질문 본문),\n"
        "  \"options\": [\"선지1\", \"선지2\", ...] (mcq일 때만),\n"
        "  \"answer\": 정답 (mcq는 1-5 정수, cloze는 문자열),\n"
        "  \"explanation\": 해설(없으면 \"\")\n"
        "}\n"
    )
    if hint_text:
        prompt += f"\n[문서 구조 힌트]\n{hint_text}\n"
    prompt += "\n[원문]\n"
    try:
        if ai_model == "🔵 Google Gemini":
            if not api_key:
                return (None, "") if return_raw else None
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(get_gemini_model_id())
            response = model.generate_content(prompt + block_text[:15000])
            raw = response.text or ""
        else:
            if not openai_api_key:
                return (None, "") if return_raw else None
            client = OpenAI(api_key=openai_api_key)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt + block_text[:15000]}],
                temperature=0.2,
                max_tokens=1200
            )
            raw = response.choices[0].message.content or ""
        data = _parse_json_from_text(raw)
        if not isinstance(data, dict):
            return (None, raw) if return_raw else None
        items = clean_parsed_items([data])
        item = items[0] if items else None
        return (item, raw) if return_raw else item
    except Exception:
        return (None, "") if return_raw else None

def should_attach_image(item):
    text = (item.get("problem") or item.get("front") or "")
    text = text.lower()
    keywords = [
        "x-ray", "xray", "ct", "mri", "us", "ultrasound", "sonography", "radiograph",
        "영상", "영상소견", "영상 소견", "사진", "그림", "figure", "fig.", "영상에서", "사진을 보고", "영상학적"
    ]
    return any(k in text for k in keywords)

def auto_attach_images_to_items(items, images, strategy="page", max_per_question=1, anchors=None, min_score=0.2, only_if_keyword=False):
    if not items or not images:
        return items
    if max_per_question < 1:
        return items

    if strategy == "sequential":
        img_idx = 0
        for item in items:
            if item.get("images"):
                continue
            attach = []
            for _ in range(max_per_question):
                if img_idx >= len(images):
                    break
                attach.append(images[img_idx]["data_uri"])
                img_idx += 1
            if attach:
                item["images"] = attach
        return items

    if strategy == "layout" and anchors:
        # build intervals per page: [qnum, start_y, end_y)
        intervals = {}
        for page, arr in anchors.items():
            if not arr:
                continue
            arr_sorted = sorted(arr, key=lambda x: x["y"])
            page_intervals = []
            for idx, a in enumerate(arr_sorted):
                start = a["y"]
                end = arr_sorted[idx + 1]["y"] if idx + 1 < len(arr_sorted) else float("inf")
                page_intervals.append({"qnum": a["qnum"], "start": start, "end": end})
            intervals[page] = page_intervals

        image_map = {}
        for img in images:
            page = img.get("page")
            y = img.get("y")
            if page not in intervals or y is None:
                continue
            for seg in intervals[page]:
                if seg["start"] <= y < seg["end"]:
                    key = (page, seg["qnum"])
                    image_map.setdefault(key, []).append(img["data_uri"])
                    break

        for item in items:
            if item.get("images"):
                continue
            if only_if_keyword and not should_attach_image(item):
                continue
            page = item.get("page")
            qnum = item.get("qnum")
            if page is None or qnum is None:
                continue
            key = (page, qnum)
            imgs = image_map.get(key) or []
            if imgs:
                item["images"] = imgs[:max_per_question]
        return items

    if strategy == "page":
        page_to_images = {}
        for img in images:
            page = img.get("page")
            page_to_images.setdefault(page, []).append(img["data_uri"])
        for item in items:
            if item.get("images"):
                continue
            if only_if_keyword and not should_attach_image(item):
                continue
            page = item.get("page")
            candidates = page_to_images.get(page) or []
            if candidates:
                item["images"] = candidates[:max_per_question]
        return items

    if strategy == "ocr":
        # build token sets per item
        item_tokens = []
        for item in items:
            text = " ".join([
                item.get("problem") or item.get("front") or "",
                " ".join(item.get("options", []) or []),
                item.get("explanation") or ""
            ])
            item_tokens.append(_tokenize_for_match(text))

        def item_key(i):
            return f"{items[i].get('page')}_{items[i].get('qnum')}_{i}"

        attached = {}
        for i, item in enumerate(items):
            attached[item_key(i)] = list(item.get("images", [])) if item.get("images") else []

        for img in images:
            ocr_text = img.get("ocr_text", "") or ""
            tokens_img = _tokenize_for_match(ocr_text)
            if not tokens_img:
                continue
            best_idx = None
            best_score = 0.0
            for i, tokens in enumerate(item_tokens):
                if not tokens:
                    continue
                if only_if_keyword and not should_attach_image(items[i]):
                    continue
                # prefer same page if available
                if img.get("page") and items[i].get("page") and img.get("page") != items[i].get("page"):
                    continue
                overlap = len(tokens_img & tokens) / max(1, len(tokens_img))
                if overlap > best_score:
                    best_score = overlap
                    best_idx = i
            if best_idx is None or best_score < min_score:
                continue
            key = item_key(best_idx)
            if img["data_uri"] in attached[key]:
                continue
            if len(attached[key]) >= max_per_question:
                continue
            attached[key].append(img["data_uri"])

        for i, item in enumerate(items):
            key = item_key(i)
            if attached.get(key):
                item["images"] = attached[key]
        return items

    return items

def extract_text_from_pdf(
    uploaded_file,
    enable_ocr=True,
    ocr_engine="auto",
    ocr_langs=("ko", "en"),
    ocr_max_pages=0,
    min_text_len=40,
    include_page_markers=False,
    ai_fallback=False,
    ai_model="",
    api_key=None,
    openai_api_key=None,
):
    """PDF에서 텍스트 추출"""
    try:
        pdf_bytes = uploaded_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        text = ""
        for i, page in enumerate(doc):
            page_text = page.get_text()
            if include_page_markers:
                text += f"=== 페이지 {i + 1} ===\n"
            text += page_text
            if include_page_markers:
                text += "\n"
        doc.close()
        if len(text.strip()) >= min_text_len:
            return text
        # OCR fallback (스캔 PDF 등)
        if not (enable_ocr or ai_fallback):
            return text
        if ai_fallback and ai_model:
            ai_text = ocr_pdf_bytes_with_ai(
                pdf_bytes,
                ai_model=ai_model,
                api_key=api_key,
                openai_api_key=openai_api_key,
                max_pages=ocr_max_pages,
                zoom=2.0,
            )
            if ai_text.strip():
                return ai_text
        engines = available_ocr_engines()
        if not engines:
            return text
        try:
            engine = engines[0] if ocr_engine == "auto" else ocr_engine
            ocr_text = ocr_pdf_bytes(pdf_bytes, engine=engine, langs=ocr_langs, max_pages=ocr_max_pages)
            return ocr_text if ocr_text.strip() else text
        except Exception:
            return text
    except Exception as e:
        raise ValueError(f"PDF 처리 실패: {str(e)}")

def extract_text_from_docx(uploaded_file):
    """Word (.docx)에서 텍스트 추출"""
    try:
        doc = Document(uploaded_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + "\n"
        return text
    except Exception as e:
        raise ValueError(f"Word 문서 처리 실패: {str(e)}")

def extract_text_from_pptx(uploaded_file):
    """PowerPoint (.pptx)에서 텍스트 추출"""
    try:
        prs = Presentation(uploaded_file)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n=== 슬라이드 {slide_num} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise ValueError(f"PowerPoint 처리 실패: {str(e)}")

def extract_text_from_hwp(uploaded_file):
    """HWP (.hwp)에서 텍스트 추출 (hwp5txt 필요)"""
    tmp_path = None
    try:
        if hasattr(uploaded_file, "read"):
            data = uploaded_file.read()
        else:
            data = uploaded_file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".hwp") as tmp:
            tmp.write(data)
            tmp_path = tmp.name

        def is_table_placeholder_text(text):
            if not text or not text.strip():
                return True
            placeholder_count = text.count("<표>")
            if placeholder_count >= 3:
                cleaned = re.sub(r"<표>", "", text)
                cleaned = re.sub(r"\s+", "", cleaned)
                if len(cleaned) < 80:
                    return True
                if not re.search(r"[①②③④⑤]|\\b정답\\b|\\b답\\b", text):
                    return True
            return False

        def extract_text_from_odt_content(xml_bytes):
            try:
                root = ET.fromstring(xml_bytes)
            except Exception:
                return ""
            ns = {
                "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
                "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
                "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
                "draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
            }
            body = root.find("office:body/office:text", ns)
            if body is None:
                return ""

            def normalize_line(line):
                line = line.replace("\u00a0", " ")
                line = re.sub(r"[ \t]+", " ", line).strip()
                return line

            def cell_lines(cell):
                lines = []
                for p in cell.findall(".//text:p", ns) + cell.findall(".//text:h", ns):
                    line = normalize_line("".join(p.itertext()))
                    if line:
                        lines.append(line)
                img_count = len(cell.findall(".//draw:image", ns))
                if img_count:
                    lines.append(f"[이미지 x{img_count}]")
                return lines

            out_lines = []
            for child in body:
                if child.tag == f"{{{ns['table']}}}table":
                    for row in child.findall("table:table-row", ns):
                        row_lines = []
                        for cell in row.findall("table:table-cell", ns):
                            lines = cell_lines(cell)
                            if lines:
                                row_lines.append("\n".join(lines))
                        if row_lines:
                            out_lines.append("\n".join(row_lines))
                            out_lines.append("")
                elif child.tag in (f"{{{ns['text']}}}p", f"{{{ns['text']}}}h"):
                    line = normalize_line("".join(child.itertext()))
                    if line:
                        out_lines.append(line)
            return "\n".join(out_lines).strip()

        def extract_text_from_hwp5odt(path):
            odt_path = None
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".odt") as tmp_odt:
                    odt_path = tmp_odt.name
                def run_odt(cmd):
                    result = subprocess.run(cmd, capture_output=True, text=True)
                    if result.returncode != 0:
                        raise ValueError(result.stderr.strip() or "hwp5odt 변환 실패")
                    if not os.path.exists(odt_path) or os.path.getsize(odt_path) == 0:
                        raise ValueError("ODT 변환 결과가 비어있습니다.")
                if shutil.which("hwp5odt"):
                    run_odt(["hwp5odt", "--output", odt_path, path])
                else:
                    try:
                        import importlib.util
                        if importlib.util.find_spec("hwp5.hwp5odt") is not None:
                            run_odt([sys.executable, "-m", "hwp5.hwp5odt", "--output", odt_path, path])
                        else:
                            return ""
                    except Exception:
                        return ""
                with zipfile.ZipFile(odt_path) as zf:
                    xml_bytes = zf.read("content.xml")
                return extract_text_from_odt_content(xml_bytes)
            except Exception:
                # ODT 보조 파싱 실패 시 hwp5txt 원문으로 안전하게 폴백
                return ""
            finally:
                if odt_path and os.path.exists(odt_path):
                    try:
                        os.unlink(odt_path)
                    except Exception:
                        pass

        def run_hwp5txt(cmd):
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True
            )
            if result.returncode != 0:
                raise ValueError(result.stderr.strip() or "hwp5txt 변환 실패")
            text = result.stdout
            if not text.strip():
                raise ValueError("HWP 텍스트가 비어있습니다.")
            return text

        if shutil.which("hwp5txt"):
            text = run_hwp5txt(["hwp5txt", tmp_path])
            if not is_table_placeholder_text(text):
                return text
            odt_text = extract_text_from_hwp5odt(tmp_path)
            if odt_text:
                return odt_text
            return text

        # fallback: python -m hwp5.hwp5txt (pyhwp 설치되어 있으나 PATH에 없을 때)
        try:
            import importlib.util
            if importlib.util.find_spec("hwp5.hwp5txt") is not None:
                text = run_hwp5txt([sys.executable, "-m", "hwp5.hwp5txt", tmp_path])
                if not is_table_placeholder_text(text):
                    return text
                odt_text = extract_text_from_hwp5odt(tmp_path)
                if odt_text:
                    return odt_text
                return text
        except Exception:
            pass

        raise ValueError(
            "hwp5txt 실행 파일을 찾을 수 없습니다. "
            "pyhwp 설치 후 다시 시도하세요. (예: `python -m pip install pyhwp`)"
        )
    except Exception as e:
        raise ValueError(f"HWP 처리 실패: {str(e)}")
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

def extract_text_from_file(uploaded_file, **kwargs):
    """파일 형식에 따라 자동으로 텍스트 추출"""
    file_ext = Path(uploaded_file.name).suffix.lower()
    
    if file_ext == ".pdf":
        return extract_text_from_pdf(uploaded_file, **kwargs)
    elif file_ext == ".docx":
        return extract_text_from_docx(uploaded_file)
    elif file_ext == ".pptx":
        return extract_text_from_pptx(uploaded_file)
    elif file_ext == ".hwp":
        return extract_text_from_hwp(uploaded_file)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {file_ext}")

def parse_uploaded_question_file(uploaded_file, mode_hint="auto"):
    """사용자 업로드 문항 파일 파싱 (json/txt/tsv)"""
    ext = Path(uploaded_file.name).suffix.lower()
    content_bytes = uploaded_file.read()
    if ext == ".json":
        try:
            data = json.loads(content_bytes.decode("utf-8"))
        except Exception:
            data = json.loads(content_bytes.decode("utf-8-sig"))
        items = []
        if isinstance(data, dict) and ("text" in data or "cloze" in data):
            for it in data.get("text", []):
                norm = normalize_mcq_item(it)
                if norm:
                    items.append(norm)
            for it in data.get("cloze", []):
                norm = normalize_cloze_item(it)
                if norm:
                    items.append(norm)
        elif isinstance(data, list):
            for it in data:
                if isinstance(it, dict) and (it.get("type") == "cloze" or "front" in it or ("content" in it and "{{c1::" in str(it.get("content", "")))):
                    norm = normalize_cloze_item(it)
                else:
                    norm = normalize_mcq_item(it)
                if norm:
                    items.append(norm)
        elif isinstance(data, dict):
            if data.get("type") == "cloze" or "front" in data or ("content" in data and "{{c1::" in str(data.get("content", ""))):
                norm = normalize_cloze_item(data)
            else:
                norm = normalize_mcq_item(data)
            if norm:
                items.append(norm)
        return items

    # text/tsv/hwp
    if ext == ".hwp":
        text = extract_text_from_hwp(content_bytes)
    else:
        text = content_bytes.decode("utf-8", errors="ignore")
    if mode_hint == "auto":
        if "{{c1::" in text:
            mode_hint = MODE_CLOZE
        elif "정답" in text and not re.search(r"①|②|③|④|⑤", text):
            mode_hint = MODE_CLOZE
        else:
            mode_hint = MODE_MCQ

    if mode_hint == MODE_CLOZE and "{{c1::" not in text:
        qa_parsed = parse_qa_to_cloze(text)
        if qa_parsed:
            return qa_parsed
    parsed = parse_generated_text_to_structured(text, mode_hint)
    if isinstance(parsed, list) and parsed:
        return parsed
    # fallback: fuzzy parser for messy past exam text
    fuzzy = parse_exam_text_fuzzy(text)
    return fuzzy if isinstance(fuzzy, list) else []

# ============================================================================
# AI 콘텐츠 생성
# ============================================================================
def detect_term_language_mode(style_text: str):
    """기출 스타일 텍스트에서 '용어 표기' 혼용을 추정한다.

    Returns:
        (mode, pattern)
        - mode: "ko" | "en" | "mixed"
        - pattern: "ko(en)" | "en(ko)" | ""
    """
    s = str(style_text or "").strip()
    if not s:
        return ("mixed", "")

    # Pattern-based mixed style detection
    ko_en = len(re.findall(r"[가-힣]{2,}\s*\([A-Za-z][^)]{2,}\)", s))
    en_ko = len(re.findall(r"[A-Za-z]{2,}(?:[ -][A-Za-z]{2,})*\s*\([가-힣]{2,}[^)]*\)", s))
    if ko_en or en_ko:
        if ko_en >= en_ko and ko_en > 0:
            return ("mixed", "ko(en)")
        if en_ko > ko_en:
            return ("mixed", "en(ko)")
        return ("mixed", "")

    # Fallback: character ratio (Hangul vs Latin)
    hangul_chars = len(re.findall(r"[가-힣]", s))
    latin_chars = len(re.findall(r"[A-Za-z]", s))
    denom = hangul_chars + latin_chars
    if denom == 0:
        return ("mixed", "")

    ratio_ko = hangul_chars / denom
    ratio_en = latin_chars / denom
    if ratio_ko >= 0.85:
        return ("ko", "")
    if ratio_en >= 0.85:
        return ("en", "")
    return ("mixed", "")

def detect_question_flavor_scores(text):
    s = str(text or "")
    if not s.strip():
        return {"basic": 0, "case": 0}
    case_patterns = [
        r"\d+\s*세",
        r"환자",
        r"내원",
        r"주호소",
        r"증상",
        r"진단",
        r"치료",
        r"처치",
        r"검사",
        r"혈압",
        r"맥박",
        r"호흡수",
        r"체온",
        r"응급",
    ]
    basic_patterns = [
        r"기전",
        r"정의",
        r"분류",
        r"구조",
        r"위치",
        r"유래",
        r"발생",
        r"막전위",
        r"이온",
        r"효소",
        r"대사",
        r"경로",
        r"수송",
        r"계산",
        r"equation",
        r"pathway",
        r"origin",
    ]
    case_score = sum(len(re.findall(p, s, flags=re.IGNORECASE)) for p in case_patterns)
    basic_score = sum(len(re.findall(p, s, flags=re.IGNORECASE)) for p in basic_patterns)
    return {"basic": basic_score, "case": case_score}

def resolve_generation_flavor(flavor_choice, raw_text="", style_text="", subject=""):
    choice = str(flavor_choice or "").strip().lower()
    if "basic" in choice:
        return "basic"
    if "case" in choice:
        return "case"
    if "mix" in choice:
        return "mix"

    subj = str(subject or "").lower()
    basic_subject_keywords = ["해부", "생리", "생화학", "면역", "발생", "조직", "약리", "기초", "anatom", "physio", "biochem", "immun"]
    case_subject_keywords = ["내과", "외과", "소아", "산부", "정신", "응급", "가정", "신경과", "진단", "임상", "internal", "surgery", "pedi", "obgyn"]

    basic_subj = any(k in subj for k in basic_subject_keywords)
    case_subj = any(k in subj for k in case_subject_keywords)
    if basic_subj and not case_subj:
        return "basic"
    if case_subj and not basic_subj:
        return "case"

    style_scores = detect_question_flavor_scores(str(style_text or "")[:12000])
    raw_scores = detect_question_flavor_scores(str(raw_text or "")[:12000])
    basic_score = style_scores["basic"] * 2 + raw_scores["basic"]
    case_score = style_scores["case"] * 2 + raw_scores["case"]
    return "basic" if basic_score >= case_score else "case"

def build_flavor_instructions(selected_mode, resolved_flavor, mix_basic_ratio=70):
    flavor = str(resolved_flavor or "").lower()
    if flavor not in {"basic", "case", "mix"}:
        return ""
    basic_ratio = max(0, min(100, int(mix_basic_ratio or 70)))
    case_ratio = 100 - basic_ratio
    if flavor == "basic":
        return """
[문항 성격 지시: 기초의학형]
- 임상 진단/처방 중심 증례형을 배제하세요.
- 기전(Mechanism), 해부학적 위치/주행, 분류, 정의, 계산(공식/수치 해석) 중심으로 출제하세요.
- 한국어 설명 + 핵심 의학 용어는 영어(또는 한영 병기)로 작성하세요.
"""
    if flavor == "case":
        return """
[문항 성격 지시: 케이스형]
- 환자 정보(연령/성별/증상/검사 소견)를 포함한 임상 증례형으로 출제하세요.
- 진단, 다음 검사, 치료 선택/금기 판단을 중심으로 구성하세요.
- 단순 정의 암기형 문항 비율을 낮추세요.
"""
    return f"""
[문항 성격 지시: 혼합형]
- 전체 문항을 기초의학형 약 {basic_ratio}%, 케이스형 약 {case_ratio}% 비율로 구성하세요.
- 기초의학형: 기전/해부학/분류/계산 중심
- 케이스형: 임상 증례 기반 진단/검사/치료 판단 중심
"""

def build_style_instructions(style_text):
    if not style_text:
        return ""
    excerpt = style_text[:8000]
    mode, pattern = detect_term_language_mode(style_text)
    if mode == "ko":
        term_rule = "- 용어 표기: 가능한 한 한국어 용어를 사용(표준 약어/단위는 허용). 영어 풀네임 병기는 최소화."
    elif mode == "en":
        term_rule = "- 용어 표기: 핵심 의학 용어는 영어로 표기. 불필요한 한국어 번역/병기는 최소화."
    else:
        if pattern == "ko(en)":
            term_rule = "- 용어 표기: 한국어 용어 뒤에 (영어)로 병기하는 스타일을 유지. 예: 노신경(radial nerve)"
        elif pattern == "en(ko)":
            term_rule = "- 용어 표기: 영어 용어 뒤에 (한국어)로 병기하는 스타일을 유지."
        else:
            term_rule = "- 용어 표기: 한국어/영어 혼용 스타일을 유지(기출문제 표현 우선)."
    return f"""
[기출문제 스타일 참고]
{excerpt}

[스타일 지시]
- 위 기출문제의 질문 구조, 난이도, 문장 길이, 선지 톤/표현을 최대한 모사
- 내용은 강의록 기반으로 생성
- 출력 형식 규칙은 반드시 유지
{term_rule}
"""

def generate_content_gemini(
    text_content,
    selected_mode,
    num_items=5,
    api_key=None,
    style_text=None,
    gemini_model_id=None,
    audit_user_id=None,
    resolved_flavor=None,
    mix_basic_ratio=70,
):
    """Gemini를 이용해 콘텐츠 생성"""
    if not api_key:
        return "⚠️ 왼쪽 사이드바에 Gemini API 키를 먼저 입력해주세요."
    
    if not text_content or len(text_content.strip()) < 10:
        return "⚠️ 추출된 텍스트가 너무 짧습니다. 다시 시도해주세요."
    
    mode_mcq = globals().get("MODE_MCQ", "📝 객관식 문제 (Case Study)")
    mode_cloze = globals().get("MODE_CLOZE", "🧩 빈칸 뚫기 (Anki Cloze)")
    mode_short = globals().get("MODE_SHORT", "🧠 단답형 문제")
    prompt_short = globals().get("PROMPT_SHORT", PROMPT_CLOZE)
    prompt_essay = globals().get("PROMPT_ESSAY", PROMPT_CLOZE)
    style_builder = globals().get("build_style_instructions", lambda _style_text: "")
    flavor_builder = globals().get("build_flavor_instructions", lambda *_args, **_kwargs: "")
    style_block = style_builder(style_text)
    flavor_block = flavor_builder(selected_mode, resolved_flavor, mix_basic_ratio=mix_basic_ratio)
    if selected_mode == mode_mcq:
        system_prompt = PROMPT_MCQ.replace("5문제", f"{num_items}문제") + style_block + flavor_block
    elif selected_mode == mode_cloze:
        system_prompt = PROMPT_CLOZE + style_block + flavor_block + f"\n\n[요청] 총 {num_items}개 항목을 출력하세요. 한 줄에 하나의 항목만 작성하세요."
    elif selected_mode == mode_short:
        system_prompt = prompt_short + style_block + flavor_block + f"\n\n[요청] 총 {num_items}개 항목을 출력하세요."
    else:
        system_prompt = prompt_essay + style_block + flavor_block + f"\n\n[요청] 총 {num_items}개 항목을 출력하세요."
    
    try:
        model_name = gemini_model_id or get_gemini_model_id()
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(model_name)
        prompt_text = f"{system_prompt}\n\n[강의록 내용]:\n{text_content[:30000]}"
        generation_config = {
            "temperature": LLM_TEMPERATURE,
            "top_p": 1.0,
        }
        response = model.generate_content(prompt_text, generation_config=generation_config)
        result_text = response.text
        payload = {
            "model": model_name,
            "temperature": LLM_TEMPERATURE,
            "seed": None,
            "prompt_hash": _hash_text(prompt_text),
            "prompt_text": prompt_text,
            "input_hash": _hash_text(text_content[:30000]),
            "output_text": result_text,
            "usage_tokens": _gemini_usage_tokens(response),
            "prompt_version": PROMPT_VERSION,
        }
        try:
            append_audit_log("gen.question", payload, user_id=audit_user_id)
        except TypeError:
            append_audit_log("gen.question", payload)
        return result_text
    except Exception as e:
        return f"❌ Gemini 생성 실패: {str(e)}"

def generate_content_openai(
    text_content,
    selected_mode,
    num_items=5,
    openai_api_key=None,
    style_text=None,
    audit_user_id=None,
    resolved_flavor=None,
    mix_basic_ratio=70,
):
    """ChatGPT를 이용해 콘텐츠 생성"""
    if not openai_api_key:
        return "⚠️ 왼쪽 사이드바에 OpenAI API 키를 먼저 입력해주세요."
    
    if not text_content or len(text_content.strip()) < 10:
        return "⚠️ 추출된 텍스트가 너무 짧습니다. 다시 시도해주세요."
    
    mode_mcq = globals().get("MODE_MCQ", "📝 객관식 문제 (Case Study)")
    mode_cloze = globals().get("MODE_CLOZE", "🧩 빈칸 뚫기 (Anki Cloze)")
    mode_short = globals().get("MODE_SHORT", "🧠 단답형 문제")
    prompt_short = globals().get("PROMPT_SHORT", PROMPT_CLOZE)
    prompt_essay = globals().get("PROMPT_ESSAY", PROMPT_CLOZE)
    style_builder = globals().get("build_style_instructions", lambda _style_text: "")
    flavor_builder = globals().get("build_flavor_instructions", lambda *_args, **_kwargs: "")
    style_block = style_builder(style_text)
    flavor_block = flavor_builder(selected_mode, resolved_flavor, mix_basic_ratio=mix_basic_ratio)
    if selected_mode == mode_mcq:
        system_prompt = PROMPT_MCQ.replace("5문제", f"{num_items}문제") + style_block + flavor_block
    elif selected_mode == mode_cloze:
        system_prompt = PROMPT_CLOZE + style_block + flavor_block + f"\n\n[요청] 총 {num_items}개 항목을 출력하세요. 한 줄에 하나의 항목만 작성하세요."
    elif selected_mode == mode_short:
        system_prompt = prompt_short + style_block + flavor_block + f"\n\n[요청] 총 {num_items}개 항목을 출력하세요."
    else:
        system_prompt = prompt_essay + style_block + flavor_block + f"\n\n[요청] 총 {num_items}개 항목을 출력하세요."
    
    try:
        import sys
        print(f"[OPENAI DEBUG] API 키 길이: {len(openai_api_key)}", file=sys.stderr)
        print(f"[OPENAI DEBUG] 텍스트 길이: {len(text_content[:30000])}", file=sys.stderr)
        
        openai_client = OpenAI(api_key=openai_api_key)
        print(f"[OPENAI DEBUG] OpenAI 클라이언트 생성 완료", file=sys.stderr)
        
        prompt_text = f"{system_prompt}\n\n[강의록 내용]:\n{text_content[:30000]}"
        openai_params = {
            "model": "gpt-4o-mini",
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"[강의록 내용]:\n{text_content[:30000]}"}
            ],
            "temperature": LLM_TEMPERATURE,
            "max_tokens": 4000,
        }
        if LLM_SEED is not None:
            openai_params["seed"] = LLM_SEED
        response = openai_client.chat.completions.create(
            **openai_params
        )
        
        result = response.choices[0].message.content
        print(f"[OPENAI DEBUG] 응답 길이: {len(result)}", file=sys.stderr)
        
        # MCQ는 JSON으로 파싱, Cloze는 그대로 반환
        if selected_mode == mode_mcq:
            result = convert_json_mcq_to_text(result, num_items)
        
        payload = {
            "model": "gpt-4o-mini",
            "temperature": LLM_TEMPERATURE,
            "seed": LLM_SEED,
            "prompt_hash": _hash_text(prompt_text),
            "prompt_text": prompt_text,
            "input_hash": _hash_text(text_content[:30000]),
            "output_text": result,
            "usage_tokens": _openai_usage_tokens(response),
            "prompt_version": PROMPT_VERSION,
        }
        try:
            append_audit_log("gen.question", payload, user_id=audit_user_id)
        except TypeError:
            append_audit_log("gen.question", payload)
        return result
    except Exception as e:
        import traceback
        error_msg = f"❌ ChatGPT 생성 실패: {str(e)}\n\n스택 트레이스:\n{traceback.format_exc()}"
        print(error_msg, file=sys.stderr)
        return error_msg

def convert_json_mcq_to_text(json_text, num_items):
    """JSON 형식의 MCQ를 기존 텍스트 형식으로 변환"""
    import json
    import sys
    
    try:
        # JSON 파싱
        data = json.loads(json_text)
        if not isinstance(data, list):
            data = [data]
        
        print(f"[JSON PARSE] {len(data)}개 MCQ 파싱 성공", file=sys.stderr)
        
        # 텍스트 형식으로 변환
        result_lines = []
        for idx, item in enumerate(data[:num_items], 1):
            problem = item.get("problem", f"[문제] {idx}번")
            options = item.get("options", [])
            answer = item.get("answer", 1)  # 1~5 숫자
            explanation = item.get("explanation", "")
            
            # problem에 [문제]가 없으면 추가
            if "[문제]" not in problem:
                problem = f"[문제] {problem}"
            
            # MCQ 블록 구성
            block = problem + "\n\n"
            circ = ['①', '②', '③', '④', '⑤']
            for i, opt in enumerate(options[:5]):
                block += f"{circ[i]} {opt}\n"
            
            # 정답과 설명 추가
            ans_num = str(answer) if isinstance(answer, int) and 1 <= answer <= 5 else "1"
            block += f"\n정답: {{{{c1::{ans_num}}}}}\n해설: {explanation}"
            
            result_lines.append(block)
        
        # '---'으로 구분
        final_result = "\n---\n".join(result_lines)
        print(f"[JSON CONVERT] {len(result_lines)}개 MCQ 변환 완료", file=sys.stderr)
        
        return final_result
    
    except json.JSONDecodeError as e:
        print(f"[JSON ERROR] JSON 파싱 실패: {str(e)}", file=sys.stderr)
        # JSON 파싱 실패시 원본 반환 (다른 파싱 로직이 처리할 것)
        return json_text
    except Exception as e:
        print(f"[CONVERT ERROR] 변환 실패: {str(e)}", file=sys.stderr)
        return json_text


def generate_content(
    text_content,
    selected_mode,
    ai_model,
    num_items=5,
    api_key=None,
    openai_api_key=None,
    style_text=None,
    gemini_model_id=None,
    audit_user_id=None,
    resolved_flavor=None,
    mix_basic_ratio=70,
):
    """선택된 AI 모델을 사용해 콘텐츠 생성"""
    if ai_model == "🔵 Google Gemini":
        return generate_content_gemini(
            text_content,
            selected_mode,
            num_items=num_items,
            api_key=api_key,
            style_text=style_text,
            gemini_model_id=gemini_model_id,
            audit_user_id=audit_user_id,
            resolved_flavor=resolved_flavor,
            mix_basic_ratio=mix_basic_ratio,
        )
    else:  # ChatGPT
        return generate_content_openai(
            text_content,
            selected_mode,
            num_items=num_items,
            openai_api_key=openai_api_key,
            style_text=style_text,
            audit_user_id=audit_user_id,
            resolved_flavor=resolved_flavor,
            mix_basic_ratio=mix_basic_ratio,
        )

def split_text_into_chunks(text, chunk_size=8000, overlap=500):
    """문자 단위로 텍스트를 분할 (중첩 포함)"""
    if chunk_size <= 0:
        return [text]
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        if end >= text_len:
            break
        start = end - overlap if end - overlap > start else end
    return chunks

def generate_content_in_chunks(
    text_content,
    selected_mode,
    ai_model,
    num_items=5,
    chunk_size=8000,
    overlap=500,
    api_key=None,
    openai_api_key=None,
    style_text=None,
    show_progress=True,
    gemini_model_id=None,
    audit_user_id=None,
    resolved_flavor=None,
    mix_basic_ratio=70,
):
    """텍스트를 청크로 나누어 모델 호출을 여러 번 수행
    
    Returns:
        - 객관식: 구조화된 dict 리스트 (각 dict는 {type, problem, options, answer, explanation})
        - 빈칸/단답/서술: 구조화된 dict 리스트 (각 dict는 {type, response_type, front, answer, explanation})
    """
    import sys
    chunks = split_text_into_chunks(text_content, chunk_size=chunk_size, overlap=overlap)
    total_chunks = len(chunks)
    
    print(f"[CHUNKS DEBUG] 총 청크 수: {total_chunks}", file=sys.stderr)
    
    if total_chunks == 0:
        return []
    
    base = num_items // total_chunks
    rem = num_items % total_chunks
    items_per_chunk = [base + (1 if i < rem else 0) for i in range(total_chunks)]

    results = [None] * total_chunks
    progress_bar = st.progress(0) if show_progress else None

    with concurrent.futures.ThreadPoolExecutor(max_workers=min(4, total_chunks)) as ex:
        futures = {}
        for idx, chunk in enumerate(chunks):
            n = items_per_chunk[idx]
            if n <= 0:
                results[idx] = ""
                continue
            futures[
                ex.submit(
                    generate_content,
                    chunk,
                    selected_mode,
                    ai_model,
                    n,
                    api_key,
                    openai_api_key,
                    style_text,
                    gemini_model_id,
                    audit_user_id,
                    resolved_flavor,
                    mix_basic_ratio,
                )
            ] = idx

        completed = 0
        for fut in concurrent.futures.as_completed(futures):
            idx = futures[fut]
            try:
                res = fut.result()
            except Exception as e:
                res = f"❌ 청크 처리 실패: {str(e)}"
            results[idx] = res if isinstance(res, str) else str(res)
            completed += 1
            if progress_bar is not None:
                progress_bar.progress(int(completed / total_chunks * 100))

    # 모든 청크 결과 결합
    combined = "\n".join([r for r in results if r])
    
    print(f"[COMBINED DEBUG] 청크 결과 개수: {len([r for r in results if r])}/{total_chunks}, 총 길이: {len(combined)}", file=sys.stderr)

    # 결합된 텍스트를 구조화된 형식으로 파싱
    structured_list = parse_generated_text_to_structured(combined, selected_mode)
    
    # 중복 제거
    seen = set()
    deduped = []
    for item in structured_list:
        key = str(item)  # 또는 더 정교한 키 생성
        if key not in seen:
            seen.add(key)
            deduped.append(item)
    
    # 필요한 개수만 반환
    return deduped[:num_items]

# ============================================================================
# 사이드바 설정
# ============================================================================
with st.sidebar:
    render_generation_recovery_panel()
    st.header("👤 계정")
    if st.session_state.auth_user_id:
        who = st.session_state.get("auth_email") or st.session_state.auth_user_id
        st.success(f"로그인됨: {who}")
        if is_admin_user():
            st.caption("운영자 권한: 활성")
        if st.button("로그아웃", key="auth_logout_btn"):
            reset_runtime_state_for_auth_change()
            st.session_state.auth_user_id = ""
            st.session_state.auth_access_token = ""
            st.session_state.auth_email = ""
            st.rerun()
    else:
        st.info("메인 화면에서 로그인 또는 회원가입을 진행하세요.")
    if not get_configured_admin_users():
        st.caption("운영자 계정 설정: AXIOMA_ADMIN_USERS=admin_id1,admin2")

    st.markdown("---")
    st.header("⚙️ 설정 & 모드")
    if st.session_state.auth_user_id:
        st.session_state.ai_model = st.radio(
            "🤖 AI 모델 선택",
            ["🔵 Google Gemini", "🟢 OpenAI ChatGPT"]
        )

        st.markdown("---")

        if st.session_state.ai_model == "🔵 Google Gemini":
            st.session_state.api_key = st.text_input("Gemini API Key 입력", type="password")
            st.session_state.gemini_model_id = st.text_input(
                "Gemini 모델 ID",
                value=st.session_state.gemini_model_id,
                help="예: gemini-2.0-flash, gemini-2.0-flash-lite"
            )
            st.session_state.openai_api_key = None
        else:
            st.session_state.api_key = None
            st.session_state.openai_api_key = st.text_input("OpenAI API Key 입력", type="password")

        st.markdown("---")
        st.session_state.chunk_size = st.slider("청크 크기 (문자 수)", 2000, 30000, 8000, 500)
        st.session_state.overlap = st.slider("청크 중첩 (문자 수)", 0, 5000, 500, 100)

        st.markdown("---")
        st.subheader("⚙️ 필터링 옵션")
        st.session_state.enable_filter = st.checkbox("품질 필터 사용", value=True)
        st.session_state.min_length = st.slider("최소 문자 수", 10, 200, 30)
        st.session_state.auto_tag_enabled = st.checkbox("자동 난이도/카테고리 태깅", value=True)
        st.session_state.explanation_default = st.checkbox("해설 기본 열기", value=st.session_state.explanation_default)
    else:
        st.caption("로그인 후 AI 키 및 생성 설정을 변경할 수 있습니다.")

    st.session_state.theme_enabled = False if LOCK_SAFE else True
    st.session_state.theme_mode = resolved_theme_mode
    st.session_state.theme_bg = "Gradient"

# 블록 외에서도 접근 가능하도록 로컬 변수에 할당
ai_model = st.session_state.get("ai_model", "🔵 Google Gemini")
api_key = st.session_state.get("api_key")
openai_api_key = st.session_state.get("openai_api_key")
chunk_size = st.session_state.get("chunk_size", 8000)
overlap = st.session_state.get("overlap", 500)
enable_filter = st.session_state.get("enable_filter", True)
min_length = st.session_state.get("min_length", 30)
auto_tag_enabled = st.session_state.get("auto_tag_enabled", True)

# Apply theme (skip if disabled)
THEME_ENABLED = should_apply_custom_theme(
    st.session_state.get("theme_enabled"),
    st.session_state.get("theme_mode"),
)
if THEME_ENABLED:
    apply_theme(st.session_state.theme_mode, st.session_state.theme_bg)
if MOBILE_CLIENT:
    apply_mobile_exam_styles()

if not st.session_state.get("auth_user_id"):
    render_auth_landing_page()
    st.stop()

def get_main_page_config(admin_mode):
    pages = [
        ("home", "🏠 홈"),
        ("generate", "📚 문제 생성"),
        ("convert", "🧾 기출문제 변환"),
        ("exam", "🎯 실전 시험"),
    ]
    if admin_mode:
        pages.append(("admin", "🛠️ 운영"))
    return pages

# ============================================================================
# 메인 UI: 라우팅 구조 (선택한 페이지만 렌더링)
# ============================================================================
admin_mode = is_admin_user()
main_pages = get_main_page_config(admin_mode)
main_labels = [label for _, label in main_pages]
label_to_page = {label: page for page, label in main_pages}
if "main_nav_label" not in st.session_state or st.session_state.main_nav_label not in main_labels:
    st.session_state.main_nav_label = main_labels[0]
active_label = st.radio("페이지", main_labels, horizontal=True, key="main_nav_label")
active_page = label_to_page.get(active_label, "home")

# ============================================================================
# PAGE: 홈
# ============================================================================
if active_page == "home":
    st.title("🏠 홈")
    show_action_notice()

    stats = get_question_stats()
    bank = load_questions()
    all_questions = bank.get("text", []) + bank.get("cloze", [])
    acc = compute_overall_accuracy(all_questions)
    acc_text = f"{acc['accuracy']:.1f}%" if acc else "—"

    if not st.session_state.get("theme_enabled"):
        st.info("Safe mode에서 테마가 비활성화되었습니다.")

    st.header("Axioma Qbank")
    st.write("강의록과 기출문제를 연결해 학습-시험-복습 흐름을 만듭니다.")
    st.write(f"전체 정답률: {acc_text}")
    st.write(f"저장된 객관식: {stats['total_text']} · 저장된 빈칸: {stats['total_cloze']}")

    with st.expander("🔐 초기 이용자용: API 키 발급 가이드", expanded=False):
        st.caption("문항 생성/변환/AI 보조 기능은 아래 모델 키가 필요합니다.")
        key_tabs = st.tabs(["Google Gemini", "OpenAI"])
        with key_tabs[0]:
            st.markdown(
                """
                1. [Google AI Studio](https://aistudio.google.com/app/apikey) 접속
                2. Google 계정 로그인 후 **Create API key** 클릭
                3. API 키 복사 후 앱 사이드바의 **Gemini API Key 입력**에 붙여넣기
                4. 모델은 `gemini-2.0-flash` 또는 `gemini-2.5-flash` 사용
                """
            )
            st.info("팁: 키는 환경변수/시크릿 관리 도구 대신 앱 세션에만 임시 저장됩니다. 브라우저 세션이 바뀌면 다시 입력해 주세요.")
        with key_tabs[1]:
            st.markdown(
                """
                1. [OpenAI API keys](https://platform.openai.com/api-keys) 접속
                2. 계정 로그인 후 **Create new secret key** 클릭
                3. key를 복사해 사이드바의 **OpenAI API Key 입력**에 붙여넣기
                4. 모델은 기본 `gpt-4o-mini`(권장) 또는 프로젝트에서 지정한 모델과 일치하도록 설정
                """
            )
            st.info("OpenAI 키는 사용량 과금이 발생할 수 있으니 프로젝트 단가/할당량을 먼저 확인하세요.")

    # 홈에서 바로 시험/학습 세션 준비
    st.markdown("---")
    st.subheader("빠른 시작 (분과/단원)")
    if all_questions:
        quick_subject_unit_map = collect_subject_unit_map(all_questions)
        quick_subjects_all = sorted(quick_subject_unit_map.keys())
        quick_subjects = st.multiselect(
            "학습할 분과",
            quick_subjects_all,
            default=quick_subjects_all[:1],
            key="home_quick_subjects",
        )

        quick_unit_filter = {}
        if quick_subjects:
            with st.expander("단원 선택", expanded=True):
                for subj in quick_subjects:
                    units = quick_subject_unit_map.get(subj, ["미분류"])
                    if not units:
                        units = ["미분류"]
                    key_name = f"home_unit_filter_{subj}"
                    prev_units = st.session_state.get(key_name, units)
                    selected_units = st.multiselect(
                        f"{subj} 단원",
                        options=units,
                        default=prev_units if set(prev_units) <= set(units) else units,
                        key=key_name,
                    )
                    if not selected_units:
                        selected_units = list(units)
                    quick_unit_filter[subj] = selected_units
        else:
            st.caption("분과를 먼저 선택하면 단원 체크박스가 나타납니다.")

        quick_mode = st.radio("모드", ["시험모드", "학습모드"], horizontal=True, key="home_quick_mode")
        quick_type = st.selectbox("문항 유형", ["객관식", "빈칸"], key="home_quick_type")

        filtered = filter_questions_by_subject_unit_hierarchy(all_questions, quick_subjects, quick_unit_filter)
        if filtered:
            quick_max = min(50, len(filtered))
            quick_min = 1 if quick_max < 5 else 5
            quick_num = st.slider("문항 수", quick_min, quick_max, min(10, quick_max), key="home_quick_num")
            if st.button("선택 조건으로 세션 준비", use_container_width=True, key="home_quick_prepare"):
                started = start_exam_session_from_items(filtered[:quick_num], quick_type, quick_mode)
                if started:
                    st.session_state.exam_mode_entry_anchor = "home"
                    st.session_state.last_action_notice = f"홈에서 {started}개 문항으로 {quick_mode}를 준비했습니다. 실전 시험 탭으로 이동해 시작하세요."
                    st.rerun()
                else:
                    st.warning("선택한 타입에 맞는 문항이 없습니다. 문항 유형(객관식/빈칸)을 다시 확인해 주세요.")
        else:
            st.info("선택한 분과/단원에 해당하는 문항이 없습니다.")
    else:
        st.info("문항이 없습니다. 먼저 문제를 생성/변환해 저장해 주세요.")

    # 통계
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("저장된 객관식", stats["total_text"])
    with col2:
        st.metric("저장된 빈칸", stats["total_cloze"])
    with col3:
        st.metric("전체 문항 정답률", acc_text)

    st.markdown("---")
    st.subheader("분과/단원 한눈에 보기")
    if all_questions:
        subject_overview = summarize_subject_review_status(all_questions)
        subject_unit_map = collect_subject_unit_map(all_questions)
        subject_rows = []
        for row in subject_overview:
            subj = row.get("분과", "General")
            units = subject_unit_map.get(subj, [])
            unit_text = ", ".join(units[:3]) + (" ..." if len(units) > 3 else "")
            subject_rows.append({
                "분과": subj,
                "총문항": row.get("총문항", 0),
                "복습대상": row.get("복습대상", 0),
                "오답문항": row.get("오답문항", 0),
                "연관 단원": unit_text,
            })
        safe_dataframe(subject_rows, use_container_width=True, hide_index=True)
    else:
        st.info("저장된 문항이 없습니다. 먼저 문제를 생성/변환해보세요.")

    st.markdown("---")
    st.subheader("학습 대시보드")
    wrong_items, total_wrong = get_wrong_note_stats(all_questions)
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("오답 누적 문항", len(wrong_items))
    with col2:
        st.metric("오답 누적 횟수", total_wrong)
    with col3:
        st.metric("전체 문항", len(all_questions))

    # 오답노트 필터
    subjects_all = sorted({(q.get("subject") or "General") for q in all_questions}) if all_questions else []
    diffs_all = sorted({(q.get("difficulty") or "미지정") for q in all_questions}) if all_questions else []
    sel_subjects = st.multiselect("오답노트 분과 필터", subjects_all, default=subjects_all)
    sel_diffs = st.multiselect("오답노트 난이도 필터", diffs_all, default=diffs_all)
    st.session_state.wrong_priority = st.selectbox(
        "오답노트 우선순위",
        ["오답 횟수", "오답률", "최근 오답"],
        index=["오답 횟수", "오답률", "최근 오답"].index(st.session_state.wrong_priority)
    )
    if st.session_state.wrong_priority == "최근 오답":
        st.session_state.wrong_weight_recent = st.slider(
            "가중치: 최근 오답",
            0.0, 1.0, st.session_state.wrong_weight_recent, 0.05
        )
        st.session_state.wrong_weight_count = 1.0 - st.session_state.wrong_weight_recent
        st.caption(f"오답 횟수 가중치: {st.session_state.wrong_weight_count:.2f}")
    filtered_wrong = [
        q for q in wrong_items
        if (q.get("subject") or "General") in sel_subjects
        and (q.get("difficulty") or "미지정") in sel_diffs
    ]

    if filtered_wrong:
        if st.button("📌 오답노트 세션 준비", use_container_width=True, key="prepare_wrong_session"):
            # 오답 문항으로 학습 세션 준비 (실전 시험 탭에서 진행)
            parsed_selected = []
            for raw in sort_wrong_first(
                filtered_wrong,
                mode=st.session_state.wrong_priority,
                weight_recent=st.session_state.wrong_weight_recent,
                weight_count=st.session_state.wrong_weight_count
            ):
                if raw.get("type") == "cloze":
                    parsed_selected.append(parse_cloze_content(raw))
                else:
                    parsed_selected.append(parse_mcq_content(raw))
            st.session_state.exam_questions = parsed_selected[:50]
            st.session_state.current_question_idx = 0
            st.session_state.user_answers = {}
            st.session_state.exam_started = True
            st.session_state.exam_finished = False
            st.session_state.exam_mode = "학습모드"
            st.session_state.revealed_answers = set()
            st.session_state.auto_advance_guard = None
            st.session_state.exam_stats_applied = False
            st.session_state.graded_questions = set()
            st.success("오답노트 세션이 준비되었습니다. 🎯 실전 시험 탭으로 이동해 시작하세요.")
    else:
        st.info("선택한 필터에 해당하는 오답 문항이 없습니다.")

    # FSRS / SRS 상태
    st.caption(f"복습 엔진: {'FSRS' if FSRS_AVAILABLE else '기본 SRS'}")

    if all_questions:
        with st.expander("📊 분과별 복습 큐(기본 화면)", expanded=False):
            subject_rows = summarize_subject_review_status(all_questions)
            if subject_rows:
                safe_dataframe(subject_rows, use_container_width=True, hide_index=True)
    elif not FSRS_AVAILABLE:
        st.info("FSRS 미설치: 기본 SRS로 동작 중입니다.")

    st.markdown("---")
    st.subheader("🧾 시험 기록")
    history = load_exam_history()
    if not history:
        st.info("저장된 시험 기록이 없습니다.")
    else:
        labels = []
        for idx, h in enumerate(history):
            ts = h.get("finished_at", "")
            acc = h.get("accuracy", 0)
            labels.append(f"{idx + 1}. {ts} | {h.get('type')} | {acc}%")
        sel = st.selectbox("기록 선택", labels, index=0)
        sel_idx = labels.index(sel)
        h = history[sel_idx]
        st.write(f"문항 수: {h.get('num_questions')} / 정답: {h.get('correct')} / 정확도: {h.get('accuracy')}%")
        if h.get("subjects"):
            st.caption(f"분과: {', '.join(h.get('subjects'))}")
        if h.get("units"):
            st.caption(f"단원: {', '.join(h.get('units'))}")

        for i, item in enumerate(h.get("items", []), 1):
            status_icon = "✅" if item.get("is_correct") else "❌"
            title = f"{status_icon} 문제 {i}"
            with st.expander(title, expanded=False):
                st.markdown(item.get("front") or "")
                if item.get("type") == "mcq":
                    opts = item.get("options") or []
                    letters = ["A", "B", "C", "D", "E"]
                    for idx_opt, opt in enumerate(opts[:5]):
                        st.write(f"{letters[idx_opt]}. {opt}")
                    user = item.get("user")
                    correct_num = item.get("correct")
                    user_display = letters[user - 1] if isinstance(user, int) and 1 <= user <= 5 else "응답 없음"
                    correct_display = letters[correct_num - 1] if isinstance(correct_num, int) and 1 <= correct_num <= 5 else "?"
                else:
                    user_display = item.get("user") or "응답 없음"
                    correct_display = item.get("answer") or ""

                st.divider()
                st.write(f"**당신의 답:** {user_display}")
                st.write(f"**정답:** {correct_display}")
                if item.get("explanation"):
                    show_exp = st.checkbox("해설 보기", value=st.session_state.explanation_default, key=f"hist_exp_{sel_idx}_{i}")
                    if show_exp:
                        st.markdown(format_explanation_text(item.get("explanation")))
                if item.get("id"):
                    note_key = f"hist_note_{sel_idx}_{i}"
                    st.text_area("메모", value=item.get("note", ""), key=note_key, height=80)
                    if st.button("메모 저장", key=f"save_hist_note_{sel_idx}_{i}"):
                        saved = update_question_note(item["id"], st.session_state.get(note_key, ""))
                        if saved:
                            st.success("메모 저장됨")

    with st.expander("🧹 데이터 관리", expanded=False):
        st.caption("주의: 삭제 작업은 되돌릴 수 없습니다.")
        confirm = st.checkbox("삭제 작업을 이해했습니다.")
        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("객관식 전체 삭제", use_container_width=True, disabled=not confirm):
                with st.spinner("객관식 문항 삭제 중..."):
                    clear_question_bank(mode="mcq")
                st.session_state.last_action_notice = "객관식 문항을 삭제했습니다."
                st.session_state.exam_started = False
                st.session_state.exam_questions = []
                st.session_state.user_answers = {}
                st.rerun()
        with col2:
            if st.button("빈칸 전체 삭제", use_container_width=True, disabled=not confirm):
                with st.spinner("빈칸 문항 삭제 중..."):
                    clear_question_bank(mode="cloze")
                st.session_state.last_action_notice = "빈칸 문항을 삭제했습니다."
                st.session_state.exam_started = False
                st.session_state.exam_questions = []
                st.session_state.user_answers = {}
                st.rerun()
        with col3:
            if st.button("전체 문항 삭제", use_container_width=True, disabled=not confirm):
                with st.spinner("전체 문항 삭제 중..."):
                    clear_question_bank(mode="all")
                st.session_state.last_action_notice = "모든 문항을 삭제했습니다."
                st.session_state.exam_started = False
                st.session_state.exam_questions = []
                st.session_state.user_answers = {}
                st.rerun()
        if st.button("시험 기록 삭제", use_container_width=True, disabled=not confirm):
            clear_exam_history()
            st.session_state.last_action_notice = "시험 기록을 삭제했습니다."
            st.rerun()

        st.markdown("---")
        subjects = sorted({(q.get("subject") or "General") for q in all_questions}) if all_questions else []
        sel_subjects_del = st.multiselect("분과별 삭제", subjects)
        if sel_subjects_del:
            if st.button("선택 분과 삭제", use_container_width=True, disabled=not confirm):
                data = load_questions()
                before_text = len(data.get("text", []))
                before_cloze = len(data.get("cloze", []))
                data["text"] = [q for q in data.get("text", []) if (q.get("subject") or "General") not in sel_subjects_del]
                data["cloze"] = [q for q in data.get("cloze", []) if (q.get("subject") or "General") not in sel_subjects_del]
                save_questions(data)
                deleted = (before_text - len(data.get("text", []))) + (before_cloze - len(data.get("cloze", [])))
                st.session_state.last_action_notice = f"{deleted}개 문항 삭제됨 (분과: {', '.join(sel_subjects_del)})"
                st.rerun()
        unit_subject = st.selectbox("단원별 삭제 대상 분과", ["선택 없음"] + subjects, key="unit_delete_subject")
        if unit_subject != "선택 없음":
            subject_units = sorted({
                (q.get("unit") or "미분류")
                for q in all_questions
                if (q.get("subject") or "General") == unit_subject
            })
            selected_units = st.multiselect(
                "삭제할 단원",
                subject_units,
                key=f"unit_delete_units_{unit_subject}",
            )
            delete_mode_label = st.radio(
                "단원 삭제 문항 유형",
                ["전체", "객관식", "빈칸"],
                horizontal=True,
                key=f"unit_delete_mode_{unit_subject}",
            )
            if selected_units:
                mode_map = {"전체": "all", "객관식": "mcq", "빈칸": "cloze"}
                if st.button("선택 단원 삭제", use_container_width=True, disabled=not confirm, key=f"delete_units_btn_{unit_subject}"):
                    deleted = delete_questions_by_subject_units(
                        {unit_subject: selected_units},
                        mode=mode_map.get(delete_mode_label, "all"),
                    )
                    st.session_state.last_action_notice = (
                        f"{deleted}개 문항 삭제됨 (분과: {unit_subject}, 단원: {', '.join(selected_units)})"
                    )
                    st.session_state.exam_started = False
                    st.session_state.exam_questions = []
                    st.session_state.user_answers = {}
                    st.rerun()

    with st.expander("🗑️ 객관식 선택 삭제", expanded=False):
        bank_now = load_questions()
        mcq_list = bank_now.get("text", [])
        if not mcq_list:
            st.info("객관식 문항이 없습니다.")
        else:
            st.caption("개별 문항을 선택해 삭제할 수 있습니다.")
            st.markdown("---")
            subj = st.selectbox(
                "분과 필터",
                ["전체"] + sorted({(q.get("subject") or "General") for q in mcq_list})
            )
            search = st.text_input("문항 검색", value="")
            filtered = []
            for q in mcq_list:
                if subj != "전체" and (q.get("subject") or "General") != subj:
                    continue
                text = q.get("problem", "")
                if search and search.lower() not in text.lower():
                    continue
                filtered.append(q)
            filtered = filtered[:200]

            def _fallback_mcq_multiselect():
                id_to_q = {q.get("id"): q for q in filtered if q.get("id")}
                options = list(id_to_q.keys())

                def format_item(qid):
                    q = id_to_q.get(qid) or {}
                    subj_name = q.get("subject") or "General"
                    title = (q.get("problem") or "")[:80]
                    return f"{qid[:8]} | {subj_name} | {title}"

                selected_ids = st.multiselect("개별 문항 선택", options, format_func=format_item)
                return selected_ids

            selected_ids = []
            if hasattr(st, "data_editor"):
                rows = []
                for q in filtered:
                    qid = q.get("id")
                    if not qid:
                        continue
                    rows.append({
                        "선택": False,
                        "id": qid,
                        "분과": q.get("subject") or "General",
                        "문항": (q.get("problem") or "")[:120],
                    })
                try:
                    edited = st.data_editor(
                        rows,
                        use_container_width=True,
                        hide_index=True,
                        column_config={
                            "id": st.column_config.TextColumn("ID", width="small"),
                            "분과": st.column_config.TextColumn("분과", width="small"),
                            "문항": st.column_config.TextColumn("문항", width="large"),
                        },
                        disabled=["id", "분과", "문항"],
                        key="mcq_delete_editor"
                    )
                    selected_ids = [r["id"] for r in edited if r.get("선택")]
                except Exception:
                    st.warning("데이터 에디터를 사용할 수 없어 목록 방식으로 대체합니다.")
                    selected_ids = _fallback_mcq_multiselect()
            else:
                selected_ids = _fallback_mcq_multiselect()

            confirm_sel = st.checkbox("개별 삭제 확인", key="confirm_item_delete")
            if selected_ids:
                if st.button("선택 문항 삭제", disabled=not confirm_sel):
                    deleted = delete_mcq_by_ids(selected_ids)
                    st.session_state.last_action_notice = f"{deleted}개 문항 삭제됨"
                    st.rerun()

            st.markdown("---")
            st.caption("세트(배치) 단위 삭제")
            batches = get_mcq_batches(mcq_list)
            if batches:
                batch_labels = []
                for b, cnt in sorted(batches.items(), key=lambda x: x[0]):
                    batch_labels.append(f"{b} ({cnt}개)")
                sel_batch = st.selectbox("세트 선택", ["선택 없음"] + batch_labels)
                confirm_batch = st.checkbox("세트 삭제 확인", key="confirm_batch_delete")
                if sel_batch != "선택 없음":
                    batch_id = sel_batch.split(" (")[0]
                    if st.button("세트 삭제", disabled=not confirm_batch):
                        deleted = delete_mcq_by_batch(batch_id)
                        st.session_state.last_action_notice = f"{deleted}개 문항 삭제됨 (세트: {batch_id})"
                        st.rerun()
            else:
                st.caption("세트 정보가 없습니다.")

    with st.expander("🛠️ 문항 개별 수정", expanded=False):
        bank_edit = load_questions()
        edit_type = st.radio(
            "문항 유형",
            ["객관식", "빈칸"],
            horizontal=True,
            key="edit_question_type",
        )
        source = bank_edit["text"] if edit_type == "객관식" else bank_edit["cloze"]
        if not source:
            st.info("수정 가능한 문항이 없습니다.")
        else:
            subjects = sorted({(q.get("subject") or "General") for q in source})
            subject_filter = st.selectbox("분과 필터", ["전체"] + subjects, key="edit_subject_filter")
            unit_filter = st.selectbox(
                "단원 필터",
                ["전체"] + sorted({(q.get("unit") or "미분류") for q in source if (q.get("subject") or "General") == subject_filter or subject_filter == "전체"}),
                key="edit_unit_filter"
            )
            keyword = st.text_input("문항 검색", value="", key="edit_keyword")

            candidates = []
            for q in source:
                if subject_filter != "전체" and (q.get("subject") or "General") != subject_filter:
                    continue
                if unit_filter != "전체" and (q.get("unit") or "미분류") != unit_filter:
                    continue
                text = q.get("problem") if edit_type == "객관식" else q.get("front", "")
                if keyword and keyword.lower() not in (text or "").lower():
                    continue
                candidates.append(q)

            if not candidates:
                st.info("필터 조건에 맞는 문항이 없습니다.")
            else:
                id_to_q = {q.get("id"): q for q in candidates if q.get("id")}

                def _format_question(qid):
                    q = id_to_q.get(qid) or {}
                    stem = (q.get("problem") if edit_type == "객관식" else q.get("front", "")) or ""
                    return f"{qid[:8]} | {(q.get('subject') or 'General')} | {(q.get('unit') or '미분류')} | {stem[:60]}"

                selected_id = st.selectbox(
                    "수정할 문항",
                    options=list(id_to_q.keys()),
                    format_func=_format_question,
                    key="selected_question_to_edit"
                )
                selected = id_to_q.get(selected_id)
                if selected:
                    st.markdown(f"**문항 ID:** `{selected_id}`")
                    edited_subject = st.text_input("과목", value=selected.get("subject") or "General", key=f"edit_subject_{selected_id}")
                    edited_unit = st.text_input("단원", value=selected.get("unit") or "미분류", key=f"edit_unit_{selected_id}")
                    edited_difficulty = st.text_input("난이도", value=selected.get("difficulty") or "", key=f"edit_difficulty_{selected_id}")
                    if edit_type == "객관식":
                        edited_problem = st.text_area("문항", value=selected.get("problem", ""), height=180, key=f"edit_problem_{selected_id}")
                        edited_options_raw = st.text_area(
                            "선지 (줄바꿈 구분)",
                            value="\n".join(selected.get("options") or []),
                            height=160,
                            key=f"edit_options_{selected_id}"
                        )
                        edited_answer = st.number_input(
                            "정답 번호(1~5)",
                            min_value=1,
                            max_value=max(1, len([l for l in (selected.get('options') or [])])),
                            value=int(selected.get("answer") or 1),
                            step=1,
                            key=f"edit_answer_{selected_id}"
                        )
                    else:
                        edited_problem = st.text_area("문항", value=selected.get("front", ""), height=180, key=f"edit_front_{selected_id}")
                        edited_answer = st.text_area("정답", value=selected.get("answer", ""), height=80, key=f"edit_answer_cloze_{selected_id}")
                    edited_explanation = st.text_area(
                        "해설",
                        value=selected.get("explanation", ""),
                        height=120,
                        key=f"edit_explanation_{selected_id}"
                    )
                    edited_note = st.text_area(
                        "메모",
                        value=selected.get("note", ""),
                        height=80,
                        key=f"edit_note_{selected_id}"
                    )

                    if st.button("문항 수정 저장", use_container_width=True, key="save_question_edit"):
                        patch = {
                            "subject": edited_subject,
                            "unit": edited_unit,
                            "difficulty": edited_difficulty,
                            "explanation": edited_explanation,
                            "note": edited_note
                        }
                        if edit_type == "객관식":
                            options_lines = [s.strip() for s in edited_options_raw.splitlines() if s.strip()]
                            patch["problem"] = edited_problem
                            patch["options"] = options_lines
                            patch["answer"] = int(edited_answer)
                        else:
                            patch["front"] = edited_problem
                            patch["answer"] = edited_answer
                        if update_question_by_id(selected_id, patch):
                            st.success("문항이 저장되었습니다.")
                            st.rerun()
                        else:
                            st.error("문항 저장에 실패했습니다.")

    st.markdown("---")
    st.subheader("학습 시각화")
    colv1, colv2 = st.columns([1, 1])
    with colv1:
        if st.button("학습 시각화 불러오기", key="load_home_visuals", use_container_width=True):
            st.session_state.home_visual_loaded = True
            st.rerun()
    with colv2:
        if st.session_state.home_visual_loaded:
            if st.button("학습 시각화 숨기기", key="hide_home_visuals", use_container_width=True):
                st.session_state.home_visual_loaded = False
                st.rerun()

    if not st.session_state.home_visual_loaded:
        st.caption("성능 최적화를 위해 시각화는 기본 숨김 상태입니다. 필요할 때만 불러오세요.")
    else:
        colp1, colp2, colp3 = st.columns([1, 1, 1])
        with colp1:
            st.session_state.profile_name = st.text_input(
                "설정 프리셋 이름",
                value=st.session_state.profile_name,
                help="히트맵 구간/색상 등 개인 설정을 저장해두는 기능입니다.",
            )
        with colp2:
            if st.button("프리셋 불러오기"):
                profile_name = (st.session_state.profile_name or "").strip()
                loaded = apply_profile_settings(profile_name)
                st.session_state.last_action_notice = "프로필 설정을 불러왔습니다." if loaded else "해당 프로필이 없습니다."
        with colp3:
            if st.button("프리셋 저장"):
                profile_name = (st.session_state.profile_name or "").strip()
                if not profile_name:
                    profile_name = "default"
                    st.session_state.profile_name = profile_name
                persist_profile_settings(profile_name)
                st.session_state.last_action_notice = "프로필 설정을 저장했습니다."

        st.caption("프리셋은 히트맵 구간/색상 등 개인 설정을 저장해두는 기능입니다.")
        acc = compute_overall_accuracy(all_questions)
        heat = compute_activity_heatmap(all_questions, days=365)
        with st.expander("히트맵 구간/색상 설정", expanded=False):
            st.caption("문항 수 구간을 조정하면 색 농도가 바뀝니다.")
            b1 = st.number_input("구간 1 (1회)", min_value=1, value=1)
            b2 = st.number_input("구간 2 (2~)", min_value=2, value=3)
            b3 = st.number_input("구간 3 (4~)", min_value=3, value=6)
            b4 = st.number_input("구간 4 (7~)", min_value=4, value=10)
            st.session_state.heatmap_bins = [0, b1, b2, b3, b4]
            st.session_state.heatmap_colors = [
                "#ffffff",
                st.color_picker("색상 1", value=st.session_state.heatmap_colors[1]),
                st.color_picker("색상 2", value=st.session_state.heatmap_colors[2]),
                st.color_picker("색상 3", value=st.session_state.heatmap_colors[3]),
                st.color_picker("색상 4", value=st.session_state.heatmap_colors[4]),
                st.color_picker("색상 5", value=st.session_state.heatmap_colors[5]),
            ]
        col_left, col_right = st.columns([1, 2])
        with col_left:
            st.markdown("**전체 정답률**")
            if acc:
                try:
                    import pandas as pd
                    import altair as alt

                    df = pd.DataFrame([
                        {"label": "Correct", "value": acc["correct"]},
                        {"label": "Wrong", "value": acc["wrong"]},
                    ])
                    base = alt.Chart(df).mark_arc(innerRadius=60, outerRadius=100).encode(
                        theta=alt.Theta("value:Q"),
                        color=alt.Color("label:N", scale=alt.Scale(range=["#34d399", "#f87171"]), legend=None),
                        tooltip=["label:N", "value:Q"]
                    )
                    text = alt.Chart(pd.DataFrame([{"text": f"{acc['accuracy']:.1f}%"}])).mark_text(
                        size=26, font="IBM Plex Sans", fontWeight="600"
                    ).encode(text="text:N")
                    st.altair_chart((base + text).properties(width=220, height=220), use_container_width=False)
                    st.caption(f"{acc['correct']}/{acc['total']} 정답")
                except Exception:
                    st.metric("전체 정답률", f"{acc['accuracy']:.1f}%")
            else:
                st.info("아직 풀이 기록이 없습니다.")

        with col_right:
            st.markdown("**학습 활동 히트맵 (최근 365일)**")
            if heat:
                try:
                    import pandas as pd
                    import altair as alt

                    df = pd.DataFrame(heat)
                    df["dow_label"] = df["dow"].map({0: "Mon", 1: "Tue", 2: "Wed", 3: "Thu", 4: "Fri", 5: "Sat", 6: "Sun"})
                    df["week_index"] = df["week_index"].astype(str)
                    b = st.session_state.heatmap_bins
                    labels = ["0", f"1-{b[1]}", f"{b[1] + 1}-{b[2]}", f"{b[2] + 1}-{b[3]}", f"{b[3] + 1}-{b[4]}", f"{b[4] + 1}+"]
                    df["bucket"] = pd.cut(
                        df["count"],
                        bins=[-0.1, 0, b[1], b[2], b[3], b[4], 9999],
                        labels=labels
                    )
                    heatmap = (
                        alt.Chart(df)
                        .mark_rect(cornerRadius=0)
                        .encode(
                            x=alt.X("week_index:O", axis=None),
                            y=alt.Y("dow_label:O", sort=["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"], axis=None),
                            color=alt.Color(
                                "bucket:N",
                                scale=alt.Scale(
                                    domain=labels,
                                    range=st.session_state.heatmap_colors
                                ),
                                legend=None
                            ),
                            tooltip=["date:T", "count:Q", "accuracy:Q"]
                        )
                        .properties(width=alt.Step(12), height=alt.Step(12))
                    )
                    st.altair_chart(heatmap, use_container_width=True)
                except Exception:
                    safe_dataframe(heat, use_container_width=True, hide_index=True)

if active_page == "admin" and admin_mode:
        st.title("🛠️ 운영자 콘솔")
        st.caption("사용자별 API 사용량, 호출 건수, 추정 비용을 확인합니다.")

        all_users = list_local_user_ids()
        if not all_users:
            st.info("로컬 사용자 데이터가 없습니다.")
        else:
            selected = st.selectbox("대상 사용자", ["전체"] + all_users, index=0, key="admin_user_filter")
            days = st.slider("조회 기간(일)", 1, 365, 30, 1, key="admin_days_filter")
            cutoff = datetime.now(timezone.utc) - timedelta(days=days)

            rows = []
            target_users = all_users if selected == "전체" else [selected]
            for uid in target_users:
                for row in read_audit_rows_for_user(uid):
                    ts_raw = str(row.get("timestamp") or "")
                    try:
                        ts = datetime.fromisoformat(ts_raw.replace("Z", "+00:00"))
                        if ts.tzinfo is None:
                            ts = ts.replace(tzinfo=timezone.utc)
                    except Exception:
                        continue
                    if ts < cutoff:
                        continue
                    enriched = dict(row)
                    enriched["user_id"] = uid
                    rows.append(enriched)

            if not rows:
                st.warning("선택한 조건에서 조회된 로그가 없습니다.")
            else:
                summary = summarize_usage_rows(rows)
                total_est, breakdown = estimate_cost_usd_from_summary(summary)
                total_calls = sum(x.get("calls", 0) for x in summary.values())
                total_tokens = sum(x.get("tokens", 0) for x in summary.values())

                m1, m2, m3 = st.columns(3)
                m1.metric("총 API 호출", f"{total_calls:,}")
                m2.metric("총 토큰", f"{total_tokens:,}")
                m3.metric("추정 비용(USD)", f"${total_est:.4f}")

                st.markdown("### 모델별 사용량")
                safe_dataframe(breakdown, use_container_width=True, hide_index=True)

                st.markdown("### 최근 로그")
                latest = sorted(rows, key=lambda x: str(x.get("timestamp") or ""), reverse=True)[:50]
                latest_view = [
                    {
                        "timestamp": r.get("timestamp"),
                        "user_id": r.get("user_id"),
                        "event": r.get("event"),
                        "model": r.get("model"),
                        "usage_tokens": r.get("usage_tokens"),
                    }
                    for r in latest
                ]
                safe_dataframe(latest_view, use_container_width=True, hide_index=True)

# ============================================================================
# PAGE: 문제 생성
# ============================================================================
if active_page == "generate":
    st.title("📚 문제 생성 & 저장")

    st.subheader("⚡ 빠른 시작")
    st.markdown("### 3단계로 시작하기")
    st.markdown("1) 자료 업로드 → 2) 모드/문항 수 설정 → 3) 문제 생성 시작")
    st.markdown("API 키가 없다면 사이드바 입력 후 다시 진행하세요.")
    with st.expander("🔒 데이터 처리/보안 안내", expanded=False):
        st.markdown(
            """
            - 업로드 자료는 **현재 로그인한 본인 계정의 생성 작업에만 사용**됩니다.
            - 다른 사용자 계정의 문제 생성에 업로드 원문이 재사용되지 않습니다.
            - 원문 파일 자체는 영구 저장하지 않고, 처리 후 텍스트/문항 결과만 계정 데이터로 저장됩니다.
            - 권리를 보유했거나 사용 허락을 받은 자료만 업로드해 주세요.
            """
        )
    with st.expander("❓ 자주 묻는 질문 (베타)", expanded=False):
        st.markdown(
            """
            **Q1. 손글씨도 읽나요?**  
            A. 인식은 가능하지만, 필기체/저화질 스캔에서는 정확도가 떨어질 수 있습니다.

            **Q2. 한 번에 몇 문제를 요청하는 게 좋나요?**  
            A. 권장값은 파일당 10~15문항입니다. 20문항 이상부터 처리 시간이 길어질 수 있습니다.

            **Q3. 속도는 무엇에 영향받나요?**  
            A. 파일 페이지 수, 스캔 품질(OCR 여부), 요청 문항 수, 스타일 파일 사용 여부의 영향을 받습니다.
            """
        )
    st.caption("처리 단계: 텍스트 추출 → OCR/AI 폴백 → 문항 생성 → 저장")

    ai_model_key_ready = bool(api_key) if ai_model == "🔵 Google Gemini" else bool(openai_api_key)
    if not ai_model_key_ready:
        st.warning("현재 AI 모델 키가 비어 있습니다. 사이드바에서 API 키를 입력하면 바로 시작할 수 있습니다.")

    st.markdown("---")

    # 파일 업로드
    uploaded_files = st.file_uploader(
        "강의 자료 업로드",
        type=["pdf", "docx", "pptx", "hwp"],
        accept_multiple_files=True,
        key="gen_upload_files",
    )
    uploaded_file = uploaded_files[0] if uploaded_files else None
    style_file = st.file_uploader(
        "기출문제 스타일 업로드 (선택)",
        type=["pdf", "docx", "pptx", "hwp", "txt", "tsv", "json"],
        key="style_upload",
    )
    if uploaded_files:
        total_size = sum(getattr(f, "size", 0) for f in uploaded_files)
        st.caption(
            f"업로드 용량: {len(uploaded_files)}개 파일 / {total_size / (1024 * 1024):.1f} MB "
            f"(권장: 한 번에 1~5개 파일)"
        )
    gen_copyright_ok = render_copyright_ack("gen")
    if (uploaded_files or style_file) and not gen_copyright_ok:
        st.warning("파일 분석/문제 생성을 시작하려면 저작권 확인 체크를 완료하세요.")

    mode = MODE_MCQ
    num_items = 10
    subject_input = "General"
    unit_input = "미분류"
    flavor_choice = "선택하세요"
    mix_basic_ratio = 70

    raw_text_cached = None
    style_text = None
    uploaded_bytes = uploaded_file.getvalue() if uploaded_file else b""
    uploaded_signature = build_upload_signature(uploaded_file.name, uploaded_bytes) if uploaded_file else ""
    style_bytes = style_file.getvalue() if style_file else b""
    style_signature = build_upload_signature(style_file.name, style_bytes) if style_file else ""

    if uploaded_file and gen_copyright_ok:
        raw_text_cached = get_generation_prewarm_text("raw", uploaded_signature)
        raw_error = get_generation_prewarm_error("raw", uploaded_signature)
        if raw_text_cached is None and not raw_error:
            try:
                with st.spinner("사전 준비 중: 강의자료 텍스트 추출"):
                    raw_text_cached = extract_text_from_file(
                        make_uploaded_file_from_bytes(uploaded_file.name, uploaded_bytes),
                        ai_model=ai_model,
                        ai_fallback=True,
                        api_key=api_key,
                        openai_api_key=openai_api_key,
                    )
                set_generation_prewarm_text("raw", uploaded_signature, raw_text_cached)
            except Exception as e:
                set_generation_prewarm_error("raw", uploaded_signature, str(e))
                raw_error = str(e)
        if raw_text_cached:
            est_chunks = len(split_text_into_chunks(raw_text_cached, chunk_size=chunk_size, overlap=overlap))
            st.caption(f"사전 준비 완료(첫 파일): 본문 {len(raw_text_cached):,}자 | 예상 청크 {est_chunks}개")
        elif raw_error:
            st.warning(f"사전 준비 실패(첫 파일): {raw_error}")

    if style_file and gen_copyright_ok:
        style_text = get_generation_prewarm_text("style", style_signature)
        style_error = get_generation_prewarm_error("style", style_signature)
        if style_text is None and not style_error:
            try:
                ext = Path(style_file.name).suffix.lower()
                if ext in [".txt", ".tsv", ".json"]:
                    style_text = style_bytes.decode("utf-8", errors="ignore")
                else:
                    style_text = extract_text_from_file(
                        make_uploaded_file_from_bytes(style_file.name, style_bytes)
                    )
                set_generation_prewarm_text("style", style_signature, style_text)
            except Exception as e:
                set_generation_prewarm_error("style", style_signature, str(e))
                style_error = str(e)
        if style_error:
            st.warning(f"기출문제 스타일 파일 처리 실패: {style_error}")
    elif style_file and not gen_copyright_ok:
        st.caption("권리 확인 체크 전에는 스타일 파일을 분석하지 않습니다.")

    if style_text:
        detected_mode, pattern = detect_term_language_mode(style_text)
        label = "혼용"
        if detected_mode == "ko":
            label = "한국어 용어 중심"
        elif detected_mode == "en":
            label = "영어 용어 중심"
        elif pattern:
            label = f"혼용 ({pattern})"
        st.caption(f"스타일 자동 감지: 용어 표기 = {label}")

    runtime_context = get_generation_runtime_context() if ai_model_key_ready else {}

    if uploaded_files:
        if len(uploaded_files) == 1:
            st.info(f"📄 **{uploaded_files[0].name}** ({uploaded_files[0].size:,} bytes)")
        else:
            st.info(f"📄 선택 파일: {len(uploaded_files)}개 (첫 파일: {uploaded_files[0].name})")

        st.markdown("### 설정")
        col1, col2 = st.columns(2)
        with col1:
            mode = st.radio("모드", [MODE_MCQ, MODE_CLOZE, MODE_SHORT, MODE_ESSAY])
        with col2:
            num_items = st.slider("생성 개수", 1, 50, 10)
            if num_items >= 20:
                st.caption("안내: 20개 이상 요청 시 처리 시간이 길어질 수 있습니다.")

        flavor_choice = st.selectbox(
            "문항 성격",
            ["선택하세요", "자동 판별(Auto)", "기초의학형(Basic)", "케이스형(Case)", "혼합(Mix)"],
            index=0,
            key="generation_flavor_choice",
        )
        if flavor_choice == "혼합(Mix)":
            st.caption("혼합 비율: Basic 70% / Case 30%")

        col_subj, col_unit = st.columns(2)
        with col_subj:
            subject_input = st.text_input("과목명 (예: 순환기내과)", value="General")
        with col_unit:
            unit_input = st.text_input("단원명 (선택)", value="미분류")
        estimated_minutes = estimate_generation_runtime_minutes(
            total_bytes=sum(getattr(f, "size", 0) for f in uploaded_files),
            num_files=len(uploaded_files),
            num_items=num_items,
            has_style_file=bool(style_file),
        )
        low = max(1, int(round(estimated_minutes * 0.7)))
        high = max(low + 1, int(round(estimated_minutes * 1.6)))
        st.caption(f"예상 처리 시간(대기열 추가+생성): 약 {low}~{high}분")
        if flavor_choice == "자동 판별(Auto)" and uploaded_file:
            preview_flavor = resolve_generation_flavor(
                flavor_choice,
                raw_text=raw_text_cached or "",
                style_text=style_text or "",
                subject=subject_input,
            )
            st.caption(f"자동 판별 예상: `{preview_flavor}`")

        if gen_copyright_ok:
            col_p1, col_p2 = st.columns([1, 1])
            with col_p1:
                if uploaded_file and st.button("사전 준비 다시 실행(첫 파일)", use_container_width=True, key="regen_prewarm_main"):
                    clear_generation_prewarm_error("raw", uploaded_signature)
                    cache_map = st.session_state.get("generation_prewarm_cache", {})
                    cache_key = _prewarm_cache_key("raw", uploaded_signature)
                    if cache_key in cache_map:
                        del cache_map[cache_key]
                    st.session_state["generation_prewarm_cache"] = cache_map
                    st.rerun()
            with col_p2:
                if style_file and st.button("스타일 사전 준비 다시 실행", use_container_width=True, key="regen_prewarm_style"):
                    clear_generation_prewarm_error("style", style_signature)
                    cache_map = st.session_state.get("generation_prewarm_cache", {})
                    cache_key = _prewarm_cache_key("style", style_signature)
                    if cache_key in cache_map:
                        del cache_map[cache_key]
                    st.session_state["generation_prewarm_cache"] = cache_map
                    st.rerun()

        if not ai_model_key_ready:
            st.button("🚀 업로드 파일들을 대기열에 추가", use_container_width=True, disabled=True, help="API 키를 먼저 입력해 주세요.")
        elif not gen_copyright_ok:
            st.button("🚀 업로드 파일들을 대기열에 추가", use_container_width=True, disabled=True, help="저작권 확인 체크를 완료해 주세요.")
        elif flavor_choice == "선택하세요":
            st.button("🚀 업로드 파일들을 대기열에 추가", use_container_width=True, disabled=True, help="문항 성격을 선택해 주세요.")
        elif st.button("🚀 업로드 파일들을 대기열에 추가", use_container_width=True):
            try:
                queue_items = load_generation_queue_items()
                added = 0
                skipped = 0
                skipped_short = 0
                skipped_duplicate = 0

                style_text_for_queue = style_text
                if style_file and not style_text_for_queue:
                    ext = Path(style_file.name).suffix.lower()
                    if ext in [".txt", ".tsv", ".json"]:
                        style_text_for_queue = style_bytes.decode("utf-8", errors="ignore")
                    else:
                        style_text_for_queue = extract_text_from_file(
                            make_uploaded_file_from_bytes(style_file.name, style_bytes)
                        )
                    set_generation_prewarm_text("style", style_signature, style_text_for_queue)

                with st.spinner("업로드 파일을 대기열용으로 준비 중..."):
                    file_payloads = []
                    seen_payload_signatures = set()
                    for uf in uploaded_files:
                        file_bytes = uf.getvalue()
                        file_sig = build_upload_signature(uf.name, file_bytes)
                        if file_sig in seen_payload_signatures:
                            skipped_duplicate += 1
                            continue
                        seen_payload_signatures.add(file_sig)
                        file_payloads.append((uf.name, file_sig, file_bytes))

                    extracted_texts = {}
                    pending = []
                    for file_name, file_sig, file_bytes in file_payloads:
                        raw_text = get_generation_prewarm_text("raw", file_sig)
                        if raw_text:
                            extracted_texts[file_sig] = raw_text
                        else:
                            pending.append((file_name, file_sig, file_bytes))

                    if pending:
                        with concurrent.futures.ThreadPoolExecutor(max_workers=min(4, len(pending))) as ex:
                            futures = {}
                            for file_name, file_sig, file_bytes in pending:
                                fut = ex.submit(
                                    extract_text_from_file,
                                    make_uploaded_file_from_bytes(file_name, file_bytes),
                                    ai_model=ai_model,
                                    ai_fallback=True,
                                    api_key=api_key,
                                    openai_api_key=openai_api_key,
                                )
                                futures[fut] = (file_name, file_sig)
                            for fut in concurrent.futures.as_completed(futures):
                                file_name, file_sig = futures[fut]
                                try:
                                    raw_text = fut.result()
                                    extracted_texts[file_sig] = raw_text
                                    set_generation_prewarm_text("raw", file_sig, raw_text)
                                except Exception as e:
                                    set_generation_prewarm_error("raw", file_sig, str(e))
                                    st.warning(f"{file_name}: 텍스트 추출 실패 ({str(e)})")

                    for file_name, file_sig, _ in file_payloads:
                        raw_text = extracted_texts.get(file_sig) or get_generation_prewarm_text("raw", file_sig)
                        raw_text = (raw_text or "").strip()
                        if not raw_text:
                            skipped += 1
                            continue
                        if len(raw_text) < 20:
                            skipped += 1
                            skipped_short += 1
                            continue
                        resolved_flavor = resolve_generation_flavor(
                            flavor_choice,
                            raw_text=raw_text,
                            style_text=style_text_for_queue,
                            subject=subject_input,
                        )
                        if is_duplicate_generation_queue_item(
                            queue_items,
                            source_signature=file_sig,
                            flavor_choice=flavor_choice,
                            mode=mode,
                            num_items=num_items,
                            subject=subject_input,
                            unit=unit_input,
                        ):
                            skipped_duplicate += 1
                            continue
                        queue_items.append(
                            build_generation_queue_item(
                                source_name=file_name,
                                source_signature=file_sig,
                                raw_text=raw_text,
                                style_text=style_text_for_queue,
                                flavor_choice=flavor_choice,
                                resolved_flavor=resolved_flavor,
                                mix_basic_ratio=mix_basic_ratio,
                                mode=mode,
                                num_items=num_items,
                                subject=subject_input,
                                unit=unit_input,
                                ai_model=ai_model,
                                chunk_size=chunk_size,
                                overlap=overlap,
                                quality_filter=enable_filter,
                                min_length=min_length,
                            )
                        )
                        added += 1

                if added <= 0:
                    if skipped_short:
                        st.warning(f"대기열에 추가할 텍스트 문서는 추출되었지만, 유효 분량이 부족합니다. PDF(강의록)은 AI 폴백을 켜고 재시도하세요. 건너뜀: {skipped}개")
                    else:
                        st.warning(
                            "대기열에 추가할 텍스트 문서가 없습니다. "
                            f"(분석 파일 {len(uploaded_files)}개, 건너뜀 {skipped}개, 중복 {skipped_duplicate}개)"
                        )
                elif not save_generation_queue_items(queue_items):
                    st.error("대기열 저장 실패: 사용자 설정 저장 중 오류가 발생했습니다.")
                else:
                    queue_items, _ = start_next_generation_queue_job_if_idle(
                        queue_items,
                        api_key=api_key,
                        openai_api_key=openai_api_key,
                        runtime_context=runtime_context,
                    )
                    save_generation_queue_items(queue_items)
                    st.session_state.generation_failure = ""
                    msg = f"대기열 추가 완료: {added}개"
                    if skipped:
                        msg += f" (건너뜀 {skipped}개)"
                    if skipped_duplicate:
                        msg += f" (중복 제외 {skipped_duplicate}개)"
                    st.session_state.last_action_notice = msg
                    st.rerun()
            except Exception as e:
                import traceback
                err_msg = f"❌ 오류: {str(e)}"
                st.error(err_msg)
                st.error(f"상세 오류:\n{traceback.format_exc()}")
                st.session_state.generation_failure = err_msg

    queue_items = load_generation_queue_items()
    queue_items, revived = revive_stale_running_queue_items(queue_items)
    queue_items, queue_notices = reconcile_generation_queue_with_async(
        queue_items,
        default_quality_filter=enable_filter,
        default_min_length=min_length,
    )
    for notice in queue_notices:
        st.success(notice)
    queue_items, auto_started = start_next_generation_queue_job_if_idle(
        queue_items,
        api_key=api_key,
        openai_api_key=openai_api_key,
        runtime_context=runtime_context,
    )
    if revived or queue_notices or auto_started:
        save_generation_queue_items(queue_items)
    if auto_started:
        st.rerun()

    st.markdown("### 🧾 생성 대기열")
    if not queue_items:
        st.info("현재 대기열이 비어 있습니다.")
    else:
        status_label = {
            "queued": "대기",
            "running": "생성중",
            "done": "완료",
            "failed": "실패",
            "cancelled": "취소",
        }
        for idx, item in enumerate(queue_items, 1):
            jid = str(item.get("id"))
            left, right = st.columns([4, 1])
            with left:
                st.markdown(
                    f"**{idx}. {item.get('source_name', 'unknown')}**  \n"
                    f"- 상태: `{status_label.get(item.get('status'), item.get('status'))}` | "
                    f"문항성격: `{item.get('resolved_flavor') or item.get('flavor_choice') or '-'}` | "
                    f"모드: `{item.get('mode', '')}` | 문항수: `{item.get('num_items', 0)}` | "
                    f"과목/단원: `{item.get('subject', 'General')} / {item.get('unit', '미분류')}`"
                )
                if item.get("status") == "done":
                    st.caption(f"자동 저장: {int(item.get('saved_count', 0))}개")
                if item.get("status") in {"failed", "cancelled"} and item.get("error"):
                    st.caption(f"사유: {item.get('error')}")
            with right:
                if item.get("status") == "queued":
                    if st.button("대기 취소", key=f"queue_cancel_{jid}", use_container_width=True):
                        changed, queue_items = remove_generation_queue_job(queue_items, jid)
                        if changed and save_generation_queue_items(queue_items):
                            st.rerun()
                elif item.get("status") == "running":
                    if st.button("실행 취소", key=f"queue_stop_{jid}", use_container_width=True):
                        async_job = st.session_state.get("generation_async_job")
                        if isinstance(async_job, dict) and str(async_job.get("queue_id")) == jid:
                            future = async_job.get("future")
                            cancelled = False
                            if future is not None:
                                try:
                                    cancelled = future.cancel()
                                except Exception:
                                    cancelled = False
                            async_job["status"] = "cancelled" if cancelled else "error"
                            async_job["error"] = "사용자 취소" if cancelled else "이미 실행 중인 작업은 취소할 수 없습니다."
                            st.session_state["generation_async_job"] = async_job
                            st.rerun()
                else:
                    if st.button("목록 삭제", key=f"queue_remove_{jid}", use_container_width=True):
                        changed, queue_items = remove_generation_queue_job(queue_items, jid)
                        if changed and save_generation_queue_items(queue_items):
                            st.rerun()
        if st.button("완료/취소/실패 항목 정리", key="queue_prune_finished", use_container_width=True):
            before = len(queue_items)
            queue_items = [x for x in queue_items if x.get("status") in {"queued", "running"}]
            if len(queue_items) != before and save_generation_queue_items(queue_items):
                st.rerun()

    st.markdown("---")
    st.info("기출문제 파일 변환은 **🧾 기출문제 변환** 탭에서 진행합니다.")

# ============================================================================
# PAGE: 기출문제 변환
# ============================================================================
if active_page == "convert":
    st.title("🧾 기출문제 전용 변환")
    st.caption("HWP/PDF/DOCX/PPTX/TXT/TSV 파일을 기출문제 형식으로 변환하여 저장합니다.")
    convert_copyright_ok = render_copyright_ack("convert")
    if not convert_copyright_ok:
        st.warning("저작권 확인 체크를 완료해야 파일 변환을 실행할 수 있습니다.")

    with st.expander("🧩 HWP+PDF 듀얼 업로드(수동 최소화)", expanded=False):
        st.caption("HWP에서 문항 텍스트를 추출하고, PDF에서 이미지/페이지 정보를 연결합니다.")
        col_dual1, col_dual2 = st.columns(2)
        with col_dual1:
            dual_hwp = st.file_uploader("HWP 업로드 (문항 텍스트)", type=["hwp"], key="dual_hwp_upload")
        with col_dual2:
            dual_pdf = st.file_uploader("PDF 업로드 (이미지/레이아웃)", type=["pdf"], key="dual_pdf_upload")

        dual_subject = st.text_input("기본 과목명", value="General", key="dual_subject")
        dual_unit = st.text_input("기본 단원명 (선택)", value="DualUpload", key="dual_unit")

        dual_threshold = st.slider("자동 매칭 신뢰도 기준", 0.05, 0.6, 0.2, step=0.05, key="dual_threshold")

        if st.button("🔗 듀얼 자동 매칭 실행", use_container_width=True, key="dual_run", disabled=not convert_copyright_ok):
            if not dual_hwp or not dual_pdf:
                st.error("HWP와 PDF를 모두 업로드해주세요.")
            else:
                try:
                    dual_hwp.seek(0)
                    dual_pdf.seek(0)
                    hwp_text = extract_text_from_hwp(dual_hwp)
                    pdf_bytes = dual_pdf.getvalue()
                    page_texts = extract_pdf_page_texts(pdf_bytes)
                    images = extract_images_from_pdf_bytes(pdf_bytes)
                    anchors = extract_pdf_question_anchors(pdf_bytes)

                    # 1) HWP 텍스트로 문항 파싱
                    items = parse_exam_text_fuzzy(hwp_text)
                    items = clean_parsed_items(items)

                    # 2) 문항-페이지 매칭
                    scores = match_questions_to_pages(items, page_texts)

                    # 3) 이미지 연결 (페이지 기반)
                    items = auto_attach_images_to_items(
                        items,
                        images,
                        strategy="page",
                        max_per_question=1,
                        anchors=anchors,
                        min_score=0.2,
                        only_if_keyword=False
                    )

                    st.session_state.past_exam_items = items
                    st.session_state.past_exam_images = images
                    st.session_state.past_exam_anchors = anchors
                    st.session_state.dual_exam_text = hwp_text
                    st.session_state.dual_exam_images = images
                    st.session_state.dual_exam_page_text = page_texts
                    st.session_state.dual_match_scores = scores

                    st.success(f"듀얼 매칭 완료: {len(items)}개 문항")
                    st.rerun()
                except Exception as e:
                    st.error(f"듀얼 매칭 실패: {str(e)}")

        if st.session_state.dual_match_scores:
            weak = [i for i, v in st.session_state.dual_match_scores.items() if v.get("score", 0) < dual_threshold]
            st.caption(f"자동 매칭 신뢰도 낮음: {len(weak)}개 문항 → 아래 편집 탭에서 수동 보정하세요.")

        if st.button("📝 HWP 텍스트만 추출", use_container_width=True, key="dual_text_only", disabled=not convert_copyright_ok):
            if not dual_hwp:
                st.error("HWP 파일을 업로드해주세요.")
            else:
                try:
                    dual_hwp.seek(0)
                    hwp_text = extract_text_from_hwp(dual_hwp)
                    hwp_text = preclean_exam_text(hwp_text)
                    items = parse_exam_text_fuzzy(hwp_text)
                    items = clean_parsed_items(items)
                    st.session_state.past_exam_items = items
                    st.session_state.past_exam_images = []
                    st.session_state.past_exam_anchors = {}
                    st.session_state.dual_exam_text = hwp_text
                    st.success(f"HWP 텍스트 추출 완료: {len(items)}개 문항")
                    st.rerun()
                except Exception as e:
                    st.error(f"HWP 텍스트 추출 실패: {str(e)}")

    uploaded_exam = st.file_uploader(
        "기출문제 파일 업로드",
        type=["hwp", "pdf", "docx", "pptx", "txt", "tsv"],
        key="past_exam_upload"
    )

    if uploaded_exam and not convert_copyright_ok:
        st.warning("저작권 확인 체크를 완료하면 업로드 파일을 변환할 수 있습니다.")
    elif uploaded_exam:
        file_ext = Path(uploaded_exam.name).suffix.lower()
        ocr_enabled = True
        ocr_engine = "auto"
        ocr_langs = ("ko", "en")
        ocr_max_pages = 0
        uploaded_bytes = uploaded_exam.getvalue()

        if file_ext == ".pdf":
            with st.expander("🧠 OCR 설정 (스캔 PDF용)", expanded=False):
                ocr_enabled = st.checkbox(
                    "텍스트가 부족하면 OCR 자동 실행",
                    value=True,
                    key="past_exam_ocr_enable"
                )
                ocr_engine = st.selectbox(
                    "OCR 엔진",
                    ["auto", "easyocr"],
                    index=0,
                    key="past_exam_ocr_engine"
                )
                lang_choice = st.selectbox(
                    "언어",
                    ["한국어+영어", "영어"],
                    index=0,
                    key="past_exam_ocr_lang"
                )
                ocr_langs = ("ko", "en") if lang_choice == "한국어+영어" else ("en",)
                ocr_max_pages = st.number_input(
                    "OCR 페이지 제한 (0=전체)",
                    min_value=0,
                    max_value=500,
                    value=0,
                    step=1,
                    key="past_exam_ocr_pages"
                )

        if st.session_state.past_exam_file != uploaded_exam.name:
            st.session_state.past_exam_file = uploaded_exam.name
            st.session_state.past_exam_text = ""
            st.session_state.past_exam_items = []
            st.session_state.past_exam_images = []
            st.session_state.past_exam_anchors = {}
            st.session_state.ai_parse_raw = ""

        if not st.session_state.past_exam_text:
            try:
                if hasattr(uploaded_exam, "seek"):
                    uploaded_exam.seek(0)
                st.session_state.past_exam_text = extract_text_from_file(
                    uploaded_exam,
                    enable_ocr=ocr_enabled,
                    ocr_engine=ocr_engine,
                    ocr_langs=ocr_langs,
                    ocr_max_pages=ocr_max_pages,
                    include_page_markers=(file_ext == ".pdf")
                )
            except Exception as e:
                st.error(f"❌ 기출문제 파일 처리 실패: {str(e)}")

        if not st.session_state.past_exam_images and uploaded_bytes:
            try:
                if file_ext == ".pdf":
                    st.session_state.past_exam_images = extract_images_from_pdf_bytes(uploaded_bytes)
                    st.session_state.past_exam_anchors = extract_pdf_question_anchors(uploaded_bytes)
                elif file_ext == ".hwp":
                    st.session_state.past_exam_images = extract_images_from_hwp_bytes(uploaded_bytes)
            except Exception:
                st.session_state.past_exam_images = []

        if file_ext == ".pdf":
            engines = available_ocr_engines()
            if len(st.session_state.past_exam_text.strip()) < 200 and not engines:
                st.warning("PDF에서 텍스트가 거의 추출되지 않았습니다. OCR이 필요합니다. `python -m pip install easyocr` 설치 후 다시 시도하세요.")
            if st.button("🔁 원문 다시 추출", use_container_width=True, key="past_exam_reextract"):
                try:
                    if hasattr(uploaded_exam, "seek"):
                        uploaded_exam.seek(0)
                    st.session_state.past_exam_text = extract_text_from_file(
                        uploaded_exam,
                        enable_ocr=ocr_enabled,
                        ocr_engine=ocr_engine,
                        ocr_langs=ocr_langs,
                        ocr_max_pages=ocr_max_pages,
                        include_page_markers=True
                    )
                    st.session_state.past_exam_items = []
                    st.session_state.past_exam_images = extract_images_from_pdf_bytes(uploaded_bytes)
                    st.session_state.past_exam_anchors = extract_pdf_question_anchors(uploaded_bytes)
                except Exception as e:
                    st.error(f"❌ 원문 재추출 실패: {str(e)}")

        col1, col2 = st.columns(2)
        with col1:
            exam_subject = st.text_input("기본 과목명", value="General", key="past_exam_subject")
        with col2:
            default_unit = Path(uploaded_exam.name).stem[:50] if uploaded_exam else "미분류"
            exam_unit = st.text_input("기본 단원명 (선택)", value=default_unit, key="past_exam_unit")

        parse_mode = st.radio(
            "변환 방식",
            ["자동(기출 파서)", "Cloze(정답: 기반)", "객관식(선지 기준)"],
            horizontal=True,
            key="past_exam_mode"
        )

        st.markdown("**이미지 자동 연결**")
        auto_attach = st.checkbox("문항에 이미지 자동 연결", value=True, key="auto_attach_images")
        max_imgs = st.slider("문항당 최대 이미지 수", 0, 3, 1, key="auto_attach_max_images")
        only_attach_keyword = st.checkbox("이미지 키워드가 있는 문항만 연결", value=True, key="auto_attach_keyword_only")

        if file_ext == ".pdf":
            attach_label = st.selectbox(
                "자동 연결 방식",
                ["레이아웃 기반(권장)", "OCR 기반(텍스트 포함 이미지)", "페이지 기반"],
                index=0,
                key="auto_attach_mode"
            )
            if attach_label.startswith("OCR"):
                attach_strategy = "ocr"
                ocr_img_limit = st.slider("OCR 이미지 개수 제한", 5, 80, 20, key="ocr_img_limit")
                ocr_min_score = st.slider("매칭 기준(0~1)", 0.05, 0.6, 0.2, step=0.05, key="ocr_min_score")
            elif attach_label.startswith("페이지"):
                attach_strategy = "page"
            else:
                attach_strategy = "layout" if st.session_state.past_exam_anchors else "page"
            use_ai_match = st.checkbox("AI 이미지 매칭(보정)", value=False, key="ai_match_images")
            ai_match_limit = st.slider("AI 매칭 이미지 수", 1, 30, 8, key="ai_match_limit")
        else:
            attach_strategy = "sequential"

        st.text_area(
            "추출된 원문 (필요시 수정 가능)",
            value=st.session_state.past_exam_text,
            height=240,
            key="past_exam_text_area"
        )

        with st.expander("🤖 AI 파서 (문항 분리/정제)", expanded=False):
            st.caption("겹쳐진 문항을 분리하거나 주관식 문항을 구조화하고 싶을 때 사용합니다.")
            ai_parse_limit = st.slider("최대 문항 수", 10, 200, 60, step=10, key="ai_parse_limit")
            parse_mode_ai = st.radio("AI 파서 방식", ["전체 텍스트", "블록 분할"], horizontal=True, key="ai_parse_mode")
            hint_text = st.text_area(
                "문서 구조 힌트 (선택)",
                value="",
                placeholder="예: 2열 표 → 좌측 문항, 우측 정답/해설. 1열 표 → 문항→정답→해설 순서.",
                key="ai_parse_hint"
            )
            if file_ext == ".pdf":
                st.caption("PDF 레이아웃 파서는 2열(좌:문항/우:정답·해설) 또는 1열 구조에 최적화되어 있습니다.")
                use_ai_layout = st.checkbox(
                    "AI로 레이아웃 파서 실행(추천)",
                    value=True,
                    key="use_ai_layout_parser"
                )
                if st.button("📐 PDF 레이아웃 파서 실행", use_container_width=True, key="layout_parse_run"):
                    with st.spinner("PDF 레이아웃 분석 중..."):
                        layout_items = []
                        if use_ai_layout:
                            if st.session_state.ai_model == "🔵 Google Gemini" and not api_key:
                                st.error("Gemini API 키가 필요합니다. 사이드바에서 입력해주세요.")
                            elif st.session_state.ai_model == "🟢 OpenAI ChatGPT" and not openai_api_key:
                                st.error("OpenAI API 키가 필요합니다. 사이드바에서 입력해주세요.")
                            else:
                                layout_items = parse_pdf_layout_ai(
                                    uploaded_bytes,
                                    ai_model=st.session_state.ai_model,
                                    api_key=api_key,
                                    openai_api_key=openai_api_key,
                                    hint_text=hint_text
                                )
                        if not layout_items:
                            layout_items = parse_pdf_layout(uploaded_bytes)
                        if layout_items:
                            if auto_attach and st.session_state.past_exam_images:
                                layout_items = auto_attach_images_to_items(
                                    layout_items,
                                    st.session_state.past_exam_images,
                                    strategy=attach_strategy,
                                    max_per_question=max_imgs,
                                    anchors=st.session_state.past_exam_anchors,
                                    min_score=st.session_state.get("ocr_min_score", 0.2),
                                    only_if_keyword=only_attach_keyword
                                )
                            if st.session_state.get("ai_match_images") and st.session_state.past_exam_images:
                                layout_items = ai_match_images_to_items(
                                    layout_items,
                                    st.session_state.past_exam_images,
                                    ai_model=st.session_state.get("ai_model", "🔵 Google Gemini"),
                                    api_key=api_key,
                                    openai_api_key=openai_api_key,
                                    max_images=st.session_state.get("ai_match_limit", 8)
                                )
                            st.session_state.past_exam_items = layout_items
                            st.success(f"레이아웃 파서 완료: {len(layout_items)}개 문항")
                            st.rerun()
                        else:
                            st.warning("레이아웃 파서 결과가 비어있습니다. OCR 후 다시 시도하거나 AI 파서를 사용하세요.")
            if parse_mode_ai == "블록 분할":
                block_limit = st.slider("블록 처리 개수", 5, 200, 50, step=5, key="ai_block_limit")
            if st.button("AI 파서로 재분할", use_container_width=True, key="ai_parse_run"):
                if st.session_state.ai_model == "🔵 Google Gemini" and not api_key:
                    st.error("Gemini API 키가 필요합니다. 사이드바에서 입력해주세요.")
                elif st.session_state.ai_model == "🟢 OpenAI ChatGPT" and not openai_api_key:
                    st.error("OpenAI API 키가 필요합니다. 사이드바에서 입력해주세요.")
                else:
                    with st.spinner("AI 파서 실행 중..."):
                        source_text = st.session_state.get("past_exam_text_area", "")
                        if parse_mode_ai == "블록 분할":
                            blocks = split_exam_blocks(source_text)
                            ai_items = []
                            raw_chunks = []
                            for block in blocks[:block_limit]:
                                item, raw = ai_parse_exam_block(
                                    block,
                                    ai_model=st.session_state.ai_model,
                                    api_key=api_key,
                                    openai_api_key=openai_api_key,
                                    hint_text=hint_text,
                                    return_raw=True
                                )
                                if raw:
                                    raw_chunks.append(raw)
                                if item:
                                    ai_items.append(item)
                            ai_items = clean_parsed_items(ai_items)
                            st.session_state.ai_parse_raw = "\n\n---\n\n".join(raw_chunks)
                        else:
                            ai_items, raw = ai_parse_exam_text(
                                source_text,
                                ai_model=st.session_state.ai_model,
                                api_key=api_key,
                                openai_api_key=openai_api_key,
                                max_items=ai_parse_limit,
                                hint_text=hint_text,
                                return_raw=True
                            )
                            st.session_state.ai_parse_raw = raw
                        if ai_items:
                            if auto_attach and st.session_state.past_exam_images:
                                ai_items = auto_attach_images_to_items(
                                    ai_items,
                                    st.session_state.past_exam_images,
                                    strategy=attach_strategy,
                                    max_per_question=max_imgs,
                                    anchors=st.session_state.past_exam_anchors,
                                    min_score=st.session_state.get("ocr_min_score", 0.2)
                                )
                            st.session_state.past_exam_items = ai_items
                            st.success(f"AI 파서 완료: {len(ai_items)}개 문항")
                            st.rerun()
                        else:
                            st.warning("AI 파서 결과가 비어있습니다. 문서 구조 힌트를 더 구체적으로 입력하거나, 블록 분할 모드를 사용해보세요.")
                            raw = st.session_state.get("ai_parse_raw", "")
                            if raw:
                                with st.expander("AI 파서 원문 결과(디버그)", expanded=False):
                                    st.code(raw[:6000])

        if st.session_state.past_exam_images:
            with st.expander("🖼️ 추출된 이미지", expanded=False):
                st.caption(f"총 {len(st.session_state.past_exam_images)}개 이미지")
                cols = st.columns(4)
                for i, img in enumerate(st.session_state.past_exam_images):
                    with cols[i % 4]:
                        st.image(img.get("data_uri"), caption=f"#{i + 1}")

        if st.button("🔎 변환 미리보기", use_container_width=True, key="past_exam_preview"):
            source_text = st.session_state.get("past_exam_text_area", "").strip()
            if not source_text:
                st.error("추출된 텍스트가 비어 있습니다.")
            else:
                if parse_mode == "Cloze(정답: 기반)":
                    items = parse_qa_to_cloze(source_text)
                    if not items:
                        items = parse_generated_text_to_structured(source_text, "🧩 빈칸 뚫기 (Anki Cloze)")
                elif parse_mode == "객관식(선지 기준)":
                    if file_ext == ".pdf":
                        use_ai_layout = st.session_state.get("use_ai_layout_parser", True)
                        if use_ai_layout and ((st.session_state.ai_model == "🔵 Google Gemini" and api_key) or (st.session_state.ai_model == "🟢 OpenAI ChatGPT" and openai_api_key)):
                            items = [i for i in parse_pdf_layout_ai(
                                uploaded_bytes,
                                ai_model=st.session_state.ai_model,
                                api_key=api_key,
                                openai_api_key=openai_api_key,
                                hint_text=st.session_state.get("ai_parse_hint", "")
                            ) if i.get("type") == "mcq"]
                        else:
                            items = [i for i in parse_pdf_layout(uploaded_bytes) if i.get("type") == "mcq"]
                    else:
                        items = [i for i in parse_exam_text_fuzzy(source_text) if i.get("type") == "mcq"]
                    if not items:
                        items = parse_generated_text_to_structured(source_text, "📝 객관식 문제 (Case Study)")
                else:
                    if file_ext == ".pdf":
                        use_ai_layout = st.session_state.get("use_ai_layout_parser", True)
                        if use_ai_layout and ((st.session_state.ai_model == "🔵 Google Gemini" and api_key) or (st.session_state.ai_model == "🟢 OpenAI ChatGPT" and openai_api_key)):
                            items = parse_pdf_layout_ai(
                                uploaded_bytes,
                                ai_model=st.session_state.ai_model,
                                api_key=api_key,
                                openai_api_key=openai_api_key,
                                hint_text=st.session_state.get("ai_parse_hint", "")
                            )
                        else:
                            items = parse_pdf_layout(uploaded_bytes)
                    else:
                        items = parse_exam_text_fuzzy(source_text)
                    if not items:
                        items = parse_exam_text_fuzzy(source_text)
                    if not items:
                        items = parse_generated_text_to_structured(source_text, "📝 객관식 문제 (Case Study)")
                        if not items:
                            items = parse_qa_to_cloze(source_text)
                if items and auto_attach and st.session_state.past_exam_images:
                    if attach_strategy == "ocr":
                        st.session_state.past_exam_images = ocr_images_for_matching(
                            st.session_state.past_exam_images,
                            engine="easyocr",
                            langs=("ko", "en"),
                            max_images=st.session_state.get("ocr_img_limit", 20)
                        )
                    items = auto_attach_images_to_items(
                        items,
                        st.session_state.past_exam_images,
                        strategy=attach_strategy,
                        max_per_question=max_imgs,
                        anchors=st.session_state.past_exam_anchors,
                        min_score=st.session_state.get("ocr_min_score", 0.2),
                        only_if_keyword=only_attach_keyword
                    )
                if items and st.session_state.get("ai_match_images") and st.session_state.past_exam_images:
                    if st.session_state.ai_model == "🔵 Google Gemini" and not api_key:
                        st.error("Gemini API 키가 필요합니다. 사이드바에서 입력해주세요.")
                    elif st.session_state.ai_model == "🟢 OpenAI ChatGPT" and not openai_api_key:
                        st.error("OpenAI API 키가 필요합니다. 사이드바에서 입력해주세요.")
                    else:
                        items = ai_match_images_to_items(
                            items,
                            st.session_state.past_exam_images,
                            ai_model=st.session_state.get("ai_model", "🔵 Google Gemini"),
                            api_key=api_key,
                            openai_api_key=openai_api_key,
                            max_images=st.session_state.get("ai_match_limit", 8)
                        )
                st.session_state.past_exam_items = items if items else []

        items = st.session_state.get("past_exam_items", [])
        if items:
            st.success(f"✅ 변환된 문항: {len(items)}개")
            with st.expander("📋 변환 결과 미리보기 (상위 5개)", expanded=True):
                for i, item_data in enumerate(items[:5], 1):
                    if item_data.get("type") == "mcq":
                        st.markdown(f"**문제 {i}** (객관식)")
                        st.write(f"**문항:** {item_data.get('problem', '')[:150]}...")
                        st.write(f"**선지:** {', '.join(item_data.get('options', [])[:3])}...")
                        st.write(f"**정답:** {item_data.get('answer', '?')} 번")
                    else:
                        st.markdown(f"**문제 {i}** (빈칸)")
                        st.write(f"**내용:** {item_data.get('front', '')[:150]}...")
                        st.write(f"**정답:** {item_data.get('answer', '?')}")
                    st.divider()

            with st.expander("🛠️ 문항 편집", expanded=False):
                total_items = len(items)
                if total_items > 0:
                    start_idx = st.number_input("시작 문항", min_value=1, max_value=total_items, value=1, step=1, key="edit_start_idx")
                    end_idx = st.number_input("끝 문항", min_value=start_idx, max_value=total_items, value=min(start_idx + 9, total_items), step=1, key="edit_end_idx")
                    image_options = list(range(len(st.session_state.past_exam_images)))

                    def image_label(i):
                        img = st.session_state.past_exam_images[i]
                        page = img.get("page")
                        return f"#{i + 1} | p{page}" if page else f"#{i + 1}"

                    for i in range(start_idx - 1, end_idx):
                        item = items[i]
                        with st.container():
                            qnum_label = f"q{item.get('qnum')}" if item.get("qnum") else "q?"
                            page_label = f"p{item.get('page')}" if item.get("page") else "p?"
                            st.markdown(f"#### 문항 {i + 1} 편집 ({item.get('type')}) · {qnum_label} · {page_label}")
                            item_type = st.selectbox(
                                "유형",
                                ["mcq", "cloze"],
                                index=0 if item.get("type") == "mcq" else 1,
                                key=f"edit_type_{i}"
                            )
                            if item_type == "mcq":
                                st.text_area("문항", value=item.get("problem", ""), height=120, key=f"edit_problem_{i}")
                                opts = item.get("options", [])
                                st.text_area("선지 (한 줄에 하나)", value="\n".join(opts), height=140, key=f"edit_options_{i}")
                                ans_default = int(item.get("answer", 1)) if str(item.get("answer", "")).isdigit() else 1
                                st.selectbox("정답", [1, 2, 3, 4, 5], index=max(0, min(ans_default - 1, 4)), key=f"edit_answer_{i}")
                            else:
                                st.text_area("문항", value=item.get("front", ""), height=120, key=f"edit_front_{i}")
                                st.text_input("정답", value=item.get("answer", ""), key=f"edit_answer_{i}")
                            st.text_area("해설", value=item.get("explanation", ""), height=120, key=f"edit_expl_{i}")
                            if image_options:
                                current_images = item.get("images", [])
                                current_indices = [idx for idx, img in enumerate(st.session_state.past_exam_images) if img.get("data_uri") in current_images]

                                img_pages = sorted({img.get("page") for img in st.session_state.past_exam_images if img.get("page")})
                                page_options = ["전체"] + [f"p{p}" for p in img_pages]
                                page_filter = st.selectbox("이미지 페이지 필터", page_options, key=f"img_page_filter_{i}")
                                per_page = st.slider("페이지당 이미지", 4, 24, 8, key=f"img_per_page_{i}")

                                filtered_indices = []
                                for idx_img, img in enumerate(st.session_state.past_exam_images):
                                    page = img.get("page")
                                    if page_filter != "전체":
                                        wanted = int(page_filter.replace("p", ""))
                                        if page != wanted:
                                            continue
                                    filtered_indices.append(idx_img)

                                total_imgs = len(filtered_indices)
                                total_pages = max(1, (total_imgs + per_page - 1) // per_page)
                                page_idx = st.number_input("이미지 페이지", 1, total_pages, 1, key=f"img_page_idx_{i}")
                                start = (page_idx - 1) * per_page
                                end = start + per_page
                                subset = filtered_indices[start:end]

                                cols = st.columns(4)
                                for j, idx_img in enumerate(subset):
                                    img = st.session_state.past_exam_images[idx_img]
                                    with cols[j % 4]:
                                        st.image(img.get("data_uri"), width=140, caption=image_label(idx_img))
                                        st.checkbox(
                                            "선택",
                                            value=idx_img in current_indices,
                                            key=f"edit_img_{i}_{idx_img}"
                                        )
                            st.checkbox("이 문항 삭제", key=f"edit_delete_{i}")
                            st.divider()

                    if st.button("✅ 편집 내용 적용", use_container_width=True, key="apply_edits"):
                        new_items = []
                        for i in range(total_items):
                            if st.session_state.get(f"edit_delete_{i}"):
                                continue
                            item = items[i]
                            item_type = st.session_state.get(f"edit_type_{i}", item.get("type"))
                            if item_type == "mcq":
                                problem = st.session_state.get(f"edit_problem_{i}", item.get("problem", "")).strip()
                                options_text = st.session_state.get(f"edit_options_{i}", "\n".join(item.get("options", [])))
                                options = [o.strip() for o in options_text.splitlines() if o.strip()]
                                answer = st.session_state.get(f"edit_answer_{i}", item.get("answer", 1))
                                updated = {
                                    **item,
                                    "type": "mcq",
                                    "problem": problem,
                                    "options": options,
                                    "answer": int(answer) if str(answer).isdigit() else 1,
                                }
                            else:
                                front = st.session_state.get(f"edit_front_{i}", item.get("front", "")).strip()
                                answer = st.session_state.get(f"edit_answer_{i}", item.get("answer", "")).strip()
                                updated = {
                                    **item,
                                    "type": "cloze",
                                    "front": front,
                                    "answer": answer,
                                }
                            updated["explanation"] = st.session_state.get(f"edit_expl_{i}", item.get("explanation", "")).strip()
                            if image_options:
                                current_images = item.get("images", [])
                                current_indices = [idx for idx, img in enumerate(st.session_state.past_exam_images) if img.get("data_uri") in current_images]
                                sel_set = set(current_indices)
                                for idx_img in image_options:
                                    key = f"edit_img_{i}_{idx_img}"
                                    if key in st.session_state:
                                        if st.session_state.get(key):
                                            sel_set.add(idx_img)
                                        else:
                                            sel_set.discard(idx_img)
                                updated["images"] = [st.session_state.past_exam_images[idx]["data_uri"] for idx in sorted(sel_set)]
                            new_items.append(updated)
                        st.session_state.past_exam_items = new_items
                        st.success("편집 내용이 반영되었습니다.")
                        st.rerun()

            col_save, col_down = st.columns(2)
            with col_save:
                if st.button("💾 문항 저장", use_container_width=True, key="past_exam_save"):
                    current_items = st.session_state.get("past_exam_items", [])
                    added = add_questions_to_bank_auto(
                        current_items,
                        subject=exam_subject,
                        unit=exam_unit,
                        quality_filter=enable_filter,
                        min_length=min_length
                    )
                    st.success(f"✅ {added}개 문항 저장 완료")
            with col_down:
                download_data = json.dumps(items, ensure_ascii=False, indent=2)
                st.download_button(
                    label="📥 JSON으로 다운로드",
                    data=download_data,
                    file_name="converted_exam_questions.json",
                    mime="application/json",
                    use_container_width=True,
                    key="past_exam_download"
                )
        elif uploaded_exam:
            st.info("변환 미리보기를 눌러 문항을 생성하세요.")

if active_page == "exam":
    st.title("🎯 실전 모의고사")
    st.caption("이 탭은 API 키 없이도 저장된 문항으로 학습/시험이 가능합니다.")
    
    bank = load_questions()
    
    if not bank["text"] and not bank["cloze"]:
        st.warning("📌 저장된 문제가 없습니다. 먼저 **📚 문제 생성** 탭에서 문제를 생성하세요.")
    else:
        st.info("기출문제 파일 변환은 **🧾 기출문제 변환** 탭에서 진행합니다.")
        if st.session_state.get("exam_mode_entry_anchor") and st.session_state.get("exam_questions"):
            st.success(
                f"생성 결과로 {len(st.session_state.exam_questions)}개 문항이 준비되어 있습니다. "
                f"아래 버튼으로 즉시 학습/시험을 이어서 시작할 수 있습니다."
            )
            col_resume1, col_resume2 = st.columns(2)
            with col_resume1:
                if st.button("✅ 준비된 세션 이어 풀기", use_container_width=True, key="resume_prepared_exam"):
                    st.session_state.exam_started = True
                    st.session_state.exam_finished = False
                    st.session_state.current_question_idx = 0
                    st.session_state.exam_mode_entry_anchor = ""
                    st.rerun()
            with col_resume2:
                if st.button("🗑 준비 세션 초기화", use_container_width=True, key="clear_prepared_exam"):
                    st.session_state.exam_started = False
                    st.session_state.exam_finished = False
                    st.session_state.exam_questions = []
                    st.session_state.current_question_idx = 0
                    st.session_state.exam_mode_entry_anchor = ""
                    st.rerun()

        # 시험/학습 설정
        if MOBILE_CLIENT:
            st.markdown("<div class='mobile-exam-caption'>모바일 풀이 모드: 터치 중심 UI</div>", unsafe_allow_html=True)
            mode_choice = st.radio("모드", ["시험모드", "학습모드"], horizontal=False)
            exam_type = st.selectbox("문항 유형", ["객관식", "빈칸"])
            mobile_image_width = max(220, min(640, int(st.session_state.image_display_width)))
            st.session_state.image_display_width = st.slider(
                "문항 이미지 크기(px)",
                220,
                640,
                mobile_image_width,
                step=10,
                key="image_display_width_slider"
            )
        else:
            c_mode, c_type, c_img = st.columns([1.2, 1, 1])
            with c_mode:
                mode_choice = st.radio("모드", ["시험모드", "학습모드"], horizontal=True)
            with c_type:
                exam_type = st.selectbox("문항 유형", ["객관식", "빈칸"])
            with c_img:
                st.session_state.image_display_width = st.slider(
                    "문항 이미지 크기(px)",
                    240,
                    900,
                    st.session_state.image_display_width,
                    step=20,
                    key="image_display_width_slider"
                )

        questions_all = bank["text"] if exam_type == "객관식" else bank["cloze"]
        subject_unit_map = collect_subject_unit_map(questions_all)
        all_subjects = sorted(subject_unit_map.keys())
        if all_subjects:
            subject_keyword = st.text_input("분과 검색", value="", placeholder="분과명 입력", key="exam_subject_search")
            subject_pool = [s for s in all_subjects if subject_keyword.lower() in s.lower()]
            if not subject_pool:
                subject_pool = ["(검색 결과 없음)"]
            if "exam_subject_multi" not in st.session_state:
                st.session_state["exam_subject_multi"] = list(all_subjects)
            selected_subjects = st.multiselect(
                "분과 선택",
                options=subject_pool if subject_pool != ["(검색 결과 없음)"] else [],
                key="exam_subject_multi"
            )
            if not selected_subjects:
                # 빈 선택은 전체 보기로 복구해 실수로 인한 빈 화면을 방지
                selected_subjects = all_subjects

            unit_filter_by_subject = {}
            selected_units = []
            if selected_subjects:
                with st.expander("단원 선택 (분과별)", expanded=True):
                    for subj in selected_subjects:
                        units = subject_unit_map.get(subj, ["미분류"])
                        if not units:
                            units = ["미분류"]
                        unit_key = f"unit_filter_{subj}"
                        previous = st.session_state.get(unit_key, units)
                        default_units = previous if set(previous) <= set(units) else units
                        selected_units_for_subject = st.multiselect(
                            f"{subj} 단원",
                            options=units,
                            default=default_units,
                            key=unit_key
                        )
                        if not selected_units_for_subject:
                            selected_units_for_subject = list(units)
                        unit_filter_by_subject[subj] = selected_units_for_subject
                        selected_units.extend(selected_units_for_subject)
            else:
                unit_filter_by_subject = {}
                selected_units = []
            filtered_questions = filter_questions_by_subject_unit_hierarchy(questions_all, selected_subjects, unit_filter_by_subject)
        else:
            selected_subjects = []
            selected_units = []
            filtered_questions = []


        learning_session_mode = "탐색형(단원 전체)"
        bookmarked_only = False
        exam_distribution_mode = "비례(보유 문항 기준)"
        exam_group_mode = "분과+단원"
        exam_seed = None
        if mode_choice == "학습모드":
            due_only = st.checkbox("오늘 복습만", value=False)
            bookmarked_only = st.checkbox("북마크 문항만", value=False)
            learning_session_mode = st.radio(
                "학습 세션 방식",
                ["탐색형(단원 전체)", "랜덤형(문항 수 선택)"],
                horizontal=True,
            )
            st.session_state.auto_next = st.checkbox("자동 다음 문제", value=st.session_state.auto_next)
            if due_only:
                filtered_questions = [q for q in filtered_questions if srs_due(q)]
            if bookmarked_only:
                filtered_questions = [q for q in filtered_questions if bool(q.get("bookmarked"))]
            if not FSRS_AVAILABLE:
                st.info("FSRS 미설치: 기본 복습 주기(SRS)로 동작합니다.")
        else:
            st.session_state.auto_next = False
            exam_distribution_mode = st.radio(
                "출제 비율",
                ["비례(보유 문항 기준)", "균등(선택 그룹 기준)"],
                horizontal=True,
            )
            exam_group_mode = st.radio(
                "출제 그룹",
                ["분과+단원", "분과"],
                horizontal=True,
            )
            use_seed = st.checkbox("랜덤 시드 고정", value=False)
            if use_seed:
                exam_seed = int(st.number_input("랜덤 시드", min_value=0, value=42, step=1))

        if mode_choice == "학습모드":
            with st.expander("📅 FSRS 복습 큐", expanded=False):
                show_queue = st.checkbox("복습 큐 표시", value=False, key="show_fsrs_queue")
                if show_queue:
                    if FSRS_AVAILABLE:
                        stats = get_fsrs_stats(filtered_questions)
                        if stats:
                            col1, col2, col3, col4 = st.columns(4)
                            with col1:
                                st.metric("오늘 복습", stats["due"])
                            with col2:
                                st.metric("연체", stats["overdue"])
                            with col3:
                                st.metric("미래", stats["future"])
                            with col4:
                                st.metric("신규", stats["new"])

                        due_list = get_fsrs_queue(filtered_questions, limit=20)
                        if not due_list:
                            st.info("오늘 복습할 문항이 없습니다.")
                        else:
                            rows = []
                            for q, due_time in due_list:
                                snippet = (q.get("problem") or q.get("front") or "").strip()
                                snippet = snippet[:80] + "..." if len(snippet) > 80 else snippet
                                rows.append({
                                    "분과": q.get("subject") or "General",
                                    "문항": snippet,
                                    "Due": due_time.isoformat()
                                })
                            safe_dataframe(rows, use_container_width=True, hide_index=True)
                    else:
                        due_list = [q for q in filtered_questions if simple_srs_due(q)]
                        st.metric("오늘 복습", len(due_list))
                        if not due_list:
                            st.info("오늘 복습할 문항이 없습니다.")

            with st.expander("⚙️ FSRS 설정", expanded=False):
                if not FSRS_AVAILABLE:
                    st.info("FSRS 패키지가 설치되지 않아 설정을 사용할 수 없습니다.")
                else:
                    st.caption("FSRS 설정은 다음 복습부터 적용됩니다.")
                    desired_retention = st.slider(
                        "목표 기억 유지율",
                        0.7,
                        0.98,
                        float(st.session_state.fsrs_desired_retention),
                        0.01,
                        key="fsrs_desired_retention_slider"
                    )
                    learning_steps_text = st.text_input(
                        "학습 단계(분, 콤마)",
                        value=st.session_state.fsrs_learning_steps_text,
                        key="fsrs_learning_steps_input"
                    )
                    relearning_steps_text = st.text_input(
                        "재학습 단계(분, 콤마)",
                        value=st.session_state.fsrs_relearning_steps_text,
                        key="fsrs_relearning_steps_input"
                    )
                    max_interval = st.number_input(
                        "최대 간격(일)",
                        min_value=30,
                        max_value=365000,
                        value=int(st.session_state.fsrs_max_interval),
                        step=30,
                        key="fsrs_max_interval_input"
                    )
                    enable_fuzzing = st.checkbox(
                        "간격 랜덤화(Fuzzing) 사용",
                        value=bool(st.session_state.fsrs_enable_fuzzing),
                        key="fsrs_enable_fuzzing_input"
                    )
                    advanced = st.checkbox("고급: 파라미터 직접 입력", value=False, key="fsrs_params_toggle")
                    params_text = None
                    if advanced:
                        params_text = st.text_area(
                            "FSRS parameters (JSON 배열)",
                            value=st.session_state.fsrs_params_text,
                            height=120,
                            key="fsrs_params_input"
                        )
                        st.caption("파라미터를 잘못 입력하면 기본값으로 동작합니다.")

                    col_a, col_b = st.columns(2)
                    with col_a:
                        if st.button("✅ FSRS 설정 저장", use_container_width=True, key="fsrs_save_btn"):
                            steps = [s.strip() for s in learning_steps_text.split(",") if s.strip()]
                            relearn_steps = [s.strip() for s in relearning_steps_text.split(",") if s.strip()]
                            try:
                                params = json.loads(params_text) if advanced and params_text else list(FSRS_DEFAULT_PARAMETERS)
                                if not isinstance(params, list) or len(params) < 10:
                                    params = list(FSRS_DEFAULT_PARAMETERS)
                            except Exception:
                                params = list(FSRS_DEFAULT_PARAMETERS)
                            settings = {
                                "desired_retention": float(desired_retention),
                                "learning_steps": [int(s) for s in steps if s.isdigit()],
                                "relearning_steps": [int(s) for s in relearn_steps if s.isdigit()],
                                "maximum_interval": int(max_interval),
                                "enable_fuzzing": bool(enable_fuzzing),
                                "parameters": params,
                            }
                            save_fsrs_settings(settings)
                            st.session_state.fsrs_desired_retention = settings["desired_retention"]
                            st.session_state.fsrs_learning_steps_text = ",".join(map(str, settings["learning_steps"]))
                            st.session_state.fsrs_relearning_steps_text = ",".join(map(str, settings["relearning_steps"]))
                            st.session_state.fsrs_max_interval = settings["maximum_interval"]
                            st.session_state.fsrs_enable_fuzzing = settings["enable_fuzzing"]
                            st.session_state.fsrs_params_text = json.dumps(settings["parameters"])
                            st.success("FSRS 설정이 저장되었습니다.")
                    with col_b:
                        if st.button("↩️ 기본값으로 초기화", use_container_width=True, key="fsrs_reset_btn"):
                            settings = load_fsrs_settings()
                            st.session_state.fsrs_desired_retention = settings["desired_retention"]
                            st.session_state.fsrs_learning_steps_text = ",".join(map(str, settings["learning_steps"]))
                            st.session_state.fsrs_relearning_steps_text = ",".join(map(str, settings["relearning_steps"]))
                            st.session_state.fsrs_max_interval = settings["maximum_interval"]
                            st.session_state.fsrs_enable_fuzzing = settings["enable_fuzzing"]
                            st.session_state.fsrs_params_text = json.dumps(settings["parameters"])
                            st.success("FSRS 기본값으로 초기화했습니다.")

            with st.expander("📈 복습 리포트", expanded=False):
                show_report = st.checkbox("리포트 표시", value=False, key="show_fsrs_report")
                if show_report:
                    if FSRS_AVAILABLE:
                        report = get_fsrs_report(filtered_questions)
                        if report:
                            st.metric("총 카드", report["total"])
                            st.metric("최근 7일 리뷰 수", report["review_count_7d"])
                            st.metric("평균 간격(일)", f"{report['avg_interval']:.1f}")
                            if report["last_review"]:
                                st.caption(f"마지막 리뷰: {report['last_review']}")

                            rating_rows = [{"평가": k, "건수": v} for k, v in report["rating_counts"].items()]
                            safe_dataframe(rating_rows, use_container_width=True, hide_index=True)
                        else:
                            st.info("리포트를 생성할 수 없습니다.")
                    else:
                        st.info("기본 SRS 모드에서는 상세 리포트를 제공하지 않습니다.")

        if questions_all:
            with st.expander("📤 시험지/문제집 내보내기", expanded=False):
                st.caption("선택한 분과 문항을 2열(DOCX) 형식으로 내보냅니다. 좌측: 문항, 우측: 정답/해설")
                export_title_default = f"AxiomaQbank_{exam_type}_문제집"
                export_title = st.text_input("문서 제목", value=export_title_default, key="export_docx_title")
                export_subjects = st.multiselect(
                    "내보낼 분과 선택",
                    options=all_subjects,
                    default=[s for s in selected_subjects if s in all_subjects] if selected_subjects else all_subjects,
                    key="export_subjects"
                )
                export_include_all_units = st.checkbox(
                    "선택 분과 전체 문항 사용 (단원 필터 무시)",
                    value=True,
                    key="export_include_all_units"
                )
                export_unit_filter_by_subject = {}
                if not export_include_all_units and export_subjects:
                    st.markdown("**내보낼 단원 선택**")
                    for subj in export_subjects:
                        export_units = subject_unit_map.get(subj, ["미분류"])
                        if not export_units:
                            export_units = ["미분류"]
                        export_unit_filter_by_subject[subj] = st.multiselect(
                            f"{subj} 단원 (내보내기)",
                            options=export_units,
                            default=export_units,
                            key=f"export_unit_filter_{subj}"
                        )
                export_randomize = st.checkbox("랜덤 배치 모드", value=False, key="export_randomize")
                export_seed = None
                if export_randomize:
                    export_seed = st.number_input("랜덤 시드", min_value=0, value=42, step=1, key="export_random_seed")
                if export_subjects:
                    export_candidates = collect_export_questions(
                        questions_all,
                        export_subjects,
                        export_unit_filter_by_subject,
                        include_all_units=export_include_all_units,
                        randomize=export_randomize,
                        random_seed=export_seed
                    )
                else:
                    export_candidates = []
                st.caption(f"내보내기 대상 문항: {len(export_candidates)}개")
                if st.button("DOCX 생성", key="build_docx_export", use_container_width=True):
                    if not export_candidates:
                        st.warning("내보낼 문항이 없습니다. 분과/단원 선택을 확인해주세요.")
                    else:
                        st.session_state.export_docx_bytes = build_docx_question_sheet(export_candidates, title=export_title)
                        st.success("DOCX 생성 완료")
                if st.session_state.get("export_docx_bytes"):
                    st.download_button(
                        "📥 DOCX 다운로드",
                        data=st.session_state.export_docx_bytes,
                        file_name=f"{export_title}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key="download_docx_export",
                        use_container_width=True
                    )
        else:
            st.session_state.export_docx_bytes = b""

        if not filtered_questions:
            st.warning("선택한 조건에 해당하는 문제가 없습니다.")
        else:
            max_questions = len(filtered_questions)
            max_limit = min(150, max(1, max_questions))
            if mode_choice == "학습모드" and str(learning_session_mode).startswith("탐색형"):
                num_questions = max_questions
                st.caption(f"탐색형: 선택한 범위의 {max_questions}문항 전체를 불러옵니다.")
            else:
                default_num = min(10, max_limit)
                num_questions = st.slider("문항 수", 1, max_limit, default_num)

            start_label = "📝 시험 시작" if mode_choice == "시험모드" else "📖 학습 시작"
            if st.button(start_label, use_container_width=True, key="start_exam"):
                if len(filtered_questions) < num_questions:
                    st.warning(f"문제가 부족합니다. {len(filtered_questions)}개만 출제합니다.")
                    num_questions = len(filtered_questions)

                if mode_choice == "시험모드":
                    raw_selected = select_exam_questions_balanced(
                        filtered_questions,
                        num_questions,
                        distribution_mode=exam_distribution_mode,
                        group_mode=exam_group_mode,
                        random_seed=exam_seed,
                    )
                else:
                    raw_selected = select_learning_session_questions(
                        filtered_questions,
                        learning_mode=learning_session_mode,
                        num_questions=num_questions,
                        random_seed=exam_seed,
                    )
                distribution_counts = {}
                if mode_choice == "시험모드":
                    for q_raw in raw_selected:
                        gk = _exam_group_key(q_raw, group_mode=exam_group_mode)
                        distribution_counts[gk] = int(distribution_counts.get(gk, 0)) + 1
                parsed_selected = []
                for raw in raw_selected:
                    if exam_type == "객관식":
                        parsed = parse_mcq_content(raw)
                    else:
                        parsed = parse_cloze_content(raw)
                    parsed_selected.append(parsed)

                st.session_state.exam_questions = parsed_selected
                st.session_state.current_question_idx = 0
                st.session_state.user_answers = {}
                st.session_state.exam_started = True
                st.session_state.exam_finished = False
                st.session_state.exam_mode = mode_choice
                st.session_state.exam_type = exam_type
                st.session_state.auto_advance_guard = None
                st.session_state.revealed_answers = set()
                st.session_state.exam_stats_applied = False
                st.session_state.graded_questions = set()
                st.session_state.exam_history_saved = False
                st.session_state.current_exam_meta = {
                    "mode": mode_choice,
                    "type": exam_type,
                    "subjects": selected_subjects,
                    "units": selected_units,
                    "num_questions": len(parsed_selected),
                    "learning_session_mode": learning_session_mode if mode_choice == "학습모드" else "",
                    "distribution_mode": exam_distribution_mode if mode_choice == "시험모드" else "",
                    "distribution_group_mode": exam_group_mode if mode_choice == "시험모드" else "",
                    "distribution_counts": distribution_counts if mode_choice == "시험모드" else {},
                    "seed": exam_seed,
                    "started_at": datetime.now(timezone.utc).isoformat()
                }

        # 시험/학습 진행
        if st.session_state.exam_started and st.session_state.exam_questions:
            exam_qs = st.session_state.exam_questions
            idx = st.session_state.current_question_idx

            if st.session_state.exam_finished:
                st.markdown("## 📊 결과")

                total = len(exam_qs)
                answered = len(st.session_state.user_answers)

                # 정답 채점
                correct_count = 0
                wrong_indices = []
                for i, q in enumerate(exam_qs):
                    if i not in st.session_state.user_answers:
                        continue

                    user_ans = st.session_state.user_answers[i]
                    if is_answer_correct(q, user_ans):
                        correct_count += 1
                    else:
                        wrong_indices.append(i)

                # 통계 업데이트 (시험 결과 1회만, 이미 반영된 문항은 제외)
                if not st.session_state.exam_stats_applied:
                    for i, q in enumerate(exam_qs):
                        if i in st.session_state.user_answers and q.get("id"):
                            if q.get("id") in st.session_state.graded_questions:
                                continue
                            user_ans = st.session_state.user_answers[i]
                            is_correct = is_answer_correct(q, user_ans)
                            updated_stats = update_question_stats(q["id"], is_correct)
                            if isinstance(updated_stats, dict):
                                q["stats"] = updated_stats
                            st.session_state.graded_questions.add(q.get("id"))
                    st.session_state.exam_stats_applied = True

                # 시험 기록 저장 (시험모드만)
                if st.session_state.exam_mode == "시험모드" and not st.session_state.exam_history_saved:
                    items = []
                    for i, q in enumerate(exam_qs):
                        user_ans = st.session_state.user_answers.get(i)
                        items.append({
                            "id": q.get("id"),
                            "type": q.get("type"),
                            "front": q.get("front"),
                            "options": q.get("options"),
                            "correct": q.get("correct"),
                            "answer": q.get("answer"),
                            "user": user_ans,
                            "is_correct": is_answer_correct(q, user_ans) if user_ans is not None else False,
                            "explanation": q.get("explanation"),
                            "subject": q.get("subject"),
                            "difficulty": q.get("difficulty"),
                            "note": q.get("note", ""),
                        })
                    meta = st.session_state.current_exam_meta or {}
                    session = {
                        "session_id": str(uuid.uuid4()),
                        "finished_at": datetime.now(timezone.utc).isoformat(),
                        "mode": meta.get("mode", st.session_state.exam_mode),
                        "type": meta.get("type", st.session_state.exam_type),
                        "subjects": meta.get("subjects", []),
                        "num_questions": len(exam_qs),
                        "answered": answered,
                        "correct": correct_count,
                        "accuracy": int(correct_count / answered * 100) if answered > 0 else 0,
                        "items": items
                    }
                    add_exam_history(session)
                    st.session_state.exam_history_saved = True

                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("정답", f"{correct_count}/{answered}")
                with col2:
                    st.metric("미응답", f"{total - answered}")
                with col3:
                    accuracy = int(correct_count / answered * 100) if answered > 0 else 0
                    st.metric("정확도", f"{accuracy}%")
                with col4:
                    st.metric("상태", "✅ 완료" if answered == total else "⚠️ 미완료")

                st.markdown("---")

                # 상세 보기
                letters = ['A', 'B', 'C', 'D', 'E']
                for i, q in enumerate(exam_qs, 1):
                    user_ans = st.session_state.user_answers.get(i - 1, None)
                    is_correct = False
                    correct_text = ""
                    correct_display = ""

                    if q.get('type') == 'mcq':
                        correct_num = q.get('correct')  # 숫자 형식: 1-5
                        correct_text = str(correct_num)
                        correct_display = letters[correct_num - 1] if 1 <= correct_num <= 5 else "?"
                        is_correct = (user_ans == correct_num) if user_ans else False
                        user_ans_display = letters[user_ans - 1] if user_ans and 1 <= user_ans <= 5 else "응답 없음"
                    else:
                        response_type = q.get("response_type", "cloze")
                        correct_text = q.get('answer') or ""
                        correct_display = correct_text
                        if response_type == "essay":
                            ai_grade = q.get("_ai_grade") if isinstance(q.get("_ai_grade"), dict) else {}
                            is_correct = bool(ai_grade.get("is_correct", False))
                        else:
                            is_correct = fuzzy_match(user_ans, correct_text) if user_ans and correct_text else False
                        user_ans_display = user_ans if user_ans else "응답 없음"

                    status_icon = "✅" if is_correct else "❌"
                    with st.expander(f"{status_icon} 문제 {i}: {user_ans_display}"):
                        st.markdown(q.get('front', q.get('raw', '')))

                        if q.get('type') == 'mcq':
                            st.markdown("**선택지:**")
                            opts = q.get('options') or []
                            for idx_opt, opt in enumerate(opts[:5]):
                                label = f"{letters[idx_opt]}. {opt}"
                                st.write(label)

                        st.divider()
                        st.write(f"**당신의 답:** {user_ans_display}")
                        answer_color = "🟢" if is_correct else "🔴"
                        st.write(f"{answer_color} **정답:** {correct_display}")
                        if q.get("response_type") == "essay":
                            if isinstance(q.get("_ai_grade"), dict):
                                st.write(f"AI 점수: {q['_ai_grade'].get('score', 0)} / 100")
                                feedback = q["_ai_grade"].get("feedback")
                                if feedback:
                                    st.write(f"AI 피드백: {feedback}")
                            else:
                                st.caption("서술형은 AI 채점 실행 전까지 정오 판정이 확정되지 않습니다.")
                        if q.get("explanation"):
                            show_exp = st.checkbox("해설 보기", value=st.session_state.explanation_default, key=f"show_exp_{i}")
                            if show_exp:
                                st.markdown(format_explanation_text(q.get('explanation')))
                        if q.get("subject"):
                            st.caption(f"📌 {q['subject']}")
                        if q.get("unit"):
                            st.caption(f"단원: {q.get('unit')}")
                        if q.get("difficulty"):
                            st.caption(f"난이도: {q.get('difficulty', '?')}")
                        if q.get("id"):
                            note_key = f"review_note_{i}"
                            st.text_area("메모", value=q.get("note", ""), key=note_key, height=80)
                            if st.button("메모 저장", key=f"save_review_note_{i}"):
                                saved = update_question_note(q["id"], st.session_state.get(note_key, ""))
                                if saved:
                                    q["note"] = st.session_state.get(note_key, "")
                                    st.success("메모 저장됨")

                # 오답노트
                if wrong_indices:
                    if st.button("📌 오답노트로 다시 풀기"):
                        wrong_qs = [exam_qs[i] for i in wrong_indices]
                        st.session_state.exam_questions = wrong_qs
                        st.session_state.user_answers = {}
                        st.session_state.current_question_idx = 0
                        st.session_state.exam_started = True
                        st.session_state.exam_finished = False
                        st.session_state.exam_mode = "학습모드"
                        st.session_state.revealed_answers = set()
                        st.session_state.auto_advance_guard = None
                        st.session_state.exam_stats_applied = False
                        st.session_state.graded_questions = set()
                        st.rerun()

                if st.button("🔄 다시 시작"):
                    st.session_state.exam_started = False
                    st.session_state.exam_finished = False
                    st.session_state.exam_questions = []
                    st.session_state.user_answers = {}
                    st.session_state.current_question_idx = 0
                    st.rerun()



            else:
                if idx < len(exam_qs):
                    q = exam_qs[idx]
                    st.progress((idx + 1) / len(exam_qs))
                    st.caption(f"USMLE 스타일 | Question {idx + 1} of {len(exam_qs)}")
                    nav_slot = st.empty()
                    unanswered_slot = st.empty()
                    st.markdown(f"### Question {idx + 1}")
                    if q.get("type") != "mcq":
                        rt = q.get("response_type", "cloze")
                        rt_label = "빈칸형" if rt == "cloze" else ("단답형" if rt == "short" else "서술형")
                        st.caption(f"유형: {rt_label}")

                    # 입력
                    if q.get('type') == 'mcq':
                        st.markdown(q.get('front', ''))
                        if q.get("images"):
                            st.image(q.get("images"), width=st.session_state.image_display_width)

                        st.markdown("**Select one option (A–E):**")
                        opts = q.get('options') or []
                        letters = ['A', 'B', 'C', 'D', 'E']
                        prev_ans = st.session_state.user_answers.get(idx)
                        default_index = (prev_ans - 1) if isinstance(prev_ans, int) and 1 <= prev_ans <= 5 else None
                        if opts:
                            labels_real = [f"{letters[i]}. {opts[i]}" for i in range(min(len(opts), len(letters)))]
                            st.session_state[f"labels_real_{idx}"] = labels_real
                            user_choice_label = st.radio("정답 선택:", labels_real, index=default_index, key=f"q_{idx}")
                            if user_choice_label:
                                chosen_num = letters.index(user_choice_label.split(".")[0]) + 1
                                st.session_state.user_answers[idx] = chosen_num
                            else:
                                st.session_state.user_answers.pop(idx, None)
                        else:
                            st.session_state[f"labels_real_{idx}"] = letters
                            user_choice = st.radio("정답 선택:", letters, index=default_index, key=f"q_{idx}")
                            if user_choice:
                                chosen_num = letters.index(user_choice) + 1
                                st.session_state.user_answers[idx] = chosen_num
                            else:
                                st.session_state.user_answers.pop(idx, None)

                        if not MOBILE_CLIENT:
                            st.text_input(
                                "키보드 입력 (A-E 또는 1-5)",
                                key=f"shortcut_{idx}",
                                on_change=apply_mcq_shortcut,
                                args=(idx,)
                            )

                        if idx in st.session_state.user_answers:
                            your = st.session_state.user_answers[idx]
                            your_letter = letters[your - 1] if 1 <= your <= 5 else "?"
                            st.caption(f"📍 Your answer: {your_letter}")
                    else:
                        st.markdown(q.get('front', q.get('raw', '')))
                        if q.get("images"):
                            st.image(q.get("images"), width=st.session_state.image_display_width)
                        prev_text = st.session_state.user_answers.get(idx, "")
                        response_type = q.get("response_type", "cloze")
                        if response_type == "essay":
                            user_input = st.text_area("서술형 답안 입력:", value=prev_text, key=f"cloze_{idx}", height=160)
                        elif response_type == "short":
                            user_input = st.text_input("단답형 정답 입력:", value=prev_text, key=f"cloze_{idx}")
                        else:
                            user_input = st.text_input("정답 입력 (한글/영문):", value=prev_text, key=f"cloze_{idx}")
                        if user_input:
                            st.session_state.user_answers[idx] = user_input
                        elif idx in st.session_state.user_answers:
                            st.session_state.user_answers.pop(idx, None)

                        if response_type == "essay" and user_input:
                            if st.button("🧠 AI 채점 (서술형)", key=f"grade_essay_{idx}"):
                                if st.session_state.ai_model == "🔵 Google Gemini" and not api_key:
                                    st.error("Gemini API 키가 필요합니다. 사이드바에서 입력해주세요.")
                                elif st.session_state.ai_model == "🟢 OpenAI ChatGPT" and not openai_api_key:
                                    st.error("OpenAI API 키가 필요합니다. 사이드바에서 입력해주세요.")
                                else:
                                    with st.spinner("AI 채점 중..."):
                                        grade, err = grade_essay_answer_ai(
                                            q,
                                            user_input,
                                            ai_model=st.session_state.ai_model,
                                            api_key=api_key,
                                            openai_api_key=openai_api_key
                                        )
                                    if grade:
                                        q["_ai_grade"] = grade
                                        st.success(f"AI 채점 완료: {grade.get('score', 0)}점")
                                    else:
                                        st.warning(f"AI 채점 실패: {err}")
                        if response_type == "essay" and isinstance(q.get("_ai_grade"), dict):
                            st.caption(f"AI 점수: {q['_ai_grade'].get('score', 0)} / 100")
                            feedback = q["_ai_grade"].get("feedback")
                            if feedback:
                                st.caption(f"피드백: {feedback}")

                    # 누적 풀이 정보 + 북마크
                    attempt = get_question_attempt_summary(q)
                    if attempt["attempts"] > 0:
                        last_dt = parse_iso_datetime(attempt.get("last_time"))
                        last_text = ""
                        if last_dt:
                            if last_dt.tzinfo is None:
                                last_dt = last_dt.replace(tzinfo=timezone.utc)
                            last_text = last_dt.astimezone().strftime("%Y-%m-%d %H:%M")
                        verdict = "정답" if attempt.get("last_correct") is True else ("오답" if attempt.get("last_correct") is False else "-")
                        info = f"누적 풀이 {attempt['attempts']}회 (정 {attempt['right']} / 오 {attempt['wrong']}) · 최근 {verdict}"
                        if last_text:
                            info += f" · {last_text}"
                        st.caption(info)
                    else:
                        st.caption("첫 풀이 문항")

                    if q.get("id"):
                        mark_label = "⭐ 다시보기 해제" if q.get("bookmarked") else "☆ 다시보기 저장"
                        if st.button(mark_label, key=f"bookmark_{idx}", use_container_width=False):
                            new_mark = not bool(q.get("bookmarked"))
                            if update_question_bookmark(q["id"], new_mark):
                                q["bookmarked"] = new_mark
                                st.success("다시보기 목록이 업데이트되었습니다.")
                                st.rerun()

                    # 문항 이동/미응답 (답안 반영 후 갱신)
                    answered_idx = set(st.session_state.user_answers.keys())
                    nav_options = list(range(len(exam_qs)))

                    def nav_format(i):
                        star = "⭐ " if exam_qs[i].get("bookmarked") else ""
                        status = "✅" if i in answered_idx else "○"
                        return f"{i + 1} {star}{status}"

                    if MOBILE_CLIENT:
                        nav_idx = nav_slot.select_slider(
                            "문항 이동",
                            options=nav_options,
                            value=idx,
                            format_func=nav_format,
                            key="nav_select_mobile",
                        )
                    else:
                        nav_idx = nav_slot.selectbox(
                            "문항 이동",
                            nav_options,
                            index=idx,
                            format_func=nav_format,
                            key="nav_select",
                        )
                    if nav_idx != idx:
                        st.session_state.current_question_idx = nav_idx

                    unanswered = [str(i + 1) for i in range(len(exam_qs)) if i not in answered_idx]
                    if unanswered:
                        unanswered_slot.caption(f"미응답: {', '.join(unanswered)}")

                    # 메모
                    if q.get("id"):
                        note_key = f"note_{idx}"
                        st.text_area("메모", value=q.get("note", ""), key=note_key, height=80)
                        if st.button("메모 저장", key=f"save_note_{idx}"):
                            saved = update_question_note(q["id"], st.session_state.get(note_key, ""))
                            if saved:
                                q["note"] = st.session_state.get(note_key, "")
                                st.success("메모 저장됨")

                    # 학습모드: 정답 확인 후 표시
                    if st.session_state.exam_mode == "학습모드" and idx in st.session_state.user_answers:
                        st.markdown("---")
                        reveal_key = f"reveal_{idx}"
                        if st.button("정답 확인", key=reveal_key):
                            st.session_state.revealed_answers.add(idx)

                        if idx in st.session_state.revealed_answers:
                            if q.get('type') == 'mcq':
                                correct_num = q.get('correct')
                                correct_display = letters[correct_num - 1] if isinstance(correct_num, int) and 1 <= correct_num <= 5 else "?"
                                is_correct = (st.session_state.user_answers[idx] == correct_num) if correct_num else False
                            else:
                                response_type = q.get("response_type", "cloze")
                                correct_text = q.get('answer') or ""
                                if response_type == "essay":
                                    if st.button("🧠 AI 채점 실행", key=f"learn_grade_essay_{idx}"):
                                        if st.session_state.ai_model == "🔵 Google Gemini" and not api_key:
                                            st.error("Gemini API 키가 필요합니다. 사이드바에서 입력해주세요.")
                                        elif st.session_state.ai_model == "🟢 OpenAI ChatGPT" and not openai_api_key:
                                            st.error("OpenAI API 키가 필요합니다. 사이드바에서 입력해주세요.")
                                        else:
                                            with st.spinner("AI 채점 중..."):
                                                grade, err = grade_essay_answer_ai(
                                                    q,
                                                    st.session_state.user_answers[idx],
                                                    ai_model=st.session_state.ai_model,
                                                    api_key=api_key,
                                                    openai_api_key=openai_api_key
                                                )
                                            if grade:
                                                q["_ai_grade"] = grade
                                            else:
                                                st.warning(f"AI 채점 실패: {err}")
                                    is_correct = bool(isinstance(q.get("_ai_grade"), dict) and q["_ai_grade"].get("is_correct"))
                                else:
                                    is_correct = fuzzy_match(st.session_state.user_answers[idx], correct_text) if correct_text else False
                                correct_display = correct_text

                            answer_color = "🟢" if is_correct else "🔴"
                            st.write(f"{answer_color} **정답:** {correct_display}")
                            if q.get("response_type") == "essay":
                                if isinstance(q.get("_ai_grade"), dict):
                                    st.write(f"AI 점수: {q['_ai_grade'].get('score', 0)} / 100")
                                    feedback = q["_ai_grade"].get("feedback")
                                    if feedback:
                                        st.write(f"AI 피드백: {feedback}")
                                else:
                                    st.info("서술형은 AI 채점 실행 후 정오 판정이 반영됩니다.")
                            # 학습모드 통계 업데이트 (1회)
                            if q.get("id") and q.get("id") not in st.session_state.graded_questions:
                                updated_stats = update_question_stats(q["id"], is_correct)
                                if isinstance(updated_stats, dict):
                                    q["stats"] = updated_stats
                                st.session_state.graded_questions.add(q.get("id"))
                            explanation_text = q.get("explanation") or q.get("rationale") or q.get("analysis") or ""
                            show_exp = st.checkbox("해설 보기", value=st.session_state.explanation_default, key=f"learn_exp_{idx}")
                            if show_exp:
                                if explanation_text.strip():
                                    st.markdown(format_explanation_text(explanation_text))
                                else:
                                    st.caption("해설이 없습니다.")
                                    if st.button("AI 해설 생성", key=f"ai_exp_{idx}"):
                                        if st.session_state.ai_model == "🔵 Google Gemini" and not api_key:
                                            st.error("Gemini API 키가 필요합니다. 사이드바에서 입력해주세요.")
                                        elif st.session_state.ai_model == "🟢 OpenAI ChatGPT" and not openai_api_key:
                                            st.error("OpenAI API 키가 필요합니다. 사이드바에서 입력해주세요.")
                                        else:
                                            with st.spinner("AI 해설 생성 중..."):
                                                text, err = generate_single_explanation_ai(
                                                    q,
                                                    ai_model=st.session_state.ai_model,
                                                    api_key=api_key,
                                                    openai_api_key=openai_api_key,
                                                    return_error=True
                                                )
                                            if text:
                                                q["explanation"] = text
                                                if q.get("id"):
                                                    update_question_explanation(q["id"], text)
                                                st.success("해설이 생성되었습니다.")
                                                st.markdown(format_explanation_text(text))
                                            else:
                                                msg = f"해설 생성 실패. 다시 시도해주세요."
                                                if err:
                                                    msg += f" (에러: {err})"
                                                st.warning(msg)

                            if q.get("id"):
                                st.markdown("**복습 평가**")
                                cols = st.columns(4)
                                if cols[0].button("Again", key=f"srs_again_{idx}"):
                                    rating = Rating.Again if FSRS_AVAILABLE else "Again"
                                    srs = apply_srs_rating(q["id"], rating)
                                    if srs:
                                        q["fsrs"] = srs if FSRS_AVAILABLE else q.get("fsrs")
                                        st.success(f"다음 복습: {srs.get('due')}")
                                if cols[1].button("Hard", key=f"srs_hard_{idx}"):
                                    rating = Rating.Hard if FSRS_AVAILABLE else "Hard"
                                    srs = apply_srs_rating(q["id"], rating)
                                    if srs:
                                        q["fsrs"] = srs if FSRS_AVAILABLE else q.get("fsrs")
                                        st.success(f"다음 복습: {srs.get('due')}")
                                if cols[2].button("Good", key=f"srs_good_{idx}"):
                                    rating = Rating.Good if FSRS_AVAILABLE else "Good"
                                    srs = apply_srs_rating(q["id"], rating)
                                    if srs:
                                        q["fsrs"] = srs if FSRS_AVAILABLE else q.get("fsrs")
                                        st.success(f"다음 복습: {srs.get('due')}")
                                if cols[3].button("Easy", key=f"srs_easy_{idx}"):
                                    rating = Rating.Easy if FSRS_AVAILABLE else "Easy"
                                    srs = apply_srs_rating(q["id"], rating)
                                    if srs:
                                        q["fsrs"] = srs if FSRS_AVAILABLE else q.get("fsrs")
                                        st.success(f"다음 복습: {srs.get('due')}")

                    # 학습모드 자동 다음 문제
                    if st.session_state.exam_mode == "학습모드" and st.session_state.auto_next:
                        guard = st.session_state.auto_advance_guard
                        current_answer = st.session_state.user_answers.get(idx)
                        if current_answer and idx in st.session_state.revealed_answers and guard != (idx, str(current_answer)) and idx < len(exam_qs) - 1:
                            st.session_state.auto_advance_guard = (idx, str(current_answer))
                            st.session_state.current_question_idx += 1
                            st.rerun()

                    # 네비게이션
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.button("⬅️ 이전", on_click=goto_prev_question, disabled=idx <= 0)
                    with col2:
                        st.button("다음 ➡️", on_click=goto_next_question, disabled=idx >= len(exam_qs) - 1)
                    with col3:
                        if st.session_state.exam_mode == "시험모드":
                            if idx == len(exam_qs) - 1:
                                st.button("✅ 채점", on_click=finish_exam_session)
                        else:
                            if idx == len(exam_qs) - 1:
                                st.button("✅ 세션 종료", on_click=finish_exam_session)
